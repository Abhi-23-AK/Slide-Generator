import io
import re
import os
import tempfile
import json
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn


def analyze_template(template_bytes: bytes, max_slides: int = 30) -> dict:
    """Analyze a PPTX template and return a comprehensive schema of its structure.

    Extracts detailed shape properties, paragraph settings, font styling, vertical
    alignments, word wrap flags, capacity estimates, and semantic roles for every
    shape on every slide, and returns them under the `zones` list.
    """
    print(f"[TEMPLATE] Starting analysis of custom template ({len(template_bytes)} bytes)")
    prs = Presentation(io.BytesIO(template_bytes))

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # --- Extract theme info from the first slide master ---
    theme_info = _extract_theme_info(prs)
    print(f"[TEMPLATE] Extracted theme: heading_font={theme_info['heading_font']}, body_font={theme_info['body_font']}")

    # --- Analyze each slide ---
    slides_info = []
    all_bg_colors = set()

    for idx, slide in enumerate(prs.slides):
        if idx >= max_slides:
            break
        
        print(f"\n[TEMPLATE] --- Analyzing Slide {idx + 1}/{len(prs.slides)} ---")
        slide_data = _analyze_slide(slide, idx, slide_width, slide_height)
        
        bg_color = slide_data.get("background_color")
        if bg_color and bg_color != "image":
            all_bg_colors.add(bg_color)
            
        slides_info.append(slide_data)

    # Store background colors discovered across slides in theme
    theme_info["background_colors"] = list(all_bg_colors)

    schema = {
        "slide_count": len(slides_info),
        "theme": theme_info,
        "slides": slides_info,
    }
    
    print("\n[TEMPLATE] Analysis completed successfully.")
    return schema


def _extract_theme_info(prs):
    """Extract theme fonts and colors from the first slide master."""
    theme_info = {
        "heading_font": "Calibri",
        "body_font": "Calibri",
        "colors": {
            "dk1": "#000000",
            "lt1": "#FFFFFF",
            "dk2": "#1F4E79",
            "lt2": "#D9E1F2",
            "accent1": "#4472C4",
            "accent2": "#ED7D31",
            "accent3": "#A5A5A5",
            "accent4": "#FFC000",
            "accent5": "#5B9BD5",
            "accent6": "#70AD47",
        }
    }

    try:
        if not prs.slide_masters:
            return theme_info

        master = prs.slide_masters[0]

        # Extract major and minor fonts
        font_scheme = master.element.findall('.//' + qn('a:fontScheme'))
        if font_scheme:
            major = font_scheme[0].find(qn('a:majorFont'))
            if major is not None:
                latin = major.find(qn('a:latin'))
                if latin is not None and latin.get('typeface'):
                    theme_info['heading_font'] = latin.get('typeface')

            minor = font_scheme[0].find(qn('a:minorFont'))
            if minor is not None:
                latin = minor.find(qn('a:latin'))
                if latin is not None and latin.get('typeface'):
                    theme_info['body_font'] = latin.get('typeface')

        # Extract clrScheme colors
        clr_scheme = master.element.findall('.//' + qn('a:clrScheme'))
        if clr_scheme:
            color_names = [
                'dk1', 'lt1', 'dk2', 'lt2',
                'accent1', 'accent2', 'accent3', 'accent4',
                'accent5', 'accent6'
            ]
            for color_name in color_names:
                color_el = clr_scheme[0].find(qn(f'a:{color_name}'))
                if color_el is not None:
                    # Check for solid hex value (srgbClr)
                    srgb = color_el.find(qn('a:srgbClr'))
                    if srgb is not None and srgb.get('val'):
                        theme_info['colors'][color_name] = f'#{srgb.get("val")}'
                    else:
                        # Fallback to system color (sysClr)
                        sys_clr = color_el.find(qn('a:sysClr'))
                        if sys_clr is not None and sys_clr.get('lastClr'):
                            theme_info['colors'][color_name] = f'#{sys_clr.get("lastClr")}'
    except Exception as e:
        print(f"[TEMPLATE] [Warning] Failed to read theme XML scheme: {e}")

    return theme_info


def _analyze_slide(slide, idx, slide_width, slide_height):
    """Analyze a single slide and return its structure dict."""
    # --- Extract background color ---
    bg_color_hex = _get_background_color(slide)
    print(f"[TEMPLATE] Slide {idx + 1} background color: {bg_color_hex}")

    # --- Extract all shapes on the slide ---
    shapes_list = []
    for shape_idx, shape in enumerate(slide.shapes):
        shape_data = _analyze_shape(shape, shape_idx, slide_width, slide_height)
        shapes_list.append(shape_data)

    # --- Determine semantic roles and card groups ---
    card_groups = _determine_semantic_roles(shapes_list)

    # --- Infer layout type ---
    layout_type = _infer_layout_type(idx, shapes_list)

    # --- Determine column count based on card groups or body zones ---
    has_cards = any(z["role"].startswith("card") for z in shapes_list)
    if has_cards:
        card_roles = [z["role"] for z in shapes_list if z["role"].startswith("card")]
        card_group_ids = set(int(r.split("_")[0].replace("card", "")) for r in card_roles if "_" in r)
        column_count = len(card_group_ids)
    else:
        body_shapes = [s for s in shapes_list if s['role'].startswith('body')]
        column_count = len(body_shapes) if len(body_shapes) <= 3 else 1

    # --- Detect image placeholder ---
    has_image_placeholder = any(
        s['role'].startswith('image')
        for s in shapes_list
    )

    print(f"[TEMPLATE] Slide {idx + 1} inferred layout: {layout_type} | columns: {column_count} | image_placeholder: {has_image_placeholder}")
    
    # Log shapes found
    for s in shapes_list:
        print(f"  Shape [{s['shape_id']}] name='{s['name']}' type={s['shape_type']} role={s['role']} capacity={s.get('capacity', 0)} pos=({s['left_pct']}%, {s['top_pct']}%, {s['width_pct']}%, {s['height_pct']}%)")

    return {
        "slide_index": idx,
        "inferred_layout": layout_type,
        "zones": shapes_list,
        "column_count": column_count,
        "has_image_placeholder": has_image_placeholder,
        "background_color": bg_color_hex,
        "card_groups": card_groups,
    }


def _get_background_color(slide):
    """Extract background color from slide or its layout."""
    bg_color = None
    try:
        if slide.background and slide.background.fill:
            fill = slide.background.fill
            if fill.type == 1:  # Solid fill
                bg_color = f"#{fill.fore_color.rgb}"
            elif fill.type == 5:  # Picture fill
                bg_color = "image"
    except Exception:
        pass

    if not bg_color:
        try:
            layout = slide.slide_layout
            if layout and layout.background and layout.background.fill:
                fill = layout.background.fill
                if fill.type == 1:
                    bg_color = f"#{fill.fore_color.rgb}"
                elif fill.type == 5:
                    bg_color = "image"
        except Exception:
            pass

    return bg_color


def _analyze_shape(shape, z_order, slide_width, slide_height):
    """Extract deep properties and settings of a shape."""
    # Positions as percentage of slide dimensions
    left_pct = (shape.left / slide_width) * 100 if slide_width else 0
    top_pct = (shape.top / slide_height) * 100 if slide_height else 0
    width_pct = (shape.width / slide_width) * 100 if slide_width else 0
    height_pct = (shape.height / slide_height) * 100 if slide_height else 0

    # Basic shape parameters
    shape_type_name = "unknown"
    shape_type_id = None
    if shape.shape_type is not None:
        try:
            shape_type_name = shape.shape_type.name
            shape_type_id = int(shape.shape_type)
        except Exception:
            shape_type_name = str(shape.shape_type)

    placeholder_idx = None
    placeholder_type = None
    if shape.is_placeholder:
        try:
            placeholder_idx = shape.placeholder_format.idx
            ph_type = shape.placeholder_format.type
            placeholder_type = ph_type.name if hasattr(ph_type, "name") else str(ph_type)
        except Exception:
            pass

    rotation = 0
    try:
        rotation = shape.rotation
    except AttributeError:
        pass

    # Picture verification
    is_picture = False
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            is_picture = True
    except Exception:
        pass
    if not is_picture and shape.is_placeholder:
        try:
            if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                is_picture = True
        except Exception:
            pass

    # Text frame property defaults
    font_name = None
    font_size = None
    font_color = None
    bold = False
    italic = False
    underline = False
    alignment = "left"
    line_spacing = None
    space_before = None
    space_after = None
    vertical_anchor = "top"
    word_wrap = None

    if shape.has_text_frame:
        tf = shape.text_frame
        word_wrap = tf.word_wrap
        
        # Vertical anchor alignment
        if tf.vertical_anchor is not None:
            try:
                vertical_anchor = tf.vertical_anchor.name.lower()
            except Exception:
                vertical_anchor = str(tf.vertical_anchor).lower()

        if tf.paragraphs:
            p = tf.paragraphs[0]
            
            # Paragraph alignments
            if p.alignment is not None:
                try:
                    alignment = p.alignment.name.lower()
                except Exception:
                    alignment = str(p.alignment).lower()
                    
            # Line and space heights
            if p.line_spacing is not None:
                try:
                    line_spacing = float(p.line_spacing.pt) if hasattr(p.line_spacing, 'pt') else float(p.line_spacing)
                except Exception:
                    line_spacing = str(p.line_spacing)
            if p.space_before is not None:
                try:
                    space_before = float(p.space_before.pt)
                except Exception:
                    pass
            if p.space_after is not None:
                try:
                    space_after = float(p.space_after.pt)
                except Exception:
                    pass

            # Extract run properties from the first run
            for run in p.runs:
                if run.font.name:
                    font_name = run.font.name
                if run.font.size:
                    try:
                        font_size = float(run.font.size.pt)
                    except Exception:
                        pass
                if run.font.color and run.font.color.type is not None:
                    try:
                        font_color = f"#{run.font.color.rgb}"
                    except Exception:
                        pass
                bold = bool(run.font.bold)
                italic = bool(run.font.italic)
                underline = bool(run.font.underline)
                break

    # Estimate text capacity (words)
    capacity = 0
    if shape.has_text_frame:
        try:
            width_in = shape.width.inches
            height_in = shape.height.inches
            area = width_in * height_in
            f_sz = font_size if font_size else 14.0
            
            # Text cell area estimation formula
            word_area = (f_sz ** 2 * 3.6) / 5184
            capacity = int(area / word_area) if word_area > 0 else 0
            capacity = max(5, min(500, capacity))
        except Exception:
            capacity = 40

    return {
        "shape_id": shape.shape_id,
        "name": shape.name or "",
        "shape_type": shape_type_name,
        "shape_type_id": shape_type_id,
        "placeholder_idx": placeholder_idx,
        "placeholder_type": placeholder_type,
        "z_order": z_order,
        "left_pct": round(left_pct, 2),
        "top_pct": round(top_pct, 2),
        "width_pct": round(width_pct, 2),
        "height_pct": round(height_pct, 2),
        "rotation": rotation,
        "has_text_frame": shape.has_text_frame,
        "is_picture": is_picture,
        "font_name": font_name,
        "font_size": font_size,
        "font_color": font_color,
        "bold": bold,
        "italic": italic,
        "underline": underline,
        "alignment": alignment,
        "line_spacing": line_spacing,
        "space_before": space_before,
        "space_after": space_after,
        "vertical_anchor": vertical_anchor,
        "word_wrap": word_wrap,
        "capacity": capacity,
    }


def _determine_semantic_roles(shapes_list):
    """Determine semantic roles for all shapes on a slide, detecting card groups first."""
    if not shapes_list:
        return []

    # 1. Handle decorative / very small shapes first
    # Exclude from card clustering: lines, separators, arrows, chart/graph elements,
    # tiny labels, small circles, decorative rectangles, borders.
    for s in shapes_list:
        width = s["width_pct"]
        height = s["height_pct"]
        name = s["name"].lower()
        shape_type = s["shape_type"].upper()
        
        is_dec = False
        if width < 3 or height < 3:
            is_dec = True
        elif "line" in name or "line" in shape_type:
            is_dec = True
        elif any(k in name for k in ("separator", "divider", "border", "arrow", "graphic", "drawing", "axis", "legend", "plot", "series", "scrollbar")):
            is_dec = True
        # Small circles (oval/ellipse shape name, size < 4)
        elif (any(k in name for k in ("ellipse", "oval", "circle")) or shape_type == "OVAL") and (width < 4 and height < 4):
            is_dec = True
        # Decorative rectangles (rectangles that are very narrow or border-like)
        elif (any(k in name for k in ("rectangle", "rect")) or shape_type == "RECTANGLE") and (width < 5 or height < 5):
            is_dec = True
            
        if is_dec:
            s["role"] = f"decorative_{s['shape_id']}"

    # 2. Basic classification for slide-level titles, subtitles, footers, speaker notes, tables, charts
    for s in shapes_list:
        if s.get("role"):  # Skip if already marked decorative
            continue
            
        name = s["name"].lower()
        ph_type = s.get("placeholder_type")
        top = s["top_pct"]
        width = s["width_pct"]
        height = s["height_pct"]
        has_text = s["has_text_frame"]
        shape_type = s["shape_type"]

        # Tables and Charts
        if shape_type == "TABLE" or "table" in name:
            s["role"] = "table_placeholder"
            continue
        if shape_type in ("CHART", "GRAPHIC_FRAME") or any(k in name for k in ("chart", "graph", "axis", "legend", "plot", "series")):
            s["role"] = "chart_placeholder"
            continue

        # Slide-level title
        if ph_type in ("TITLE", "CENTER_TITLE", "VERTICAL_TITLE"):
            s["role"] = "title_placeholder"
            continue
        if "title" in name and not any(k in name for k in ("subtitle", "caption", "card", "ellipse", "oval", "circle")):
            if width >= 40:
                s["role"] = "title_placeholder"
                continue

        # Slide-level subtitle
        if ph_type == "SUBTITLE":
            s["role"] = "subtitle_placeholder"
            continue
        if "subtitle" in name and "card" not in name:
            if width >= 40:
                s["role"] = "subtitle_placeholder"
                continue

        # Slide metadata / footer
        if any(k in name for k in ("footer", "page", "slide number", "date")):
            s["role"] = "footer_placeholder"
            continue
        if top > 88 and height < 10 and width < 40 and has_text:
            s["role"] = "footer_placeholder"
            continue
            
        # Speaker notes
        if "note" in name or "speaker" in name:
            s["role"] = "note_placeholder"
            continue

        # Replaceable Image Zone (Pictures occupying significant area: width > 15% or height > 15%)
        # Exclude shapes named like icons
        is_pic = s.get("is_picture", False)
        is_icon_circle = any(k in name for k in ("icon", "logo", "symbol", "ellipse", "oval", "circle", "vector"))
        if (is_pic or any(k in name for k in ("picture", "image", "photo", "placeholder"))) and not is_icon_circle:
            if width > 15 or height > 15:
                s["role"] = "image_placeholder"
                continue

    # Main Slide Title / Subtitle fallback
    has_title = any(s.get("role") in ("title_placeholder", "title") for s in shapes_list)
    if not has_title:
        top_textboxes = [s for s in shapes_list if s["has_text_frame"] and s["top_pct"] < 25 and s["width_pct"] >= 40 and not s.get("role")]
        if top_textboxes:
            top_textboxes.sort(key=lambda x: x["top_pct"])
            top_textboxes[0]["role"] = "title_placeholder"
            
    # Subtitle fallback
    has_subtitle = any(s.get("role") in ("subtitle_placeholder", "subtitle") for s in shapes_list)
    if not has_subtitle:
        title_top_list = [s["top_pct"] for s in shapes_list if s.get("role") in ("title_placeholder", "title")]
        title_top = min(title_top_list) if title_top_list else 0
        sub_textboxes = [s for s in shapes_list if s["has_text_frame"] and title_top < s["top_pct"] < 40 and s["width_pct"] >= 40 and not s.get("role")]
        if sub_textboxes:
            sub_textboxes.sort(key=lambda x: x["top_pct"])
            sub_textboxes[0]["role"] = "subtitle_placeholder"

    # Map placeholders to unique sequential roles
    title_count = 0
    subtitle_count = 0
    table_count = 0
    chart_count = 0
    footer_count = 0
    note_count = 0
    image_count = 0
    
    for s in shapes_list:
        role = s.get("role")
        if role == "title_placeholder":
            title_count += 1
            s["role"] = "title" if title_count == 1 else f"title{title_count}"
        elif role == "subtitle_placeholder":
            subtitle_count += 1
            s["role"] = "subtitle" if subtitle_count == 1 else f"subtitle{subtitle_count}"
        elif role == "table_placeholder":
            table_count += 1
            s["role"] = f"table{table_count}"
        elif role == "chart_placeholder":
            chart_count += 1
            s["role"] = f"chart{chart_count}"
        elif role == "footer_placeholder":
            footer_count += 1
            s["role"] = f"footer{footer_count}"
        elif role == "note_placeholder":
            note_count += 1
            s["role"] = f"note{note_count}"
        elif role == "image_placeholder":
            image_count += 1
            s["role"] = f"image{image_count}"

    # 3. Card Group Detection using coordinate clustering
    candidates = []
    for s in shapes_list:
        role = s.get("role")
        if role:  # Skip already assigned slide-level/decorative shapes
            continue
        # Exclude extremely large slide background/border boxes
        if s["width_pct"] > 85 and s["height_pct"] > 85:
            continue
        candidates.append(s)
        
    def get_shape_distance(s1, s2):
        l1, t1, w1, h1 = s1["left_pct"], s1["top_pct"], s1["width_pct"], s1["height_pct"]
        l2, t2, w2, h2 = s2["left_pct"], s2["top_pct"], s2["width_pct"], s2["height_pct"]
        r1, b1 = l1 + w1, t1 + h1
        r2, b2 = l2 + w2, t2 + h2
        dx = max(0, l2 - r1, l1 - r2)
        dy = max(0, t2 - b1, t1 - b2)
        return (dx**2 + (dy * 1.5)**2)**0.5

    clusters = [[s] for s in candidates]
    # Reduced clustering threshold for tighter visual proximity
    threshold = 5.0
    
    while len(clusters) > 1:
        min_dist = float("inf")
        merge_pair = None
        for i in range(len(clusters)):
            for j in range(i + 1, len(clusters)):
                for s1 in clusters[i]:
                    for s2 in clusters[j]:
                        d = get_shape_distance(s1, s2)
                        if d < min_dist:
                            min_dist = d
                            merge_pair = (i, j)
        if merge_pair is not None and min_dist < threshold:
            i, j = merge_pair
            clusters[i].extend(clusters[j])
            clusters.pop(j)
        else:
            break

    # Separate card clusters (must contain at least one text shape and not be standalone small/narrow shape)
    card_clusters = []
    standalone_shapes = []
    for c in clusters:
        # Shapes with specific icon names are never treated as text shapes
        has_text = any(
            s["has_text_frame"] and not any(k in s["name"].lower() for k in ("icon", "logo", "symbol", "circle", "ellipse", "vector"))
            for s in c
        )
        is_card = False
        if has_text:
            if len(c) >= 2:
                is_card = True
            else:
                s = c[0]
                is_small_narrow = s["capacity"] < 15 or s["width_pct"] < 25 or s["height_pct"] < 10
                if not is_small_narrow:
                    is_card = True
        if is_card:
            card_clusters.append(c)
        else:
            standalone_shapes.extend(c)

    # Sort card clusters: top-to-bottom, left-to-right
    def get_cluster_center(c):
        avg_cx = sum(s["left_pct"] + s["width_pct"]/2 for s in c) / len(c)
        avg_cy = sum(s["top_pct"] + s["height_pct"]/2 for s in c) / len(c)
        return avg_cx, avg_cy

    card_clusters.sort(key=lambda c: (round(get_cluster_center(c)[1] / 15.0), get_cluster_center(c)[0]))

    card_groups = []
    for idx, c in enumerate(card_clusters):
        card_idx = idx + 1
        
        # Separate text and non-text shapes in the card.
        # Shapes with icon names are forced to non-text (icon) shapes.
        text_shapes = [s for s in c if s["has_text_frame"] and not any(k in s["name"].lower() for k in ("icon", "logo", "symbol", "circle", "ellipse", "vector"))]
        non_text_shapes = [s for s in c if not s["has_text_frame"] or any(k in s["name"].lower() for k in ("icon", "logo", "symbol", "circle", "ellipse", "vector"))]
        
        # Sort vertically within the card
        text_shapes.sort(key=lambda s: s["top_pct"])
        non_text_shapes.sort(key=lambda s: s["top_pct"])
        
        card_info = {
            "group_id": card_idx,
            "icon_shape_id": None,
            "title_shape_id": None,
            "description_shape_id": None,
            "image_shape_id": None,
            "other_shape_ids": []
        }
        
        # Assign non-text roles (icons/images)
        icon_assigned = False
        image_assigned = False
        for s in non_text_shapes:
            name = s["name"].lower()
            is_pic = s.get("is_picture", False)
            is_circle_or_icon = (
                s["width_pct"] < 12 and s["height_pct"] < 12
            ) or any(k in name for k in ("icon", "logo", "symbol", "circle", "ellipse", "oval", "vector", "graphic"))
            
            if is_circle_or_icon and not is_pic:
                if not icon_assigned:
                    s["role"] = f"card{card_idx}_icon"
                    card_info["icon_shape_id"] = s["shape_id"]
                    icon_assigned = True
                else:
                    role_name = f"icon_{s['shape_id']}"
                    s["role"] = role_name
                    card_info["other_shape_ids"].append(s["shape_id"])
            else:
                if not image_assigned:
                    s["role"] = f"card{card_idx}_image"
                    card_info["image_shape_id"] = s["shape_id"]
                    image_assigned = True
                else:
                    role_name = f"image_{s['shape_id']}"
                    s["role"] = role_name
                    card_info["other_shape_ids"].append(s["shape_id"])
                    
        # Assign text roles: Maximum one title and one description.
        # Additional text boxes must become caption_xxx, label_xxx, footer_xxx.
        if len(text_shapes) == 1:
            s = text_shapes[0]
            if s["capacity"] < 15 or s["width_pct"] < 25 or s["height_pct"] < 10:
                s["role"] = f"card{card_idx}_title"
                card_info["title_shape_id"] = s["shape_id"]
            else:
                s["role"] = f"card{card_idx}_description"
                card_info["description_shape_id"] = s["shape_id"]
                
        elif len(text_shapes) == 2:
            s1, s2 = text_shapes
            s1["role"] = f"card{card_idx}_title"
            s2["role"] = f"card{card_idx}_description"
            card_info["title_shape_id"] = s1["shape_id"]
            card_info["description_shape_id"] = s2["shape_id"]
            
        else:
            s1 = text_shapes[0]
            s2 = text_shapes[1]
            s1["role"] = f"card{card_idx}_title"
            s2["role"] = f"card{card_idx}_description"
            card_info["title_shape_id"] = s1["shape_id"]
            card_info["description_shape_id"] = s2["shape_id"]
            
            for st in text_shapes[2:]:
                if st["top_pct"] > 80:
                    role_name = f"footer_{st['shape_id']}"
                    st["role"] = role_name
                elif st["top_pct"] < 30:
                    role_name = f"label_{st['shape_id']}"
                    st["role"] = role_name
                else:
                    role_name = f"caption_{st['shape_id']}"
                    st["role"] = role_name
                card_info["other_shape_ids"].append(st["shape_id"])
                
        card_groups.append(card_info)

    # 4. Classify remaining unassigned/standalone shapes
    icon_count = 0
    label_count = 0
    caption_count = 0
    card_counter = len(card_groups) + 1
    
    for s in shapes_list:
        if s.get("role"):
            continue
            
        name = s["name"].lower()
        has_text = s["has_text_frame"]
        is_pic = s.get("is_picture", False)
        width = s["width_pct"]
        height = s["height_pct"]
        top = s["top_pct"]
        
        is_circle_or_icon = (
            width < 12 and height < 12
        ) or any(k in name for k in ("icon", "logo", "symbol", "circle", "ellipse", "oval", "vector", "graphic"))
        
        # Standalone Icon
        if (is_circle_or_icon and not is_pic) or (not has_text and is_circle_or_icon):
            icon_count += 1
            s["role"] = f"icon{icon_count}"
            continue
            
        # Standalone Image
        if is_pic or any(k in name for k in ("picture", "image", "photo", "placeholder")):
            image_count += 1
            s["role"] = f"image{image_count}"
            continue
            
        if has_text:
            is_small = s["capacity"] < 15 or width < 25 or height < 10
            if is_small:
                if top > 80:
                    s["role"] = f"footer_{s['shape_id']}"
                elif top < 30:
                    s["role"] = f"label_{s['shape_id']}"
                else:
                    s["role"] = f"caption_{s['shape_id']}"
            else:
                s["role"] = f"card{card_counter}_description"
                card_groups.append({
                    "group_id": card_counter,
                    "icon_shape_id": None,
                    "title_shape_id": None,
                    "description_shape_id": s["shape_id"],
                    "image_shape_id": None,
                    "other_shape_ids": []
                })
                card_counter += 1
        else:
            s["role"] = f"other_{s['shape_id']}"

    # Force fallback for any shape missing a role
    for s in shapes_list:
        if not s.get("role"):
            s["role"] = f"other_{s['shape_id']}"

    # 5. Before returning verify all roles are unique. No two shapes may have: body1 or card1_title simultaneously.
    seen_roles = set()
    for s in shapes_list:
        role = s["role"]
        
        # Strip body1, body2, body3 and other body roles to prevent conflicts
        if role in ("body1", "body2", "body3") or "body" in role.lower():
            role = f"description_{s['shape_id']}"
            s["role"] = role
            
        # Guarantee uniqueness
        if role in seen_roles:
            new_role = f"{role}_{s['shape_id']}"
            print(f"[TEMPLATE WARNING] Duplicate role '{role}' detected. Renaming shape {s['shape_id']} to '{new_role}' for uniqueness.")
            s["role"] = new_role
            role = new_role
            
        seen_roles.add(role)

    # 6. Add debug logging showing shape_id, name, role, shape_type, left_pct, top_pct, width_pct, height_pct
    print("\n[TEMPLATE ROLES DEBUG]")
    for s in shapes_list:
        print(f"  Shape [{s['shape_id']}] name='{s['name']}' role='{s['role']}' shape_type='{s['shape_type']}' left_pct={s['left_pct']} top_pct={s['top_pct']} width_pct={s['width_pct']} height_pct={s['height_pct']}")

    return card_groups


def _infer_layout_type(slide_idx, zones_list):
    """Infer slide layout by analyzing zones roles, positions, and counts."""
    # Check if there are card groups
    has_cards = any(z["role"].startswith("card") for z in zones_list)
    card_roles = [z["role"] for z in zones_list if z["role"].startswith("card")]
    card_group_ids = set(int(r.split("_")[0].replace("card", "")) for r in card_roles if "_" in r)
    num_cards = len(card_group_ids)
    
    if slide_idx == 0 and any(z["role"] == "title" for z in zones_list):
        return "hero"
        
    if has_cards:
        has_icons = any(z["role"].endswith("icon") for z in zones_list)
        if has_icons:
            return "icon_cards"
        if num_cards == 2:
            return "two_column"
        if num_cards == 3:
            return "three_column"
        if num_cards == 4:
            return "four_grid"
            
    # Check other layouts
    body_count = sum(1 for z in zones_list if z["role"].startswith("body"))
    image_count = sum(1 for z in zones_list if z["role"].startswith("image"))
    
    if image_count >= 2:
        return "two_image_grid"
    if body_count == 2:
        return "two_column"
    if body_count == 3:
        return "three_column"
    if body_count >= 4:
        return "four_grid"
        
    # Process / timeline keywords check
    timeline_keywords = ["timeline", "history", "milestone", "roadmap", "journey", "step", "phase"]
    for z in zones_list:
        name = z["name"].lower()
        if any(kw in name for kw in timeline_keywords):
            return "timeline"
            
    # Process check
    if any("arrow" in z["name"].lower() or "process" in z["name"].lower() for z in zones_list):
        return "process"

    return "content"
