import os
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Design System Colors
BACKGROUND = RGBColor(241, 244, 247)     # #F1F4F7
CARD_BG = RGBColor(241, 244, 247)        # #F1F4F7
TITLE_COLOR = RGBColor(18, 42, 110)      # #122A6E (Deep Navy)
BODY_COLOR = RGBColor(18, 42, 110)       # #122A6E
LIGHT_SHADOW = RGBColor(255, 255, 255)   # #FFFFFF
DARK_SHADOW = RGBColor(205, 210, 220)    # #CDD2DC

# Pop Accent Colors
PINK = RGBColor(255, 90, 160)            # #FF5AA0
ORANGE = RGBColor(255, 180, 60)          # #FFB43C
BLUE = RGBColor(70, 190, 255)            # #46BEFF
TEAL = RGBColor(40, 220, 180)            # #28DCB4

ACCENTS = [PINK, ORANGE, BLUE, TEAL]

def _get_neuro_hero_title_font_size(title_text: str) -> int:
    t_len = len(str(title_text))
    if t_len > 50:
        return 28
    elif t_len > 35:
        return 34
    elif t_len > 20:
        return 40
    else:
        return 48

def _get_neuro_body_font_size(text):
    length = len(str(text))

    if length > 180:
        return 10
    elif length > 120:
        return 11
    elif length > 80:
        return 12
    else:
        return 14

def _parse_bullet_to_metric(bullet, idx):
    import re
    bullet = bullet.strip()
    
    # Clean prefixes
    cleaned = bullet
    for prefix in ["•", "-", "*", "✓"]:
        if cleaned.startswith(prefix):
            cleaned = cleaned[len(prefix):].strip()
            
    # Try split by colon or dash
    parts = None
    if ":" in cleaned:
        parts = cleaned.split(":", 1)
    elif " - " in cleaned:
        parts = cleaned.split(" - ", 1)
        
    if parts:
        left, right = parts[0].strip(), parts[1].strip()
        # A right part is a short metric if it has <= 3 words or is <= 15 chars
        is_short_metric = len(right.split()) <= 3 or len(right) <= 15
        if is_short_metric:
            return {"label": left, "value": right}
        else:
            # Right part is a long sentence description. We treat the whole bullet as the label.
            # Try to extract a metric from the text, otherwise use fallback.
            match = re.search(r'(\+?-?\d+\.?\d*%\s*)|(\$\d+,?\d*\s*)|(\b\d+\.?\d*[KMB]?\b\s*)', cleaned)
            if match:
                val = match.group(0).strip()
            else:
                default_values = ["84.5%", "$12,400", "98.2%", "+12.4%"]
                val = default_values[idx % len(default_values)]
            return {"label": cleaned, "value": val}
            
    # Try to extract a percentage, currency, or count from a flat bullet
    match = re.search(r'(\+?-?\d+\.?\d*%\s*)|(\$\d+,?\d*\s*)|(\b\d+\.?\d*[KMB]?\b\s*)', cleaned)
    if match:
        val = match.group(0).strip()
        label = cleaned.replace(val, "").strip()
        label = re.sub(r'\s+', ' ', label).strip().strip(',').strip('.').strip(':').strip(' - ')
        if not label:
            label = f"Metric {idx + 1}"
        return {"label": label, "value": val}
    else:
        default_values = ["84.5%", "$12,400", "98.2%", "+12.4%"]
        val = default_values[idx % len(default_values)]
        return {"label": cleaned, "value": val}


def _add_neuro_card(slide, l, t, w, h, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, fill_color=CARD_BG, rotation=0):
    """Tactile extruded card with dual shadows (top-left white, bottom-right dark)"""
    offset = Inches(0.05)
    
    # 1. White highlight (top-left)
    s_light = slide.shapes.add_shape(shape_type, l - offset, t - offset, w, h)
    s_light.fill.solid()
    s_light.fill.fore_color.rgb = LIGHT_SHADOW
    s_light.line.fill.background()
    s_light.line.color.rgb = LIGHT_SHADOW
    s_light.line.width = Pt(1)
    if rotation:
        s_light.rotation = rotation
    
    # 2. Dark shadow (bottom-right)
    s_dark = slide.shapes.add_shape(shape_type, l + offset, t + offset, w, h)
    s_dark.fill.solid()
    s_dark.fill.fore_color.rgb = DARK_SHADOW
    s_dark.line.fill.background()
    s_dark.line.color.rgb = DARK_SHADOW
    s_dark.line.width = Pt(1)
    if rotation:
        s_dark.rotation = rotation
    
    # 3. Main shape
    s_main = slide.shapes.add_shape(shape_type, l, t, w, h)
    s_main.fill.solid()
    s_main.fill.fore_color.rgb = fill_color
    s_main.line.fill.background()
    s_main.line.color.rgb = fill_color
    s_main.line.width = Pt(1)
    if rotation:
        s_main.rotation = rotation
    
    return s_main


def _add_neuro_inset_card(slide, l, t, w, h, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, fill_color=CARD_BG, rotation=0):
    """Tactile sunken card with inverted shadows"""
    offset = Inches(0.05)
    
    # Light shadow (bottom-right)
    s_light = slide.shapes.add_shape(shape_type, l + offset, t + offset, w, h)
    s_light.fill.solid()
    s_light.fill.fore_color.rgb = LIGHT_SHADOW
    s_light.line.fill.background()
    s_light.line.color.rgb = LIGHT_SHADOW
    if rotation:
        s_light.rotation = rotation
    
    # Dark shadow (top-left)
    s_dark = slide.shapes.add_shape(shape_type, l - offset, t - offset, w, h)
    s_dark.fill.solid()
    s_dark.fill.fore_color.rgb = DARK_SHADOW
    s_dark.line.fill.background()
    s_dark.line.color.rgb = DARK_SHADOW
    if rotation:
        s_dark.rotation = rotation
    
    # Main shape
    s_main = slide.shapes.add_shape(shape_type, l, t, w, h)
    s_main.fill.solid()
    s_main.fill.fore_color.rgb = fill_color
    s_main.line.fill.background()
    s_main.line.color.rgb = fill_color
    if rotation:
        s_main.rotation = rotation
    
    return s_main


def _add_neuro_circle(slide, l, t, w, h, fill_color, rotation=0):
    """Draws a tactile extruded circle/pebble"""
    return _add_neuro_card(slide, l, t, w, h, shape_type=MSO_SHAPE.OVAL, fill_color=fill_color, rotation=rotation)

def _add_neuro_image_pebble(slide, l, t, w, h, image_path, shape_type=MSO_SHAPE.OVAL):
    """Draws a tactile white-bordered organic picture frame pebble"""
    from services.pptx_builder import _set_shape_picture_fill
    
    # Outer white pebble card border
    _add_neuro_card(slide, l, t, w, h, shape_type=shape_type, fill_color=LIGHT_SHADOW)
    
    # Inner shape filled with image
    margin = Inches(0.15)
    s_img = slide.shapes.add_shape(shape_type, l + margin, t + margin, w - 2*margin, h - 2*margin)
    s_img.line.fill.background()
    s_img.line.color.rgb = LIGHT_SHADOW
    
    _set_shape_picture_fill(s_img, image_path)
    return s_img

def _add_neuro_gradient_badge(slide, l, t, w, h, text, bg_color=ORANGE):
    """Draws a bright solid colored rounded square/badge for stats"""
    from services.pptx_builder import _set_center_text
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    badge.fill.solid()
    badge.fill.fore_color.rgb = bg_color
    badge.line.fill.background()
    
    _set_center_text(badge, text, 24, LIGHT_SHADOW, font_name="Times New Roman")
    if badge.text_frame.paragraphs[0].runs:
        badge.text_frame.paragraphs[0].runs[0].font.bold = False  # never bold
    return badge

def _add_neuro_title(slide, title_text, emoji=None, align=PP_ALIGN.LEFT):
    """Adds a Poppins Light huge title with accent dots below it"""
    from services.pptx_builder import _add_textbox
    full_title = f"{emoji}  {title_text}" if emoji else title_text
    
    # Huge typography, Poppins Light, never bold
    title_box = _add_textbox(
        slide,
        full_title,
        l=Inches(0.8),
        t=Inches(0.4),
        w=Inches(11.73),
        h=Inches(0.9),
        font_size=36,
        bold=False,
        color=TITLE_COLOR,
        align=align,
        font_name="Times New Roman",
        fit_text=True
    )
    
    # Draw dots
    dot_y = Inches(1.35)
    dot_size = Inches(0.08)
    
    if align == PP_ALIGN.CENTER:
        center_x = Inches(13.33) / 2
        dot1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, center_x - Inches(0.15), dot_y, dot_size, dot_size)
        dot2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, center_x, dot_y, dot_size, dot_size)
        dot3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, center_x + Inches(0.15), dot_y, dot_size, dot_size)
    else:
        dot1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), dot_y, dot_size, dot_size)
        dot2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), dot_y, dot_size, dot_size)
        dot3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.1), dot_y, dot_size, dot_size)
        
    for d, color in zip([dot1, dot2], [PINK, ORANGE]):
        d.fill.solid()
        d.fill.fore_color.rgb = color
        d.line.fill.background()
        d.line.color.rgb = color
        
    dot3.fill.solid()
    dot3.fill.fore_color.rgb = LIGHT_SHADOW
    dot3.line.color.rgb = DARK_SHADOW
    dot3.line.width = Pt(1.5)

def _add_neuro_foliage(slide):
    """Adds abstract overlapping foliage/leaves to the background matching Image 5"""
    leaf_color_1 = RGBColor(220, 238, 230)
    leaf_color_2 = RGBColor(205, 232, 220)
    
    # Upper-right leaves
    l1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.0), Inches(0.3), Inches(2.2), Inches(3.5))
    l1.fill.solid()
    l1.fill.fore_color.rgb = leaf_color_1
    l1.line.fill.background()
    
    l2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.2), Inches(1.2), Inches(1.6), Inches(2.5))
    l2.fill.solid()
    l2.fill.fore_color.rgb = leaf_color_2
    l2.line.fill.background()
    
    # Bottom-left leaves
    l3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.5), Inches(4.8), Inches(2.2), Inches(3.0))
    l3.fill.solid()
    l3.fill.fore_color.rgb = leaf_color_1
    l3.line.fill.background()


# ─── NEUMORPHIC LAYOUT TEMPLATES ───

def _add_neuro_blob_hero(prs, data, palette):

    from services.pptx_builder import (
        _add_textbox,
        _fill_bg,
        _resolve_image_to_local_path,
        _set_speaker_notes,
        _set_center_text,
    )

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    _fill_bg(slide, BACKGROUND)

    title = data.get("title") or "Untitled"

    title_size = _get_neuro_hero_title_font_size(title)

    # huge title
    _add_textbox(
        slide,
        title,
        l=Inches(0.8),
        t=Inches(1.0),
        w=Inches(4.8),
        h=Inches(2.7),
        font_size=title_size,
        bold=False,
        color=TITLE_COLOR,
        font_name="Times New Roman",
        fit_text=False
    )

    # decorative dots
    dot_y = Inches(3.8)

    for i, c in enumerate([PINK, ORANGE, BLUE]):
        d = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.85 + 0.16*i),
            dot_y,
            Inches(0.08),
            Inches(0.08)
        )
        d.fill.solid()
        d.fill.fore_color.rgb = c
        d.line.fill.background()

    # ---------- large image blob ----------

    img_source = (
        (data.get("visual_items")[0]
         if data.get("visual_items")
         else None)
        or data.get("image_keyword")
        or title
    )

    img_path = _resolve_image_to_local_path(img_source)

    blob_x = Inches(6.0)
    blob_y = Inches(0.6)

    blob_w = Inches(6.3)
    blob_h = Inches(6.1)

    if img_path and os.path.exists(img_path):

        _add_neuro_image_pebble(
            slide,
            blob_x,
            blob_y,
            blob_w,
            blob_h,
            img_path,
            shape_type=MSO_SHAPE.TEAR
        )

    else:

        _add_neuro_card(
            slide,
            blob_x,
            blob_y,
            blob_w,
            blob_h,
            shape_type=MSO_SHAPE.TEAR
        )

    bullets = data.get("bullet_points", [])

    c1_title = bullets[0] if len(bullets) > 0 else "Modern Design"
    c2_title = bullets[1] if len(bullets) > 1 else "Creative Thinking"

    # -------- floating card 1 --------

    card1_x = Inches(4.8)
    card1_y = Inches(4.0)

    _add_neuro_card(
        slide,
        card1_x,
        card1_y,
        Inches(2.3),
        Inches(1.8)
    )

    badge1 = _add_neuro_circle(
        slide,
        card1_x + Inches(0.15),
        card1_y + Inches(0.25),
        Inches(0.5),
        Inches(0.5),
        PINK
    )

    _set_center_text(
        badge1,
        "✦",
        12,
        LIGHT_SHADOW,
        font_name="Times New Roman"
    )

    _add_textbox(
        slide,
        c1_title,
        l=card1_x + Inches(0.8),
        t=card1_y + Inches(0.45),
        w=Inches(1.3),
        h=Inches(0.9),
        font_size=12,
        bold=False,
        color=TITLE_COLOR,
        font_name="Times New Roman",
        fit_text=True
    )

    # -------- floating card 2 --------

    card2_x = Inches(7.0)
    card2_y = Inches(5.0)

    _add_neuro_card(
        slide,
        card2_x,
        card2_y,
        Inches(2.4),
        Inches(1.8)
    )

    badge2 = _add_neuro_circle(
        slide,
        card2_x + Inches(0.15),
        card2_y + Inches(0.25),
        Inches(0.5),
        Inches(0.5),
        ORANGE
    )

    _set_center_text(
        badge2,
        "✦",
        12,
        LIGHT_SHADOW,
        font_name="Times New Roman"
    )

    _add_textbox(
        slide,
        c2_title,
        l=card2_x + Inches(0.8),
        t=card2_y + Inches(0.45),
        w=Inches(1.3),
        h=Inches(0.9),
        font_size=12,
        bold=False,
        color=TITLE_COLOR,
        font_name="Times New Roman",
        fit_text=True
    )

    _set_speaker_notes(slide, data)


def _add_neuro_photo_showcase(prs, data, palette):

    from services.pptx_builder import (
        _add_textbox,
        _fill_bg,
        _resolve_image_to_local_path,
        _set_speaker_notes,
        _set_center_text,
    )

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    _fill_bg(slide, BACKGROUND)

    title = data.get("title") or "Photo Showcase"

    title_size = _get_neuro_hero_title_font_size(title)

    # title
    _add_textbox(
        slide,
        title,
        l=Inches(0.8),
        t=Inches(0.8),
        w=Inches(4.8),
        h=Inches(2.3),
        font_size=title_size,
        bold=False,
        color=TITLE_COLOR,
        font_name="Times New Roman",
        fit_text=False
    )

    bullets = data.get("bullet_points", [])

    body = bullets[0] if bullets else ""

    _add_textbox(
        slide,
        body,
        l=Inches(0.8),
        t=Inches(3.2),
        w=Inches(4.4),
        h=Inches(1.5),
        font_size=_get_neuro_body_font_size(body),
        bold=False,
        color=BODY_COLOR,
        line_spacing=1.3,
        font_name="Times New Roman",
        fit_text=False
    )

    # ---------- statistics ----------
    stat1 = "388"
    stat2 = "680"

    if len(bullets) > 1:
        import re
        nums = re.findall(r"\d+", bullets[1])
        if nums:
            stat1 = nums[0]

    if len(bullets) > 2:
        import re
        nums = re.findall(r"\d+", bullets[2])
        if nums:
            stat2 = nums[0]

    # stat1
    _add_textbox(
        slide,
        stat1,
        l=Inches(0.8),
        t=Inches(4.9),
        w=Inches(2),
        h=Inches(0.7),
        font_size=36,
        bold=False,
        color=TITLE_COLOR,
        font_name="Times New Roman",
        fit_text=True
    )

    _add_textbox(
        slide,
        "Projects",
        l=Inches(0.8),
        t=Inches(5.6),
        w=Inches(2),
        h=Inches(0.4),
        font_size=11,
        bold=False,
        color=BODY_COLOR,
        font_name="Times New Roman",
        fit_text=True
    )

    # stat2
    _add_textbox(
        slide,
        stat2,
        l=Inches(3.1),
        t=Inches(4.9),
        w=Inches(2),
        h=Inches(0.7),
        font_size=36,
        bold=False,
        color=TITLE_COLOR,
        font_name="Times New Roman",
        fit_text=True
    )

    _add_textbox(
        slide,
        "Clients",
        l=Inches(3.1),
        t=Inches(5.6),
        w=Inches(2),
        h=Inches(0.4),
        font_size=11,
        bold=False,
        color=BODY_COLOR,
        font_name="Times New Roman",
        fit_text=True
    )

    # image
    img_source = (
        (data.get("visual_items")[0]
         if data.get("visual_items")
         else None)
        or data.get("image_keyword")
        or title
    )

    img_path = _resolve_image_to_local_path(img_source)

    img_x = Inches(6.2)
    img_y = Inches(0.9)

    img_w = Inches(6.1)
    img_h = Inches(5.9)

    if img_path and os.path.exists(img_path):

        _add_neuro_image_pebble(
            slide,
            img_x,
            img_y,
            img_w,
            img_h,
            img_path,
            shape_type=MSO_SHAPE.TEAR
        )

    else:

        _add_neuro_card(
            slide,
            img_x,
            img_y,
            img_w,
            img_h,
            shape_type=MSO_SHAPE.TEAR
        )

    # floating pink badge
    badge = _add_neuro_circle(
        slide,
        Inches(5.7),
        Inches(1.5),
        Inches(1.0),
        Inches(1.0),
        PINK
    )

    _set_center_text(
        badge,
        "✦",
        20,
        LIGHT_SHADOW,
        font_name="Times New Roman"
    )

    _set_speaker_notes(slide, data)


def _add_neuro_progress_counter(prs, data, palette):

    from services.pptx_builder import (
        _add_textbox,
        _fill_bg,
        _resolve_image_to_local_path,
        _set_speaker_notes,
        _set_center_text,
    )

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    _fill_bg(slide, BACKGROUND)

    title = data.get("title") or "Progress"

    title_size = _get_neuro_hero_title_font_size(title)

    # title
    _add_textbox(
        slide,
        title,
        l=Inches(0.8),
        t=Inches(0.8),
        w=Inches(4.8),
        h=Inches(2.4),
        font_size=title_size,
        bold=False,
        color=TITLE_COLOR,
        font_name="Times New Roman",
        fit_text=False
    )

    bullets = data.get("bullet_points", [])

    body = bullets[0] if bullets else ""

    _add_textbox(
        slide,
        body,
        l=Inches(0.8),
        t=Inches(3.3),
        w=Inches(3.8),
        h=Inches(1.5),
        font_size=_get_neuro_body_font_size(body),
        bold=False,
        color=BODY_COLOR,
        font_name="Times New Roman",
        fit_text=False
    )

    # pink badge
    pink_badge = _add_neuro_circle(
        slide,
        Inches(0.8),
        Inches(5.0),
        Inches(0.9),
        Inches(0.9),
        PINK
    )

    _set_center_text(
        pink_badge,
        "✦",
        18,
        LIGHT_SHADOW,
        font_name="Times New Roman"
    )

    # image
    img_source = (
        (data.get("visual_items")[0]
         if data.get("visual_items")
         else None)
        or data.get("image_keyword")
        or title
    )

    img_path = _resolve_image_to_local_path(img_source)

    blob_x = Inches(6.0)
    blob_y = Inches(0.8)

    blob_w = Inches(6.0)
    blob_h = Inches(5.8)

    if img_path and os.path.exists(img_path):

        _add_neuro_image_pebble(
            slide,
            blob_x,
            blob_y,
            blob_w,
            blob_h,
            img_path,
            shape_type=MSO_SHAPE.TEAR
        )

    else:

        _add_neuro_card(
            slide,
            blob_x,
            blob_y,
            blob_w,
            blob_h,
            shape_type=MSO_SHAPE.TEAR
        )

    # extract number

    stat_value = "36"

    if len(bullets) > 1:

        import re

        nums = re.findall(r"\d+", bullets[1])

        if nums:
            stat_value = nums[0]

    orange_x = Inches(9.3)
    orange_y = Inches(4.4)

    _add_neuro_card(
        slide,
        orange_x,
        orange_y,
        Inches(2.2),
        Inches(1.7),
        fill_color=ORANGE
    )

    _add_textbox(
        slide,
        stat_value,
        l=orange_x,
        t=orange_y + Inches(0.2),
        w=Inches(2.2),
        h=Inches(0.8),
        font_size=36,
        bold=False,
        color=LIGHT_SHADOW,
        align=PP_ALIGN.CENTER,
        font_name="Times New Roman"
    )

    _add_textbox(
        slide,
        "DAYS",
        l=orange_x,
        t=orange_y + Inches(1.0),
        w=Inches(2.2),
        h=Inches(0.4),
        font_size=10,
        bold=False,
        color=LIGHT_SHADOW,
        align=PP_ALIGN.CENTER,
        font_name="Times New Roman"
    )

    _set_speaker_notes(slide, data)


def _add_neuro_liquid_process(prs, data, palette):

    from services.pptx_builder import (
        _add_textbox,
        _fill_bg,
        _set_speaker_notes,
        _set_center_text,
    )

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    _fill_bg(slide, BACKGROUND)

    title = data.get("title") or "Workflow"

    # centered title
    _add_textbox(
        slide,
        title,
        l=Inches(2.5),
        t=Inches(0.6),
        w=Inches(8),
        h=Inches(0.8),
        font_size=36,
        bold=False,
        color=TITLE_COLOR,
        align=PP_ALIGN.CENTER,
        font_name="Times New Roman",
        fit_text=True
    )

    process_steps = data.get("process_steps")

    steps = []

    if process_steps:

        for item in process_steps[:4]:

            steps.append(
                (
                    item.get("title", ""),
                    item.get("body", "")
                )
            )

    else:

        bullets = data.get("bullet_points", [])

        for bullet in bullets[:4]:

            if ":" in bullet:

                a, b = bullet.split(":",1)

                steps.append(
                    (
                        a.strip(),
                        b.strip()
                    )
                )

            else:

                steps.append((bullet,""))

    while len(steps) < 4:

        steps.append(("Step",""))

    coords = [

        (Inches(1.3), Inches(2.0), PINK),

        (Inches(4.2), Inches(3.5), ORANGE),

        (Inches(7.2), Inches(3.6), TEAL),

        (Inches(10.0), Inches(2.2), BLUE)

    ]

    sizes = [

        Inches(1.8),

        Inches(1.9),

        Inches(1.9),

        Inches(1.8)

    ]

    for idx, (x,y,color) in enumerate(coords):

        size = sizes[idx]

        blob = _add_neuro_circle(

            slide,

            x,

            y,

            size,

            size,

            color

        )

        _set_center_text(

            blob,

            str(idx+1),

            20,

            LIGHT_SHADOW,

            font_name="Times New Roman"

        )

        title_text = steps[idx][0]

        body_text = steps[idx][1]

        tx = x - Inches(0.4)

        ty = y + size + Inches(0.2)

        _add_textbox(

            slide,

            title_text,

            l=tx,

            t=ty,

            w=Inches(2.6),

            h=Inches(0.4),

            font_size=14,

            bold=False,

            color=TITLE_COLOR,

            align=PP_ALIGN.CENTER,

            font_name="Times New Roman",

            fit_text=True

        )

        _add_textbox(

            slide,

            body_text,

            l=tx,

            t=ty + Inches(0.4),

            w=Inches(2.6),

            h=Inches(1.2),

            font_size=10,

            bold=False,

            color=BODY_COLOR,

            align=PP_ALIGN.CENTER,

            font_name="Times New Roman",

            fit_text=True

        )

    # floating connectors

    connector1 = slide.shapes.add_connector(
        3,
        Inches(2.7),
        Inches(2.8),
        Inches(4.2),
        Inches(4.1)
    )

    connector2 = slide.shapes.add_connector(
        3,
        Inches(5.8),
        Inches(4.4),
        Inches(7.1),
        Inches(4.4)
    )

    connector3 = slide.shapes.add_connector(
        3,
        Inches(8.8),
        Inches(4.0),
        Inches(10.0),
        Inches(2.9)
    )

    for c in [connector1, connector2, connector3]:

        c.line.color.rgb = TITLE_COLOR

        c.line.width = Pt(2)

    _set_speaker_notes(slide, data)


def _add_neuro_pebble_cards(prs, data, palette):
    """Pebble cards layout: 3 asymmetric organic pebble cards with badge circles on top"""
    from services.pptx_builder import (
        _add_textbox,
        _set_center_text,
        _set_speaker_notes,
        _fill_bg,
    )

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, BACKGROUND)

    title_text = data.get("title") or ""
    _add_neuro_title(slide, title_text, data.get("icon_emoji"), PP_ALIGN.LEFT)

    grid_items = data.get("grid_items")
    cols = []
    if grid_items:
        for item in grid_items[:3]:
            cols.append((item.get("title", ""), item.get("description", "")))
    else:
        bullets = data.get("bullet_points", [])
        for bullet in bullets[:3]:
            if ":" in bullet:
                h, t = bullet.split(":", 1)
                cols.append((h.strip(), t.strip()))
            else:
                cols.append((bullet, ""))

    # Asymmetric spacing coordinates (absolute positioning)
    asymmetric_cols = [
        {"x": Inches(0.8), "y": Inches(1.8), "w": Inches(3.6), "h": Inches(4.8), "badge": PINK},
        {"x": Inches(4.8), "y": Inches(2.2), "w": Inches(3.4), "h": Inches(4.4), "badge": ORANGE},
        {"x": Inches(8.6), "y": Inches(1.6), "w": Inches(3.8), "h": Inches(5.0), "badge": BLUE},
    ]

    for idx, (heading, text) in enumerate(cols[:3]):
        coord = asymmetric_cols[idx]
        cx = coord["x"]
        cy = coord["y"]
        cw = coord["w"]
        ch = coord["h"]
        b_color = coord["badge"]

        # Teardrop shaped asymmetric card
        _add_neuro_card(slide, cx, cy, cw, ch, shape_type=MSO_SHAPE.TEAR)

        # Top circle badge
        badge = _add_neuro_circle(slide, cx + Inches(0.3), cy + Inches(0.3), Inches(0.6), Inches(0.6), fill_color=b_color)
        _set_center_text(badge, f"0{idx + 1}", 12, LIGHT_SHADOW, font_name="Times New Roman")
        if badge.text_frame.paragraphs[0].runs:
            badge.text_frame.paragraphs[0].runs[0].font.bold = False

        # Heading
        _add_textbox(
            slide,
            heading,
            l=cx + Inches(0.3),
            t=cy + Inches(1.1),
            w=cw - Inches(0.6),
            h=Inches(0.5),
            font_size=18,
            bold=False,
            color=TITLE_COLOR,
            font_name="Times New Roman",
            fit_text=True
        )

        # Description
        _add_textbox(
            slide,
            text,
            l=cx + Inches(0.3),
            t=cy + Inches(1.7),
            w=cw - Inches(0.6),
            h=ch - Inches(2.0),
            font_size=12,
            color=BODY_COLOR,
            font_name="Times New Roman",
            fit_text=True
        )

    _set_speaker_notes(slide, data)


def _add_neuro_blob_image_right(prs, data, palette):

    from services.pptx_builder import (
        _add_textbox,
        _fill_bg,
        _resolve_image_to_local_path,
        _set_speaker_notes,
        _set_center_text,
    )

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    _fill_bg(slide, BACKGROUND)

    title = data.get("title") or "Comparison"

    _add_textbox(
        slide,
        title,
        l=Inches(0.8),
        t=Inches(0.9),
        w=Inches(5),
        h=Inches(1),
        font_size=40,
        bold=False,
        color=TITLE_COLOR,
        font_name="Times New Roman",
        fit_text=True
    )

    comparison_items = data.get("comparison_items")

    left_points = []
    right_points = []

    if comparison_items:

        for item in comparison_items:

            left_points.append(item.get("left",""))

            right_points.append(item.get("right",""))

    bullets = data.get("bullet_points", [])

    if not left_points:

        half = max(1, len(bullets)//2)

        left_points = bullets[:half]

        right_points = bullets[half:]

    # ---------- left floating card ----------

    card1_x = Inches(0.9)
    card1_y = Inches(2.0)

    _add_neuro_card(
        slide,
        card1_x,
        card1_y,
        Inches(3.2),
        Inches(2.0)
    )

    _add_textbox(
        slide,
        "\n".join(left_points[:3]),
        l=card1_x + Inches(0.25),
        t=card1_y + Inches(0.25),
        w=Inches(2.7),
        h=Inches(1.6),
        font_size=12,
        bold=False,
        color=BODY_COLOR,
        font_name="Times New Roman",
        fit_text=True
    )

    # ---------- second floating card ----------

    card2_x = Inches(2.5)
    card2_y = Inches(4.2)

    _add_neuro_card(
        slide,
        card2_x,
        card2_y,
        Inches(3.4),
        Inches(1.8)
    )

    _add_textbox(
        slide,
        "\n".join(right_points[:3]),
        l=card2_x + Inches(0.25),
        t=card2_y + Inches(0.25),
        w=Inches(3.0),
        h=Inches(1.4),
        font_size=12,
        bold=False,
        color=BODY_COLOR,
        font_name="Times New Roman",
        fit_text=True
    )

    # ---------- giant image blob ----------

    img_source = (
        (data.get("visual_items")[0]
         if data.get("visual_items")
         else None)
        or data.get("image_keyword")
        or title
    )

    img_path = _resolve_image_to_local_path(img_source)

    blob_x = Inches(6.3)
    blob_y = Inches(1.0)

    blob_w = Inches(5.8)
    blob_h = Inches(5.8)

    if img_path:

        _add_neuro_image_pebble(
            slide,
            blob_x,
            blob_y,
            blob_w,
            blob_h,
            img_path,
            shape_type=MSO_SHAPE.TEAR
        )

    else:

        _add_neuro_card(
            slide,
            blob_x,
            blob_y,
            blob_w,
            blob_h,
            shape_type=MSO_SHAPE.TEAR
        )

    # floating pink badge

    badge = _add_neuro_circle(

        slide,

        Inches(5.8),

        Inches(1.8),

        Inches(1),

        Inches(1),

        PINK

    )

    _set_center_text(

        badge,

        "✦",

        18,

        LIGHT_SHADOW,

        font_name="Times New Roman"

    )

    _set_speaker_notes(slide, data)


def _add_neuro_blob_image_left(prs, data, palette):

    from services.pptx_builder import (
        _add_textbox,
        _fill_bg,
        _resolve_image_to_local_path,
        _set_speaker_notes,
    )

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    _fill_bg(slide, BACKGROUND)

    title = data.get("title") or "Timeline"

    _add_textbox(
        slide,
        title,
        l=Inches(0.8),
        t=Inches(0.8),
        w=Inches(4.8),
        h=Inches(1),
        font_size=40,
        bold=False,
        color=TITLE_COLOR,
        font_name="Times New Roman",
        fit_text=True
    )

    timeline_steps = data.get("timeline_steps")

    steps = []

    if timeline_steps:

        for item in timeline_steps[:3]:

            steps.append(
                (
                    item.get("title",""),
                    item.get("description","")
                )
            )

    else:

        bullets = data.get("bullet_points", [])

        for bullet in bullets[:3]:

            if ":" in bullet:

                a,b = bullet.split(":",1)

                steps.append((a.strip(),b.strip()))

            else:

                steps.append((bullet,""))

    # image blob

    img_source = (
        (data.get("visual_items")[0]
         if data.get("visual_items")
         else None)
        or data.get("image_keyword")
        or title
    )

    img_path = _resolve_image_to_local_path(img_source)

    blob_x = Inches(0.8)
    blob_y = Inches(1.9)

    blob_w = Inches(5.0)
    blob_h = Inches(5.0)

    if img_path:

        _add_neuro_image_pebble(
            slide,
            blob_x,
            blob_y,
            blob_w,
            blob_h,
            img_path,
            shape_type=MSO_SHAPE.TEAR
        )

    else:

        _add_neuro_card(
            slide,
            blob_x,
            blob_y,
            blob_w,
            blob_h,
            shape_type=MSO_SHAPE.TEAR
        )

    positions = [

        (Inches(6.4), Inches(1.9), PINK),

        (Inches(7.2), Inches(3.4), ORANGE),

        (Inches(6.2), Inches(5.0), TEAL)

    ]

    for idx, (x,y,color) in enumerate(positions):

        if idx >= len(steps):
            break

        _add_neuro_card(
            slide,
            x,
            y,
            Inches(5.5),
            Inches(1.3)
        )

        badge = _add_neuro_circle(

            slide,

            x + Inches(0.2),

            y + Inches(0.25),

            Inches(0.7),

            Inches(0.7),

            color

        )

        badge.text_frame.text = str(idx+1)

        _add_textbox(
            slide,
            steps[idx][0],
            l=x + Inches(1.1),
            t=y + Inches(0.15),
            w=Inches(4.0),
            h=Inches(0.3),
            font_size=13,
            bold=False,
            color=TITLE_COLOR,
            font_name="Times New Roman",
            fit_text=True
        )

        _add_textbox(
            slide,
            steps[idx][1],
            l=x + Inches(1.1),
            t=y + Inches(0.55),
            w=Inches(4.0),
            h=Inches(0.65),
            font_size=10,
            bold=False,
            color=BODY_COLOR,
            font_name="Times New Roman",
            fit_text=True
        )

    _set_speaker_notes(slide, data)


def _add_neuro_glass_badges(prs, data, palette):
    """Dashboard slide matching Image 5 style: center-left card with scattered floating badges"""
    from services.pptx_builder import (
        _add_textbox,
        _fill_bg,
        _set_speaker_notes,
        _set_center_text
    )

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, BACKGROUND)

    title = data.get("title") or "Dashboard"

    _add_textbox(
        slide,
        title,
        l=Inches(0.8),
        t=Inches(0.7),
        w=Inches(6),
        h=Inches(1),
        font_size=38,
        bold=False,
        color=TITLE_COLOR,
        font_name="Times New Roman",
        fit_text=True
    )

    # Large Center-Left Card
    _add_neuro_card(slide, Inches(0.8), Inches(1.8), Inches(6.8), Inches(4.9))

    # Metric parsing logic
    metrics = data.get("dashboard_metrics")
    stats = []
    if metrics:
        for item in metrics[:4]:
            stats.append({"label": item.get("label", ""), "value": item.get("value", "")})
    else:
        bullets = data.get("bullet_points") or []
        for idx, bullet in enumerate(bullets[:4]):
            stats.append(_parse_bullet_to_metric(bullet, idx))

    colors = [PINK, ORANGE, TEAL, BLUE]
    coords = [
        (Inches(1.2), Inches(2.2)),
        (Inches(3.0), Inches(2.8)),
        (Inches(4.8), Inches(2.0)),
        (Inches(6.6), Inches(2.9))
    ]

    for idx, stat in enumerate(stats[:4]):
        x, y = coords[idx]
        badge = _add_neuro_circle(
            slide,
            x,
            y,
            Inches(1.6),
            Inches(1.6),
            colors[idx]
        )
        _set_center_text(
            badge,
            str(stat["value"]),
            18,
            LIGHT_SHADOW,
            font_name="Times New Roman"
        )
        if badge.text_frame.paragraphs[0].runs:
            badge.text_frame.paragraphs[0].runs[0].font.bold = False

        _add_textbox(
            slide,
            stat["label"],
            l=x - Inches(0.2),
            t=y + Inches(1.65),
            w=Inches(2.0),
            h=Inches(0.5),
            font_size=10,
            color=BODY_COLOR,
            align=PP_ALIGN.CENTER,
            font_name="Times New Roman",
            fit_text=True
        )

    # Right: Executive Insights Card
    right_x = Inches(8.0)
    right_y = Inches(1.8)
    right_w = Inches(4.5)
    right_h = Inches(4.9)

    _add_neuro_card(slide, right_x, right_y, right_w, right_h)

    # Sparkle circle badge on top
    sparkle = _add_neuro_circle(slide, right_x + Inches(0.3), right_y + Inches(0.3), Inches(0.6), Inches(0.6), fill_color=BLUE)
    _set_center_text(sparkle, "✦", 16, LIGHT_SHADOW, font_name="Times New Roman")

    _add_textbox(
        slide,
        "Executive Insights",
        l=right_x + Inches(1.1),
        t=right_y + Inches(0.3),
        w=right_w - Inches(1.3),
        h=Inches(0.6),
        font_size=18,
        bold=False,
        color=TITLE_COLOR,
        font_name="Times New Roman",
        fit_text=True
    )

    # Description text
    bullets = data.get("bullet_points") or []
    insight_text = data.get("dashboard_insight") or (bullets[4] if len(bullets) > 4 else "Tactile performance indicators.")
    _add_textbox(
        slide,
        insight_text,
        l=right_x + Inches(0.3),
        t=right_y + Inches(1.2),
        w=right_w - Inches(0.6),
        h=Inches(2.6),
        font_size=12,
        color=BODY_COLOR,
        font_name="Times New Roman",
        fit_text=True
    )

    # Floating gradient-colored highlight badge at bottom
    badge_val = stats[0]["value"] if stats else "Insight"
    _add_neuro_gradient_badge(slide, right_x + Inches(0.3), right_y + Inches(3.9), right_w - Inches(0.6), Inches(0.7), badge_val, bg_color=ORANGE)

    _set_speaker_notes(slide, data)
def _add_neuro_metric_badges(prs, data, palette):

    from services.pptx_builder import (
        _add_textbox,
        _fill_bg,
        _set_speaker_notes,
        _set_center_text
    )

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    _fill_bg(slide, BACKGROUND)

    title = data.get("title") or "Highlights"

    _add_textbox(
        slide,
        title,
        l=Inches(3),
        t=Inches(0.7),
        w=Inches(7),
        h=Inches(1),
        font_size=36,
        bold=False,
        color=TITLE_COLOR,
        align=PP_ALIGN.CENTER,
        font_name="Times New Roman",
        fit_text=True
    )

    items = data.get("grid_items")
    cells = []
    if items:
        for item in items[:4]:
            cells.append({"label": item.get("title", ""), "desc": item.get("description", "")})
    else:
        bullets = data.get("bullet_points") or []
        for bullet in bullets[:4]:
            label = "Detail"
            desc = bullet
            if ":" in bullet:
                parts = bullet.split(":", 1)
                label = parts[0].strip()
                desc = parts[1].strip()
            cells.append({"label": label, "desc": desc})

    coords = [
        (Inches(1.0), Inches(2.1), PINK),
        (Inches(4.2), Inches(3.7), ORANGE),
        (Inches(7.3), Inches(2.0), TEAL),
        (Inches(10.1), Inches(3.7), BLUE)
    ]

    for idx, cell in enumerate(cells[:4]):
        x, y, color = coords[idx]
        blob = _add_neuro_circle(
            slide,
            x,
            y,
            Inches(1.5),
            Inches(1.5),
            color
        )
        _set_center_text(
            blob,
            cell["label"],
            14,
            LIGHT_SHADOW,
            font_name="Times New Roman"
        )
        if blob.text_frame.paragraphs[0].runs:
            blob.text_frame.paragraphs[0].runs[0].font.bold = False

        _add_textbox(
            slide,
            cell["desc"],
            l=x - Inches(0.45),
            t=y + Inches(1.7),
            w=Inches(2.4),
            h=Inches(1.5),
            font_size=10,
            color=BODY_COLOR,
            align=PP_ALIGN.CENTER,
            font_name="Times New Roman",
            fit_text=True
        )

    _set_speaker_notes(slide, data)


def _add_neuro_leaf_background(prs, data, palette):

    from services.pptx_builder import (
        _add_textbox,
        _fill_bg,
        _set_speaker_notes
    )

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    _fill_bg(slide, BACKGROUND)

    _add_neuro_foliage(slide)

    drawio_xml = data.get("drawio_xml")
    if drawio_xml:
        try:
            from services.visual_engine import generate_drawio_diagram
            svg_path, png_path = generate_drawio_diagram(
                drawio_xml,
                topic=data.get("title") or data.get("deck_title") or "",
                disable_fallback=True
            )
            if png_path and os.path.exists(png_path):
                print(f"[PPTX] Adding PNG to neumorphism slide with aspect ratio preservation: {png_path}")
                from PIL import Image
                with Image.open(png_path) as img:
                    img_w, img_h = img.size
                
                img_ratio = img_w / img_h
                max_w_in = 12.33
                max_h_in = 6.5
                left_margin_in = 0.5
                top_margin_in = 0.5
                
                if img_ratio > (max_w_in / max_h_in):
                    w_in = max_w_in
                    h_in = max_w_in / img_ratio
                    l_in = left_margin_in
                    t_in = top_margin_in + (max_h_in - h_in) / 2.0
                else:
                    h_in = max_h_in
                    w_in = max_h_in * img_ratio
                    t_in = top_margin_in
                    l_in = left_margin_in + (max_w_in - w_in) / 2.0
                
                print(f"[PPTX] Calculated neumorphism dimensions: l={l_in:.2f}\", t={t_in:.2f}\", w={w_in:.2f}\", h={h_in:.2f}\"")
                slide.shapes.add_picture(
                    png_path,
                    Inches(l_in),
                    Inches(t_in),
                    Inches(w_in),
                    Inches(h_in)
                )
            else:
                print(f"[PPTX] Draw.io diagram render failed (PNG not found): {png_path}")
        except Exception as e:
            print(f"[PPTX] Draw.io diagram inline render failed: {e}")
        _set_speaker_notes(slide, data)
        return

    nodes_data = data.get("architecture_nodes")
    nodes_list = []
    if nodes_data:
        for idx, item in enumerate(nodes_data[:5]):
            nodes_list.append({
                "title": item.get("title") or f"Component {idx + 1}",
                "description": item.get("description") or item.get("body") or ""
            })
    else:
        bullets = data.get("bullet_points") or []
        for idx, bullet in enumerate(bullets[:5]):
            title = f"Component {idx + 1}"
            description = bullet
            if ":" in bullet:
                parts = bullet.split(":", 1)
                title = parts[0].strip()
                description = parts[1].strip()
            nodes_list.append({"title": title, "description": description})

    coords = [
        (Inches(1.2), Inches(2.0)),
        (Inches(5.0), Inches(3.0)),
        (Inches(9.0), Inches(2.1)),
        (Inches(3.0), Inches(5.0)),
        (Inches(7.2), Inches(5.2))
    ]

    for idx, node in enumerate(nodes_list[:5]):
        x, y = coords[idx]
        _add_neuro_card(
            slide,
            x,
            y,
            Inches(3.0),
            Inches(1.5)
        )

        _add_textbox(
            slide,
            node["title"],
            l=x + Inches(0.2),
            t=y + Inches(0.2),
            w=Inches(2.5),
            h=Inches(0.4),
            font_size=14,
            bold=False,
            color=TITLE_COLOR,
            font_name="Times New Roman",
            fit_text=True
        )

        _add_textbox(
            slide,
            node["description"],
            l=x + Inches(0.2),
            t=y + Inches(0.7),
            w=Inches(2.5),
            h=Inches(0.7),
            font_size=10,
            color=BODY_COLOR,
            font_name="Times New Roman",
            fit_text=True
        )

    _set_speaker_notes(slide, data)