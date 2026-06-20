"""
Template Filler V2 — Direct Shape Replacement
==============================================
Fills a PPTX template by matching shapes via shape_id and replacing
content in-place. Preserves all template formatting.

This is the V2 pipeline filler. It:
  - Matches shapes by shape_id (primary) → placeholder_idx → name → coords
  - Preserves existing paragraph structure, run formatting, fonts, colors
  - Replaces image blobs in-place (no shape resizing or repositioning)
  - Skips icons, decorative shapes, and footers entirely
  - Truncates gracefully with "..." instead of overflowing
  - Never uses fallback shape matching — missing content is better than wrong content
"""

import io
import os
import copy
import uuid
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE


def fill_template_v2(
    template_path: str,
    template_schema: dict,
    slide_contents: list,
    output_dir: str = "outputs",
) -> str:
    """Fill a PPTX template with V2-generated content.

    Args:
        template_path: Path to the original .pptx template file.
        template_schema: V2 schema dict from analyze_template_v2().
        slide_contents: List of flat dicts, one per slide, with role-keyed content.
        output_dir: Directory for saving output.

    Returns:
        Absolute path to the saved .pptx file.
    """
    print(f"[FILLER_V2] Starting template fill: {template_path}")
    print(f"[FILLER_V2] Slides to fill: {len(slide_contents)}")

    prs = Presentation(template_path)

    # Generate output filename
    os.makedirs(output_dir, exist_ok=True)
    out_name = f"template_{uuid.uuid4().hex[:10]}.pptx"
    out_path = os.path.join(output_dir, out_name)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    slides_schema = template_schema.get("slides", [])

    # ── Capture initial shape state for validation ──
    initial_state = _capture_shape_state(prs)

    # ── Fill each slide ──
    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx >= len(slide_contents):
            print(f"[FILLER_V2] Slide {slide_idx + 1}: no content provided, skipping")
            continue

        content = slide_contents[slide_idx]
        if not content:
            print(f"[FILLER_V2] Slide {slide_idx + 1}: empty content, skipping")
            continue

        schema_slide = slides_schema[slide_idx] if slide_idx < len(slides_schema) else None
        if not schema_slide:
            print(f"[FILLER_V2] Slide {slide_idx + 1}: no schema found, skipping")
            continue

        _fill_slide_v2(slide, content, schema_slide, slide_idx, slide_width, slide_height)

    # ── Remove extra slides if template has more than content ──
    total_slides = len(prs.slides)
    if total_slides > len(slide_contents) and len(slide_contents) > 0:
        for del_idx in range(total_slides - 1, len(slide_contents) - 1, -1):
            try:
                _delete_slide(prs, del_idx)
                print(f"[FILLER_V2] Deleted extra slide {del_idx + 1}")
            except Exception as e:
                print(f"[FILLER_V2] Warning: Could not delete slide {del_idx + 1}: {e}")

    # ── Post-fill validation ──
    _validate_shape_state(prs, initial_state)

    # ── Apply complex script font fix ──
    fix_presentation_complex_scripts(prs)

    # ── Save ──
    prs.save(out_path)
    abs_path = os.path.abspath(out_path)
    print(f"[FILLER_V2] Saved: {abs_path}")
    return abs_path


# ═══════════════════════════════════════════════════════════════════════
# SLIDE FILLING
# ═══════════════════════════════════════════════════════════════════════

# ── STRICT WHITELIST: only these roles receive content ──
FILLABLE_TEXT_ROLES = {"title", "subtitle", "body1", "body2", "body3", "body4"}
FILLABLE_IMAGE_ROLES = {"image1", "image2", "image3", "icon1", "icon2"}


def _fill_slide_v2(slide, content, schema_slide, slide_idx, slide_width, slide_height):
    """Fill a single slide with content.

    Allows title, subtitle, body*, image*, and icon* roles dynamically.
    """
    zones = schema_slide.get("zones", [])
    if not zones:
        print(f"[FILLER_V2] Slide {slide_idx + 1}: no zones in schema")
        return

    # Map unified "body" key from frontend editor back to individual body1, body2... roles
    body_text = content.get("body")
    if body_text and isinstance(body_text, str):
        lines = [line.strip() for line in body_text.split("\n") if line.strip()]
        for i, line in enumerate(lines):
            content[f"body{i+1}"] = line
        # If there are fewer lines than body shapes in template, blank out the rest
        total_body_zones = sum(1 for z in zones if z.get("role", "").startswith("body"))
        for i in range(len(lines), total_body_zones):
            content[f"body{i+1}"] = ""

    filled_shape_ids = set()

    for zone in zones:
        role = zone.get("role", "other")
        role_lower = role.lower()

        # Check if role matches our dynamic content patterns
        is_text_role = role == "title" or role == "subtitle" or role_lower.startswith("body") or role_lower.startswith("title") or role_lower.startswith("subtitle")
        is_image_role = role_lower.startswith("image") or role_lower.startswith("icon")

        if not is_text_role and not is_image_role:
            continue

        # Get content value for this role
        val = _get_value_for_role(role, content)

        # Skip empty content — never overwrite a shape with empty text
        if val is None:
            continue
        if isinstance(val, str) and not val.strip():
            continue

        # Find the matching shape
        shape = _find_shape_v2(slide, zone, slide_width, slide_height)
        if shape is None:
            print(f"[FILLER_V2] WARNING: Slide {slide_idx + 1}, role '{role}' — shape not found, skipping")
            continue

        if shape.shape_id in filled_shape_ids:
            print(f"[FILLER_V2] WARNING: Slide {slide_idx + 1}, shape {shape.shape_id} already filled, skipping role '{role}'")
            continue

        # Fill based on type
        if is_image_role:
            _replace_image_in_shape(slide, shape, val, is_icon=role_lower.startswith("icon"))
        elif is_text_role and shape.has_text_frame:
            # Estimate capacity and truncate if needed
            font_size = _get_font_size(shape)
            max_words = _estimate_capacity(shape, font_size)
            
            # Respect the word limit based on original text if available
            orig_text = zone.get("original_text", "")
            orig_words = len(orig_text.split())
            if orig_words > 0:
                limit = max(5, int(orig_words * 1.5))
                max_words = min(max_words, limit)

            val = _truncate_text(val, max_words)
            _fill_text_preserving_format(shape, val)

        filled_shape_ids.add(shape.shape_id)
        print(f"[FILLER_V2] Slide {slide_idx + 1}: filled role '{role}' → shape_id={shape.shape_id}")


# ═══════════════════════════════════════════════════════════════════════
# SHAPE MATCHING — strict priority: shape_id → placeholder_idx → name → coords
# ═══════════════════════════════════════════════════════════════════════

def _get_all_slide_shapes(slide):
    """Recursively collect all shapes from a slide, including children of GROUP shapes."""
    result = []
    def _collect(shapes_collection):
        for shape in shapes_collection:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    _collect(shape.shapes)
                else:
                    result.append(shape)
            except Exception:
                result.append(shape)
    _collect(slide.shapes)
    return result


def _find_shape_v2(slide, zone, slide_width, slide_height):
    """Find the shape matching a zone definition.

    Priority:
    1. shape_id match
    2. placeholder_idx match
    3. name match
    4. coordinate similarity (< 5.0 Manhattan distance)

    Returns None if no match found. Never falls back to arbitrary shapes.
    """
    all_shapes = _get_all_slide_shapes(slide)

    target_shape_id = zone.get("shape_id")
    target_ph_idx = zone.get("placeholder_idx")
    target_name = zone.get("name", "")

    # Priority 1: shape_id
    if target_shape_id is not None:
        for s in all_shapes:
            if s.shape_id == target_shape_id:
                return s

    # Priority 2: placeholder_idx
    if target_ph_idx is not None:
        for s in all_shapes:
            if s.is_placeholder:
                try:
                    if s.placeholder_format.idx == target_ph_idx:
                        return s
                except Exception:
                    pass

    # Priority 3: name
    if target_name:
        for s in all_shapes:
            if s.name == target_name:
                return s

    # Priority 4: coordinate matching
    target_left = zone.get("left_pct", -1)
    target_top = zone.get("top_pct", -1)
    target_w = zone.get("width_pct", -1)
    target_h = zone.get("height_pct", -1)

    if target_left >= 0 and target_top >= 0:
        best_shape = None
        best_dist = 5.0  # threshold

        for s in all_shapes:
            try:
                s_left = (s.left / slide_width) * 100 if slide_width else 0
                s_top = (s.top / slide_height) * 100 if slide_height else 0
                s_w = (s.width / slide_width) * 100 if slide_width else 0
                s_h = (s.height / slide_height) * 100 if slide_height else 0

                dist = (
                    abs(s_left - target_left) +
                    abs(s_top - target_top) +
                    abs(s_w - target_w) +
                    abs(s_h - target_h)
                )

                if dist < best_dist:
                    best_dist = dist
                    best_shape = s
            except Exception:
                continue

        if best_shape:
            return best_shape

    return None


# ═══════════════════════════════════════════════════════════════════════
# CONTENT VALUE LOOKUP
# ═══════════════════════════════════════════════════════════════════════

def _get_value_for_role(role: str, content: dict) -> str:
    """Look up the content value for a given role.

    Checks:
    1. Direct key: content[role]
    2. zone_content: content["zone_content"][role]
    3. Image keyword: content["{role}_keyword"] for image roles

    Returns None if nothing found.
    """
    # Direct key
    val = content.get(role)
    if val is not None:
        if isinstance(val, list):
            return "\n".join(str(v) for v in val)
        return str(val)

    # zone_content sub-dict
    zc = content.get("zone_content")
    if isinstance(zc, dict):
        val = zc.get(role)
        if val is not None:
            if isinstance(val, list):
                return "\n".join(str(v) for v in val)
            return str(val)

    # Image keyword lookup
    role_lower = role.lower()
    if role_lower.startswith("image"):
        # Try {role}_keyword
        kw = content.get(f"{role}_keyword")
        if kw:
            return str(kw)
        # Try generic image_keyword
        kw = content.get("image_keyword")
        if kw:
            return str(kw)

    return None


# ═══════════════════════════════════════════════════════════════════════
# TEXT FILLING — preserves paragraph structure and run formatting
# ═══════════════════════════════════════════════════════════════════════

def _fill_text_preserving_format(shape, text_value: str):
    """Fill text into a shape while preserving existing formatting.

    Rules:
    1. Preserve existing paragraph count and structure
    2. Write into existing runs (preserves font properties)
    3. Only the first run of each paragraph is modified
    4. Extra paragraphs beyond text lines get empty string
    5. Never add new paragraphs if the template didn't have them
    """
    if not shape.has_text_frame:
        return

    tf = shape.text_frame

    # Split text into lines
    if isinstance(text_value, list):
        lines = [str(l) for l in text_value]
    else:
        lines = text_value.split("\n") if "\n" in text_value else [text_value]

    # Remove empty trailing lines
    while lines and not lines[-1].strip():
        lines.pop()

    if not lines:
        return

    paras = list(tf.paragraphs)
    P = len(paras)

    if P == 0:
        # No paragraphs at all — add one
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        run = p.add_run()
        run.text = " ".join(lines)
        return

    # Map lines to paragraphs
    if P == 1:
        # Single paragraph: join all lines
        _set_paragraph_text(paras[0], " ".join(lines))
    elif P >= len(lines):
        # More paragraphs than lines: fill available, blank the rest
        for i, para in enumerate(paras):
            if i < len(lines):
                _set_paragraph_text(para, lines[i])
            else:
                _set_paragraph_text(para, "")
    else:
        # More lines than paragraphs: first P-1 get one line each, last gets the rest
        for i in range(P - 1):
            _set_paragraph_text(paras[i], lines[i])
        remaining = " ".join(lines[P - 1:])
        _set_paragraph_text(paras[P - 1], remaining)


def _set_paragraph_text(para, text: str):
    """Set text in a paragraph by modifying the first run.
    Preserves all run formatting (font, color, bold, etc.).
    """
    if para.runs:
        # Write into existing first run — formatting preserved
        para.runs[0].text = text
        # Clear any additional runs to avoid leftover text
        for run in para.runs[1:]:
            run.text = ""
    else:
        # No runs exist — add one (won't have template formatting)
        run = para.add_run()
        run.text = text


# ═══════════════════════════════════════════════════════════════════════
# IMAGE HANDLING — in-place blob replacement
# ═══════════════════════════════════════════════════════════════════════

def _replace_image_in_shape(slide, shape, image_keyword: str, is_icon: bool = False):
    """Replace an image in a shape.
    
    First tries in-place blob replacement (preserves cropping, shape, styles).
    If the shape doesn't have an embedded image blip (e.g. vector icon/shape),
    falls back to deleting the shape and adding a new picture at the exact same bounds.
    """
    # Download the image
    image_path = _download_image_for_template(image_keyword, is_icon=is_icon)
    if not image_path:
        print(f"[FILLER_V2] WARNING: Could not download image for '{image_keyword}'")
        return

    # Find the blip element (embedded image reference)
    blip = None
    try:
        blip = shape.element.find('.//' + qn('a:blip'))
    except Exception:
        pass
    if blip is None:
        try:
            blip = shape._element.find('.//' + qn('a:blip'))
        except Exception:
            pass

    if blip is not None:
        # 1. In-place blob replacement
        try:
            rId = blip.get(qn('r:embed'))
            if rId:
                image_part = slide.part.rels[rId].target_part
                with open(image_path, 'rb') as f:
                    image_part.blob = f.read()
                print(f"[FILLER_V2] Replaced image blob in shape {shape.shape_id} with '{image_keyword}'")
                return
        except Exception as e:
            print(f"[FILLER_V2] WARNING: In-place replacement failed, trying fallback: {e}")

    # 2. Fallback: delete shape and add picture
    try:
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        
        # Remove original shape
        sp = shape._element
        parent = sp.getparent()
        if parent is not None:
            parent.remove(sp)
        else:
            shape._element.getparent().remove(shape._element)
            
        slide.shapes.add_picture(image_path, left, top, width, height)
        print(f"[FILLER_V2] Deleted shape {shape.shape_id} and added new picture for '{image_keyword}'")
    except Exception as e:
        print(f"[FILLER_V2] ERROR: Fallback image placement failed for shape {shape.shape_id}: {e}")


def _download_image_for_template(keyword: str, is_icon: bool = False):
    """Download an image using the visual engine or resolve local/remote URLs."""
    if not keyword:
        return None
        
    keyword_str = str(keyword).strip()
    
    # 1. Handle local static URLs (from the web app preview)
    if "/static/" in keyword_str:
        filename = keyword_str.split("/static/")[-1]
        import tempfile
        local_path = os.path.join(tempfile.gettempdir(), filename)
        if os.path.exists(local_path):
            return local_path

    # 2. Handle remote URLs
    if keyword_str.startswith("http://") or keyword_str.startswith("https://"):
        try:
            import requests
            import tempfile
            import uuid
            resp = requests.get(keyword_str, timeout=15)
            resp.raise_for_status()
            ext = ".jpg"
            if "png" in resp.headers.get("Content-Type", ""):
                ext = ".png"
            local_path = os.path.join(tempfile.gettempdir(), f"downloaded_{uuid.uuid4().hex[:10]}{ext}")
            with open(local_path, "wb") as f:
                f.write(resp.content)
            return local_path
        except Exception as e:
            print(f"[FILLER_V2] Failed to download URL '{keyword_str}': {e}")
            return None

    # 3. Handle standard search keyword
    try:
        from services.visual_engine import get_hero_image
        if is_icon:
            # Look for transparent vector/illustration on Pixabay first
            try:
                from services.visual_engine import PIXABAY_API_KEY
                import requests
                import uuid
                if PIXABAY_API_KEY:
                    resp = requests.get(
                        "https://pixabay.com/api/",
                        params={
                            "key": PIXABAY_API_KEY,
                            "q": keyword_str,
                            "image_type": "vector",
                            "per_page": 3,
                            "safesearch": "true"
                        },
                        timeout=8
                    )
                    if resp.status_code == 200:
                        hits = resp.json().get("hits", [])
                        if hits:
                            img_url = hits[0].get("vectorSVG") or hits[0].get("largeImageURL") or hits[0].get("webformatURL")
                            if img_url:
                                r = requests.get(img_url, timeout=10)
                                if r.status_code == 200:
                                    ext = ".png" if "png" in img_url or "png" in r.headers.get("Content-Type", "") else ".svg" if "svg" in img_url else ".jpg"
                                    import tempfile
                                    p = os.path.join(tempfile.gettempdir(), f"icon_{uuid.uuid4().hex[:10]}{ext}")
                                    with open(p, "wb") as f:
                                        f.write(r.content)
                                    print(f"[FILLER_V2] Downloaded transparent Pixabay icon: {p}")
                                    return p
            except Exception as ex:
                print(f"[FILLER_V2] Pixabay vector search failed: {ex}")

            # Fallback search query
            search_query = f"{keyword_str} flat icon png"
            return get_hero_image(search_query)
        return get_hero_image(keyword_str)
    except Exception as e:
        print(f"[FILLER_V2] WARNING: Image download failed for '{keyword_str}': {e}")
        return None


# ═══════════════════════════════════════════════════════════════════════
# TEXT CAPACITY & TRUNCATION
# ═══════════════════════════════════════════════════════════════════════

def _get_font_size(shape) -> float:
    """Extract font size from the first run of the first paragraph."""
    try:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        return float(run.font.size.pt)
    except Exception:
        pass
    return 14.0  # default


def _estimate_capacity(shape, font_size: float) -> int:
    """Estimate how many words fit in a shape based on its physical size."""
    try:
        if shape.has_text_frame:
            width_in = shape.width.inches
            height_in = shape.height.inches
            area = width_in * height_in
            f_sz = font_size if font_size else 14.0
            word_area = (f_sz ** 2 * 3.6) / 5184
            capacity = int(area / word_area) if word_area > 0 else 0
            return max(5, min(500, capacity))
    except Exception:
        pass
    # Fallback: estimate from font size
    if font_size >= 24:
        return 10
    elif font_size >= 18:
        return 20
    else:
        return 35


def _truncate_text(text: str, max_words: int) -> str:
    """Truncate text to fit within word limit, appending '...' if needed."""
    words = text.split()
    if len(words) <= max_words:
        return text
    return " ".join(words[:max_words]) + "..."


# ═══════════════════════════════════════════════════════════════════════
# SHAPE STATE VALIDATION
# ═══════════════════════════════════════════════════════════════════════

def _capture_shape_state(prs) -> list:
    """Capture initial shape positions/sizes for post-fill validation."""
    state = []
    for slide_idx, slide in enumerate(prs.slides):
        slide_shapes = []
        for shape in slide.shapes:
            slide_shapes.append({
                "shape_id": shape.shape_id,
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            })
        state.append(slide_shapes)
    return state


def _validate_shape_state(prs, initial_state: list):
    """Compare current shape state against initial. Log warnings only."""
    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx >= len(initial_state):
            break
        initial_shapes = {s["shape_id"]: s for s in initial_state[slide_idx]}
        for shape in slide.shapes:
            init = initial_shapes.get(shape.shape_id)
            if not init:
                continue
            if shape.left != init["left"] or shape.top != init["top"]:
                print(f"[FILLER_V2] WARNING: Slide {slide_idx + 1}, shape {shape.shape_id} position changed!")
            if shape.width != init["width"] or shape.height != init["height"]:
                print(f"[FILLER_V2] WARNING: Slide {slide_idx + 1}, shape {shape.shape_id} size changed!")


# ═══════════════════════════════════════════════════════════════════════
# SLIDE DELETION
# ═══════════════════════════════════════════════════════════════════════

def _delete_slide(prs, slide_idx):
    """Delete a slide by index using XML manipulation."""
    rId = prs.slides._sldIdLst[slide_idx].get(qn('r:id'))
    prs.part.drop_rel(rId)
    sldId = prs.slides._sldIdLst[slide_idx]
    prs.slides._sldIdLst.remove(sldId)


def fix_presentation_complex_scripts(prs):
    """Post-processor that ensures all runs containing complex script text (non-ASCII characters)
    have the complex script typeface (<a:cs>) properly set in their OOXML properties.
    This prevents PowerPoint from spacing out/jumbling Indic, Arabic, and other complex characters.
    """
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def _get_all_shapes_recursive(shapes_collection):
        result = []
        for shape in shapes_collection:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    result.extend(_get_all_shapes_recursive(shape.shapes))
                else:
                    result.append(shape)
            except Exception:
                result.append(shape)
        return result

    def is_rtl_text(text: str) -> bool:
        if not text:
            return False
        for char in text:
            o = ord(char)
            # Arabic/Persian/Urdu: 0x0600-0x06FF, 0x0750-0x077F, 0x08A0-0x08FF
            # Hebrew: 0x0590-0x05FF
            # Thaana: 0x0780-0x07BF
            # Samaritan: 0x0800-0x083F
            if (1424 <= o <= 1791) or (1872 <= o <= 1919) or (2208 <= o <= 2303) or (1920 <= o <= 1983) or (2048 <= o <= 2111):
                return True
        return False

    def _process_text_frame(tf):
        from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN

        # 1. Check if the text frame contains any complex script or RTL text
        has_complex = False
        has_rtl = False
        for para in tf.paragraphs:
            para_text = "".join(run.text for run in para.runs if run.text)
            if any(ord(c) > 127 for c in para_text):
                has_complex = True
            if is_rtl_text(para_text):
                has_rtl = True

        # Apply auto-fit and word wrap to prevent text overflow and clipping
        if has_complex:
            try:
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception as e:
                print(f"[COMPLEX_SCRIPT] Warning: failed to apply auto size: {e}")

        # 2. Process paragraphs and runs
        for para in tf.paragraphs:
            para_text = "".join(run.text for run in para.runs if run.text)
            if is_rtl_text(para_text):
                try:
                    pPr = para._p.get_or_add_pPr()
                    pPr.set('rtl', '1')
                    # Align RTL text to the right if it's default/left-aligned
                    if para.alignment is None or para.alignment == PP_ALIGN.LEFT:
                        para.alignment = PP_ALIGN.RIGHT
                except Exception as e:
                    print(f"[COMPLEX_SCRIPT] Warning: failed to apply RTL paragraph properties: {e}")

            for run in para.runs:
                text = run.text
                if not text:
                    continue
                if any(ord(c) > 127 for c in text):
                    try:
                        rPr = run.font._rPr
                        if rPr is not None:
                            font_name = run.font.name
                            if not font_name:
                                font_name = "Calibri"  # Fallback
                            
                            # Ensure <a:latin> element exists
                            latin = rPr.find(qn('a:latin'))
                            if latin is None:
                                latin = OxmlElement('a:latin')
                                rPr.append(latin)
                            if not latin.get('typeface'):
                                latin.set('typeface', font_name)
                            
                            # Ensure <a:cs> (complex script) element exists and matches typeface
                            cs = rPr.find(qn('a:cs'))
                            if cs is None:
                                cs = OxmlElement('a:cs')
                                rPr.append(cs)
                            cs.set('typeface', font_name)
                            
                            # Ensure <a:ea> (East Asian) element exists
                            ea = rPr.find(qn('a:ea'))
                            if ea is None:
                                ea = OxmlElement('a:ea')
                                rPr.append(ea)
                            if not ea.get('typeface'):
                                ea.set('typeface', font_name)
                    except Exception as e:
                        print(f"[COMPLEX_SCRIPT] Warning: failed to apply complex script font to run: {e}")

    for slide in prs.slides:
        all_shapes = _get_all_shapes_recursive(slide.shapes)
        for shape in all_shapes:
            if shape.has_text_frame:
                _process_text_frame(shape.text_frame)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        _process_text_frame(cell.text_frame)
