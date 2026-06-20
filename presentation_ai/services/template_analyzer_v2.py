import io
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn

def analyze_template_v2(template_bytes: bytes, max_slides: int = 30) -> dict:
    """Analyze a PPTX template and return a V2 layout schema of its structure.
    
    This is the V2 pipeline analyzer. It extracts basic shape properties
    without any card clustering, layout inference, or decorative filtering.
    It maps shapes strictly to: title, subtitle, body1-4, image1-3, icon1-2, footer1.
    Small boxes and excess text shapes become other_<shape_id> and are never filled.
    """
    print(f"[TEMPLATE_V2] Starting V2 analysis of template ({len(template_bytes)} bytes)")
    prs = Presentation(io.BytesIO(template_bytes))

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Extract theme info
    theme_info = _extract_theme_info(prs)

    slides_info = []
    for idx, slide in enumerate(prs.slides):
        if idx >= max_slides:
            break
            
        print(f"[TEMPLATE_V2] Analyzing Slide {idx + 1}/{len(prs.slides)}")
        slide_data = _analyze_slide_v2(slide, idx, slide_width, slide_height)
        slides_info.append(slide_data)

    schema = {
        "slide_count": len(slides_info),
        "theme": theme_info,
        "slides": slides_info,
    }
    
    print("[TEMPLATE_V2] V2 Analysis completed successfully.")
    return schema


def _extract_theme_info(prs):
    theme_info = {
        "heading_font": "Calibri",
        "body_font": "Calibri",
        "colors": {
            "dk1": "#000000",
            "lt1": "#FFFFFF",
            "dk2": "#1F4E79",
            "lt2": "#D9E1F2",
            "accent1": "#4472C4",
            "accent2": "#ED7D31",
            "accent3": "#A5A5A5",
            "accent4": "#FFC000",
            "accent5": "#5B9BD5",
            "accent6": "#70AD47",
        }
    }
    try:
        if not prs.slide_masters:
            return theme_info
        master = prs.slide_masters[0]
        font_scheme = master.element.findall('.//' + qn('a:fontScheme'))
        if font_scheme:
            major = font_scheme[0].find(qn('a:majorFont'))
            if major is not None:
                latin = major.find(qn('a:latin'))
                if latin is not None and latin.get('typeface'):
                    theme_info['heading_font'] = latin.get('typeface')
            minor = font_scheme[0].find(qn('a:minorFont'))
            if minor is not None:
                latin = minor.find(qn('a:latin'))
                if latin is not None and latin.get('typeface'):
                    theme_info['body_font'] = latin.get('typeface')

        clr_scheme = master.element.findall('.//' + qn('a:clrScheme'))
        if clr_scheme:
            color_names = ['dk1', 'lt1', 'dk2', 'lt2', 'accent1', 'accent2', 'accent3', 'accent4', 'accent5', 'accent6']
            for name in color_names:
                color_el = clr_scheme[0].find(qn(f'a:{name}'))
                if color_el is not None:
                    srgb = color_el.find(qn('a:srgbClr'))
                    if srgb is not None and srgb.get('val'):
                        theme_info['colors'][name] = f'#{srgb.get("val")}'
                    else:
                        sys_clr = color_el.find(qn('a:sysClr'))
                        if sys_clr is not None and sys_clr.get('lastClr'):
                            theme_info['colors'][name] = f'#{sys_clr.get("lastClr")}'
    except Exception as e:
        print(f"[TEMPLATE_V2] Warning: Theme parsing failed: {e}")
    return theme_info


def _get_all_shapes_recursive(shapes, slide_width, slide_height, group_offset_left=0, group_offset_top=0):
    """Recursively flatten GROUP containers so child shapes are analyzed.

    Returns a list of (shape, abs_left_offset, abs_top_offset) tuples.
    For non-group shapes the offsets are simply passed through;
    for group children the group's position is accumulated as an offset.
    """
    result = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            new_offset_left = group_offset_left + shape.left
            new_offset_top = group_offset_top + shape.top
            result.extend(
                _get_all_shapes_recursive(
                    shape.shapes, slide_width, slide_height,
                    new_offset_left, new_offset_top,
                )
            )
        else:
            result.append((shape, group_offset_left, group_offset_top))
    return result


def _analyze_slide_v2(slide, idx, slide_width, slide_height):
    bg_color_hex = None
    try:
        if slide.background and slide.background.fill and slide.background.fill.type == 1:
            bg_color_hex = f"#{slide.background.fill.fore_color.rgb}"
    except Exception:
        pass

    shapes_list = []
    all_shapes = _get_all_shapes_recursive(slide.shapes, slide_width, slide_height)
    for shape_idx, (shape, offset_left, offset_top) in enumerate(all_shapes):
        shape_data = _analyze_shape_v2(shape, shape_idx, slide_width, slide_height, offset_left, offset_top)
        shapes_list.append(shape_data)

    # Assign V2 roles strictly following the instructions
    _assign_v2_roles(shapes_list)

    return {
        "slide_index": idx,
        "inferred_layout": "content", # simplified
        "zones": shapes_list,
        "column_count": 1,
        "has_image_placeholder": any(s["role"].startswith("image") for s in shapes_list),
        "background_color": bg_color_hex,
    }


def _analyze_shape_v2(shape, z_order, slide_width, slide_height, offset_left=0, offset_top=0):
    left_pct = ((shape.left + offset_left) / slide_width) * 100 if slide_width else 0
    top_pct = ((shape.top + offset_top) / slide_height) * 100 if slide_height else 0
    width_pct = (shape.width / slide_width) * 100 if slide_width else 0
    height_pct = (shape.height / slide_height) * 100 if slide_height else 0

    shape_type_name = "unknown"
    if shape.shape_type is not None:
        try:
            shape_type_name = shape.shape_type.name
        except Exception:
            shape_type_name = str(shape.shape_type)

    placeholder_idx = None
    placeholder_type = None
    if shape.is_placeholder:
        try:
            placeholder_idx = shape.placeholder_format.idx
            ph_type = shape.placeholder_format.type
            placeholder_type = ph_type.name if hasattr(ph_type, "name") else str(ph_type)
        except Exception:
            pass

    is_picture = False
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            is_picture = True
    except Exception:
        pass
    if not is_picture and shape.is_placeholder:
        try:
            if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                is_picture = True
        except Exception:
            pass

    font_name = None
    font_size = None
    font_color = None
    bold = False
    italic = False
    underline = False
    alignment = "left"

    if shape.has_text_frame:
        tf = shape.text_frame
        if tf.paragraphs:
            p = tf.paragraphs[0]
            if p.alignment is not None:
                try:
                    alignment = p.alignment.name.lower()
                except Exception:
                    pass
            for run in p.runs:
                if run.font.name:
                    font_name = run.font.name
                if run.font.size:
                    try:
                        font_size = float(run.font.size.pt)
                    except Exception:
                        pass
                if run.font.color and run.font.color.type is not None:
                    try:
                        font_color = f"#{run.font.color.rgb}"
                    except Exception:
                        pass
                bold = bool(run.font.bold)
                italic = bool(run.font.italic)
                underline = bool(run.font.underline)
                break

    capacity = 40
    if shape.has_text_frame:
        try:
            width_in = shape.width.inches
            height_in = shape.height.inches
            area = width_in * height_in
            f_sz = font_size if font_size else 14.0
            word_area = (f_sz ** 2 * 3.6) / 5184
            capacity = int(area / word_area) if word_area > 0 else 0
            capacity = max(5, min(500, capacity))
        except Exception:
            capacity = 40

    return {
        "shape_id": shape.shape_id,
        "name": shape.name or "",
        "shape_type": shape_type_name,
        "placeholder_idx": placeholder_idx,
        "placeholder_type": placeholder_type,
        "z_order": z_order,
        "left_pct": round(left_pct, 2),
        "top_pct": round(top_pct, 2),
        "width_pct": round(width_pct, 2),
        "height_pct": round(height_pct, 2),
        "has_text_frame": shape.has_text_frame,
        "is_picture": is_picture,
        "font_name": font_name,
        "font_size": font_size,
        "font_color": font_color,
        "bold": bold,
        "italic": italic,
        "underline": underline,
        "alignment": alignment,
        "capacity": capacity,
        "original_text": shape.text_frame.text.strip() if shape.has_text_frame else "",
        "role": "other",
    }


def _assign_v2_roles(shapes):
    """Assign strict V2 roles to slide shapes.

    ALLOWED ROLES ONLY:
        title, subtitle, body1 to bodyN,
        image1 to imageM, icon1 to iconK, footer1,
        other_<shape_id> (never filled).
    """
    footers = []
    large_pictures = []
    icons = []
    text_boxes = []

    for s in shapes:
        name_lower = s["name"].lower()
        stype_upper = s["shape_type"].upper()
        text_clean = s["original_text"].strip()

        # Check if text is page number, slide number, date, copyright, etc.
        is_page_num_or_date = False
        if s["has_text_frame"] and text_clean:
            is_slide_num = text_clean.isdigit() and int(text_clean) < 150 and s["top_pct"] > 85
            is_fraction_num = bool(re.match(r'^\d+/\d+$', text_clean)) and s["top_pct"] > 85
            
            is_page_num_or_date = (
                is_slide_num or
                is_fraction_num or 
                (bool(re.match(r'^\d{4}-\d{2}-\d{2}$', text_clean)) and s["top_pct"] > 85) or
                ("page" in text_clean.lower() and s["top_pct"] > 85) or
                "copyright" in text_clean.lower() or
                "confidential" in text_clean.lower()
            )

        # ── Footer: bottom-anchored small text or page numbers/dates ──
        if s["has_text_frame"] and (
            (s["top_pct"] > 88 and s["height_pct"] < 10 and s["width_pct"] < 40) or
            is_page_num_or_date
        ):
            footers.append(s)
            continue

        # ── Large pictures → image candidates ──
        if s["is_picture"] and (s["width_pct"] > 12 or s["height_pct"] > 12):
            large_pictures.append(s)
            continue

        # ── Icons / small decorative shapes (ONLY if they don't contain text) ──
        is_icon = not text_clean and (
            any(k in name_lower for k in ("icon", "logo", "vector", "symbol", "circle", "ellipse", "oval")) or
            stype_upper in ("OVAL", "ELLIPSE") or
            (s["width_pct"] < 12 and s["height_pct"] < 12)
        )
        if is_icon:
            icons.append(s)
            continue

        # ── Decorative sequential numbers (01, 02, 03, etc.) — never filled ──
        if s["has_text_frame"] and text_clean:
            is_decorative_number = bool(re.match(r'^\d{1,2}$', text_clean)) and len(text_clean) <= 2
            # Also check for numbers with leading zeros like '01', '02'
            if is_decorative_number:
                # Check context: if it's small and there are other shapes with similar numbers nearby, it's decorative
                if s["width_pct"] < 15 and s["height_pct"] < 15:
                    s["role"] = f"other_{s['shape_id']}"
                    continue

        # ── Text boxes (excluding chart/table/graph labels and empty text) ──
        if s["has_text_frame"] and text_clean and not any(k in name_lower for k in ("table", "chart", "graph")):
            text_boxes.append(s)
            continue

    # Qualified text boxes — no small box filter anymore!
    qualified_text = text_boxes

    # ── 1. Title: largest top textbox (top_pct < 35) ──
    title_shape = None
    top_textboxes = [s for s in qualified_text if s["top_pct"] < 35]
    if top_textboxes:
        # Sort by font size DESC, then by area DESC
        top_textboxes.sort(key=lambda s: (s.get("font_size") or 0, s["width_pct"] * s["height_pct"]), reverse=True)
        title_shape = top_textboxes[0]
        title_shape["role"] = "title"

    # ── 2. Subtitle: second largest textbox (top_pct < 45) ──
    subtitle_shape = None
    remaining_top = [s for s in qualified_text if s != title_shape and s["top_pct"] < 45]
    if remaining_top:
        remaining_top.sort(key=lambda s: (s.get("font_size") or 0, s["width_pct"] * s["height_pct"]), reverse=True)
        subtitle_shape = remaining_top[0]
        subtitle_shape["role"] = "subtitle"

    # ── 3. Body1-N: remaining textboxes sorted by reading order (top-to-bottom, left-to-right) ──
    bodies = [s for s in qualified_text if s != title_shape and s != subtitle_shape]
    # Reading order sort: group by top position rounded to nearest 5% to handle minor misalignments, then sort by left_pct
    bodies.sort(key=lambda s: (round(s["top_pct"] / 5) * 5, s["left_pct"]))
    for i, s in enumerate(bodies):
        s["role"] = f"body{i + 1}"

    # ── 4. Images: image1, image2, image3 ... (left-to-right) ──
    large_pictures.sort(key=lambda s: (round(s["top_pct"] / 10) * 10, s["left_pct"]))
    for i, s in enumerate(large_pictures):
        s["role"] = f"image{i + 1}"

    # ── 5. Icons: icon1, icon2 ... (left-to-right) ──
    icons.sort(key=lambda s: (round(s["top_pct"] / 10) * 10, s["left_pct"]))
    for i, s in enumerate(icons):
        s["role"] = f"icon{i + 1}"

    # ── 6. Footer: first only ──
    footers.sort(key=lambda s: s["left_pct"])
    for i, s in enumerate(footers):
        if i == 0:
            s["role"] = "footer1"
        else:
            s["role"] = f"other_{s['shape_id']}"

    # ── 7. Everything still marked "other" gets other_<id> ──
    for s in shapes:
        if s.get("role") == "other":
            s["role"] = f"other_{s['shape_id']}"
