import io
import os
import tempfile
import threading

thread_local = threading.local()
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_CONNECTOR, MSO_SHAPE

try:
    from services.visual_engine import (
        get_hero_image,
        get_icon,
        generate_diagram,
        generate_drawio_diagram,
        render_chart,
        get_illustration,
        build_table,
        download_all_images,
    )
    _VISUAL_ENGINE_AVAILABLE = True
except ImportError:
    _VISUAL_ENGINE_AVAILABLE = False
    print("[PPTX] visual_engine not available — using built-in visuals only")

from typing import Optional
from services.theme_engine import get_theme

FONT_SIZES = {
    'slide_title':    Pt(32),
    'section_heading': Pt(18),
    'hero_body':      Pt(24),
    'bullet_point':   Pt(16),
    'card_body':      Pt(13),
    'numbered_label': Pt(11),
    'footer':         Pt(9),
}

FONT_WEIGHTS = {
    'slide_title':    True,   # Bold
    'section_heading': True,  # Bold
    'hero_body':      False,  # Regular
    'bullet_point':   False,  # Regular
    'card_body':      False,  # Regular
}

FONT_NAME = 'Calibri'

def _parse_bullet_to_metric_pptx(bullet, idx):
    import re
    bullet = bullet.strip()
    
    # Clean prefixes
    cleaned = bullet
    for prefix in ["•", "-", "*", "✓"]:
        if cleaned.startswith(prefix):
            cleaned = cleaned[len(prefix):].strip()
            
    # Try split by colon or dash
    parts = None
    if ":" in cleaned:
        parts = cleaned.split(":", 1)
    elif " - " in cleaned:
        parts = cleaned.split(" - ", 1)
        
    if parts:
        left, right = parts[0].strip(), parts[1].strip()
        # A right part is a short metric if it has <= 3 words or is <= 15 chars
        is_short_metric = len(right.split()) <= 3 or len(right) <= 15
        if is_short_metric:
            return {"label": left, "value": right}
        else:
            # Right part is a long sentence description. We treat the whole bullet as the label.
            # Try to extract a metric from the text, otherwise use fallback.
            match = re.search(r'(\+?-?\d+\.?\d*%\s*)|(\$\d+,?\d*\s*)|(\b\d+\.?\d*[KMB]?\b\s*)', cleaned)
            if match:
                val = match.group(0).strip()
            else:
                default_values = ["84.5%", "$12,400", "98.2%", "+12.4%"]
                val = default_values[idx % len(default_values)]
            return {"label": cleaned, "value": val}
            
    # Try to extract a percentage, currency, or count from a flat bullet
    match = re.search(r'(\+?-?\d+\.?\d*%\s*)|(\$\d+,?\d*\s*)|(\b\d+\.?\d*[KMB]?\b\s*)', cleaned)
    if match:
        val = match.group(0).strip()
        label = cleaned.replace(val, "").strip()
        label = re.sub(r'\s+', ' ', label).strip().strip(',').strip('.').strip(':').strip(' - ')
        if not label:
            label = f"Metric {idx + 1}"
        return {"label": label, "value": val}
    else:
        default_values = ["84.5%", "$12,400", "98.2%", "+12.4%"]
        val = default_values[idx % len(default_values)]
        return {"label": cleaned, "value": val}

def _resolve_image_to_local_path(img_str: str) -> Optional[str]:
    if not img_str:
        return None
    img_str = str(img_str).strip()
    if not img_str:
        return None
        
    # If it is already a local path and exists, return it
    if os.path.exists(img_str):
        return img_str
        
    # If it is a base64 Data URL, decode it to a temp file
    if img_str.startswith("data:image/") and ";base64," in img_str:
        try:
            import base64
            import uuid
            header, base64_data = img_str.split(";base64,", 1)
            ext = ".jpg"
            if "png" in header:
                ext = ".png"
            elif "gif" in header:
                ext = ".gif"
            elif "webp" in header:
                ext = ".webp"
            img_data = base64.b64decode(base64_data)
            path = os.path.join(tempfile.gettempdir(), f"uploaded_{uuid.uuid4().hex[:10]}{ext}")
            with open(path, "wb") as f:
                f.write(img_data)
            print(f"[PPTX] Base64 Data URL decoded to local temp: {path}")
            return path
        except Exception as e:
            print(f"[PPTX] Failed to decode base64 Data URL: {e}")
            return None

    # If it is a local static URL, map it back to the temp folder
    if "/static/" in img_str:
        filename = img_str.split("/static/")[-1]
        filename = filename.split("?")[0].split("#")[0]
        local_path = os.path.join(tempfile.gettempdir(), filename)
        if os.path.exists(local_path):
            return local_path
            
    # If it is a remote web URL, download it to the temp directory
    if img_str.startswith("http://") or img_str.startswith("https://"):
        try:
            import requests
            import uuid
            resp = requests.get(img_str, timeout=10)
            resp.raise_for_status()
            ext = ".jpg"
            content_type = resp.headers.get("Content-Type", "")
            if "png" in content_type:
                ext = ".png"
            elif "webp" in content_type:
                ext = ".webp"
            path = os.path.join(tempfile.gettempdir(), f"downloaded_{uuid.uuid4().hex[:10]}{ext}")
            with open(path, "wb") as f:
                f.write(resp.content)
            print(f"[PPTX] Image URL downloaded to local temp: {path}")
            return path
        except Exception as e:
            print(f"[PPTX] Failed to download URL {img_str}: {e}")
            return None
            
    # Otherwise, it's a keyword. Fetch via visual engine if available
    if _VISUAL_ENGINE_AVAILABLE:
        try:
            return get_hero_image(img_str)
        except Exception as e:
            print(f"[PPTX] Failed to get hero image for keyword '{img_str}': {e}")
    return None

def _get_title_font_size(title_text: str) -> int:
    t_len = len(str(title_text))
    if t_len > 40:
        return 24
    elif t_len > 25:
        return 28
    else:
        return 34

def _get_hero_title_font_size(title_text: str) -> int:
    t_len = len(str(title_text))
    if t_len > 50:
        return 32
    elif t_len > 35:
        return 40
    elif t_len > 20:
        return 48
    else:
        return 56

TONE_PALETTES = {
    "professional": {
        "bg":         RGBColor(255, 255, 255),      # #ffffff (White background)
        "primary":    RGBColor(236, 196, 53),       # #ecc435 (Golden yellow)
        "secondary":  RGBColor(33, 37, 41),         # #212529 (Charcoal black)
        "surface":    RGBColor(248, 249, 250),      # #f8f9fa (Subtle card gray)
        "text":       RGBColor(33, 37, 41),         # #212529
        "text_muted": RGBColor(73, 80, 87),         # #495057
        "font_heading": "Calibri",
        "font_body":    "Calibri",
    },
    "creative": {
        "bg":         RGBColor(15, 7, 40),          # #0f0728 (Deep Indigo/Fuchsia Dark)
        "primary":    RGBColor(217, 70, 239),       # #d946ef (Fuchsia)
        "secondary":  RGBColor(225, 29, 72),        # #e11d48 (Rose)
        "surface":    RGBColor(31, 18, 65),         # rgba(31, 18, 65) (Glass Card)
        "text":       RGBColor(253, 244, 255),      # #fdf4ff (Pinkish White)
        "text_muted": RGBColor(216, 180, 254),      # #d8b4fe (Fuchsia Muted)
        "font_heading": "Segoe UI",
        "font_body":    "Segoe UI",
    },
    "lavish": {
        "bg":         RGBColor(7, 17, 31),
        "primary":    RGBColor(67, 223, 255),
        "secondary":  RGBColor(79, 125, 255),
        "surface":    RGBColor(13, 31, 55),
        "text":       RGBColor(248, 251, 255),
        "text_muted": RGBColor(182, 199, 230),
        "font_heading": "Segoe UI",
        "font_body":    "Segoe UI",
    },
    "technical": {
        "bg":         RGBColor(3, 7, 18),           # #030712 (Futuristic Black)
        "primary":    RGBColor(16, 185, 129),       # #10b981 (Emerald Green)
        "secondary":  RGBColor(6, 182, 212),        # #06b6d4 (Cyan)
        "surface":    RGBColor(17, 24, 39),         # rgba(17, 24, 39) (Tech Card)
        "text":       RGBColor(249, 250, 251),      # #f9fafb (White)
        "text_muted": RGBColor(156, 163, 175),      # #9ca3af (Gray Muted)
        "font_heading": "Segoe UI",
        "font_body":    "Segoe UI",
    },
    "academic": {
        "bg":         RGBColor(250, 250, 249),      # #fafaf9 (Warm Editorial Paper)
        "primary":    RGBColor(120, 53, 15),        # #78350f (Warm Amber/Brown)
        "secondary":  RGBColor(157, 23, 77),        # #9d174d (Pinkish Red)
        "surface":    RGBColor(245, 245, 244),      # #f5f5f4 (Light Stone Card)
        "text":       RGBColor(28, 25, 23),         # #1c1917 (Charcoal Dark Text)
        "text_muted": RGBColor(107, 102, 97),       # #6b6661 (Brown Muted)
        "font_heading": "Georgia",
        "font_body":    "Georgia",
    },
    "education": {
        "bg":         RGBColor(45, 45, 45),          # #2d2d2d (Dark Chalkboard)
        "primary":    RGBColor(245, 197, 24),        # #f5c518 (Bright Yellow - Chalk)
        "secondary":  RGBColor(229, 57, 53),         # #e53935 (Chalk Red)
        "surface":    RGBColor(55, 55, 55),          # #373737 (Chalkboard Card)
        "text":       RGBColor(255, 255, 255),       # #ffffff (White Chalk)
        "text_muted": RGBColor(200, 200, 200),       # #c8c8c8 (Light Chalk Muted)
        "font_heading": "Georgia",
        "font_body":    "Georgia",
    },
    "ecommerce": {
        "bg":         RGBColor(15, 10, 40),          # #0f0a28 (Deep Purple-Black)
        "primary":    RGBColor(255, 107, 107),       # #ff6b6b (Coral/Salmon)
        "secondary":  RGBColor(168, 85, 247),        # #a855f7 (Vivid Purple)
        "surface":    RGBColor(255, 255, 255),       # #ffffff (White floating card)
        "text":       RGBColor(30, 27, 56),          # #1e1b38 (Dark text on white card)
        "text_muted": RGBColor(107, 114, 128),       # #6b7280 (Gray muted)
        "accent":     RGBColor(251, 191, 36),        # #fbbf24 (Gold/amber accent)
        "font_heading": "Poppins",
        "font_body":    "Inter",
        # Gradient stops for background
        "grad_mid":    RGBColor(168, 85, 247),       # Purple
        "grad_end":    RGBColor(59, 130, 246),        # Blue
    },
    "neumorphism": {
        "bg":         RGBColor(241, 244, 247),      # #F1F4F7
        "primary":    RGBColor(18, 42, 110),        # #122A6E
        "secondary":  RGBColor(18, 42, 110),        # #122A6E
        "surface":    RGBColor(241, 244, 247),      # #F1F4F7
        "text":       RGBColor(18, 42, 110),        # #122A6E
        "text_muted": RGBColor(18, 42, 110),        # #122A6E
        "font_heading": "Poppins",
        "font_body":    "Poppins",
    }
}

def _hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip('#')
    r, g, b = int(hex_color[0:2],16), int(hex_color[2:4],16), int(hex_color[4:6],16)
    return RGBColor(r, g, b)

def get_pptx_theme(tone: str) -> dict:
    t_lower = tone.lower()
    if t_lower == "neumorphism":
        pal = TONE_PALETTES["neumorphism"]
        return {
            'bg_color':      pal["bg"],
            'heading_color': pal["text"],
            'body_color':    pal["text_muted"],
            'primary_color': pal["primary"],
            'card_bg':       pal["surface"],
            'heading_font':  pal["font_heading"],
            'body_font':     pal["font_body"],
            'heading_bold':  True,
            'heading_caps':  False,
        }
    if t_lower in ["ecommerce", "e-commerce", "education", "lavish"]:
        pal = TONE_PALETTES.get(t_lower, TONE_PALETTES.get("ecommerce" if "ecommerce" in t_lower else "professional"))
        return {
            'bg_color':      pal["bg"],
            'heading_color': pal["text"],
            'body_color':    pal["text_muted"],
            'primary_color': pal["primary"],
            'card_bg':       pal["surface"],
            'heading_font':  pal.get("font_heading", "Calibri"),
            'body_font':     pal.get("font_body", "Calibri"),
            'heading_bold':  True,
            'heading_caps':  False,
        }
    
    theme = get_theme(tone)
    return {
        'bg_color':      _hex_to_rgb(theme['background']),
        'heading_color': _hex_to_rgb(theme['text_heading']),
        'body_color':    _hex_to_rgb(theme['text_body']),
        'primary_color': _hex_to_rgb(theme['primary']),
        'card_bg':       _hex_to_rgb(theme['card_background']),
        'heading_font':  'Calibri',  # PPTX safe font
        'body_font':     'Calibri',
        'heading_bold':  True,
        'heading_caps':  tone.lower() == 'professional',
    }

def add_slide_title_with_emoji(slide, title: str, emoji: str, theme: dict):
    """Add emoji + title text to PPTX slide"""
    full_title = f"{emoji}  {title}" if emoji else title
    title_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.3), 
        Inches(12.5), Inches(1.2)
    )
    tf = title_box.text_frame
    tf.word_wrap = False
    
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = full_title
    run.font.size = Pt(32)
    run.font.bold = True
    
    # Resolve color safely from various theme structures
    color_val = RGBColor(0, 0, 0)
    if 'text_heading' in theme:
        color_val = _hex_to_rgb(theme['text_heading'])
    elif 'heading_color' in theme:
        color_val = theme['heading_color']
    elif 'text' in theme:
        if isinstance(theme['text'], str):
            color_val = _hex_to_rgb(theme['text'])
        else:
            color_val = theme['text']
            
    run.font.color.rgb = color_val
    run.font.name = 'Segoe UI Emoji'
    return title_box

DEFAULT_PALETTE = TONE_PALETTES["professional"]


def build_pptx(deck: dict, topic: str, template_bytes: bytes = None, tone: str = "Professional", output_dir="outputs") -> str:
    import time
    total_start = time.time()
    
    # 1. Sanitize topic to create a dynamic filename
    safe_name = "".join(c for c in topic if c.isalnum() or c in (' ', '-', '_')).strip()
    safe_name = safe_name[:50].replace(" ", "_")
    if not safe_name:
        safe_name = "presentation"

    os.makedirs(output_dir, exist_ok=True)
    filename = f"{output_dir}/{safe_name}.pptx"
    
    print(f"[TIMING] Starting PPTX build for '{topic}'")

    # 2. Download all images in parallel if visual engine is available
    slides = deck.get("slides", [])
    image_map = {}
    if _VISUAL_ENGINE_AVAILABLE:
        try:
            image_map = download_all_images(slides)
        except Exception as e:
            print(f"[PPTX] Parallel image download failed: {e}, falling back to sequential")

    # 3. Select color palette based on tone
    theme_info = get_pptx_theme(tone)
    palette = {
        "bg":         theme_info['bg_color'],
        "text":       theme_info['heading_color'],
        "text_muted": theme_info['body_color'],
        "primary":    theme_info['primary_color'],
        "surface":    theme_info['card_bg'],
        "font_heading": theme_info['heading_font'],
        "font_body":    theme_info['body_font'],
        "heading_bold": theme_info['heading_bold'],
        "heading_caps": theme_info['heading_caps'],
    }
    
    t_lower = tone.lower()
    if t_lower in TONE_PALETTES:
        for k, v in TONE_PALETTES[t_lower].items():
            if k not in palette:
                palette[k] = v
                
    if t_lower == "academic":
        palette["swot_colors"] = {
            "S": RGBColor(45, 122, 79),
            "W": RGBColor(232, 245, 238),
            "O": RGBColor(72, 187, 120),
            "T": RGBColor(31, 92, 58),
        }
        
    thread_local.palette = palette

    # 3. Load or create presentation
    if template_bytes:
        prs = Presentation(io.BytesIO(template_bytes))
        template_slide_count = len(prs.slides)
    else:
        prs = Presentation()
        template_slide_count = 0

    # Force 16:9 widescreen slide size
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 4. Generate slides
    for slide_index, slide_data in enumerate(slides):
        # Use pre-downloaded image if available
        if slide_index in image_map and image_map[slide_index]:
            slide_data["_preloaded_image"] = image_map[slide_index]
        stype = slide_data.get("slide_type", "content")
        
        # Add deck_title key context if not present
        if "deck_title" not in slide_data and "deck_title" in deck:
            slide_data["deck_title"] = deck["deck_title"]

        # If heading_caps is active, modify the slide title to be uppercase
        if palette.get("heading_caps", False) and "title" in slide_data:
            slide_data["title"] = str(slide_data["title"]).upper()

        # Prepend emoji to slide title
        emoji = slide_data.get("icon_emoji", "")
        if emoji and "title" in slide_data:
            slide_data["title"] = f"{emoji}  {slide_data['title']}"

        if template_bytes:
            if slide_index < template_slide_count:
                slide = prs.slides[slide_index]
                _populate_template_slide(slide, slide_data, palette)
                _set_speaker_notes(slide, slide_data)
                continue

            layout = _find_layout(prs, stype)
            if layout:
                slide = prs.slides.add_slide(layout)
                _populate_template_slide(slide, slide_data, palette)
                _set_speaker_notes(slide, slide_data)
                continue

            _build_slide_from_scratch(prs, slide_data, stype, palette, tone)
        else:
            _build_slide_from_scratch(prs, slide_data, stype, palette, tone)

    if template_bytes:
        for idx in range(len(prs.slides) - 1, len(slides) - 1, -1):
            _delete_slide(prs, idx)

    # Apply complex script formatting fix
    try:
        fix_presentation_complex_scripts(prs)
    except Exception as e:
        print(f"[PPTX] Warning: failed to apply complex script fix: {e}")

    prs.save(filename)
    
    total_elapsed = time.time() - total_start
    print(f"[TIMING] PPTX assembly completed in {total_elapsed:.2f}s")
    print(f"[TIMING] Total PPTX build time: {total_elapsed:.2f}s")
    
    return filename


def _build_slide_from_scratch(prs, slide_data, stype, palette, tone=None):
    layout_id = slide_data.get("layout_id")
    is_ecommerce = palette.get("grad_start") is not None
    
    if tone is None:
        # Detect tone via caller stack frame search
        import inspect
        tone = "professional"
        try:
            frame = inspect.currentframe().f_back
            while frame:
                if 'tone' in frame.f_locals:
                    tone = frame.f_locals['tone']
                    break
                frame = frame.f_back
        except Exception:
            pass
        finally:
            del frame

    if tone and tone.lower() == "professional":
        from services.professional_layouts import (
            _add_prof_hero_slide,
            _add_prof_1_column_slide,
            _add_prof_2_column_slide,
            _add_prof_3_column_slide,
            _add_prof_4_grid_slide,
            _add_prof_process_slide,
            _add_prof_dashboard_slide,
            _add_prof_comparison_slide,
            _add_prof_timeline_slide,
            _add_prof_architecture_slide,
        )
        if layout_id == "1-column":
            _add_prof_1_column_slide(prs, slide_data, palette)
        elif layout_id == "2-column":
            _add_prof_2_column_slide(prs, slide_data, palette)
        elif layout_id == "3-column":
            _add_prof_3_column_slide(prs, slide_data, palette)
        elif layout_id == "4-grid":
            _add_prof_4_grid_slide(prs, slide_data, palette)
        elif layout_id == "hero":
            _add_prof_hero_slide(prs, slide_data, palette)
        elif layout_id == "dashboard":
            _add_prof_dashboard_slide(prs, slide_data, palette)
        elif layout_id == "comparison":
            _add_prof_comparison_slide(prs, slide_data, palette)
        elif layout_id == "process":
            _add_prof_process_slide(prs, slide_data, palette)
        elif layout_id == "timeline":
            _add_prof_timeline_slide(prs, slide_data, palette)
        elif layout_id == "architecture":
            _add_prof_architecture_slide(prs, slide_data, palette)
        else:
            if stype == "title":
                _add_prof_hero_slide(prs, slide_data, palette)
            else:
                _add_prof_1_column_slide(prs, slide_data, palette)
        return

    if tone and tone.lower() == "neumorphism":
        from services.neumorphism_layouts import (
            _add_neuro_blob_hero,
            _add_neuro_photo_showcase,
            _add_neuro_progress_counter,
            _add_neuro_liquid_process,
            _add_neuro_pebble_cards,
            _add_neuro_blob_image_right,
            _add_neuro_blob_image_left,
            _add_neuro_glass_badges,
            _add_neuro_metric_badges,
            _add_neuro_leaf_background,
        )
        if layout_id == "1-column":
            _add_neuro_photo_showcase(prs, slide_data, palette)
        elif layout_id == "2-column":
            _add_neuro_progress_counter(prs, slide_data, palette)
        elif layout_id == "3-column":
            _add_neuro_pebble_cards(prs, slide_data, palette)
        elif layout_id == "4-grid":
            _add_neuro_metric_badges(prs, slide_data, palette)
        elif layout_id == "hero":
            _add_neuro_blob_hero(prs, slide_data, palette)
        elif layout_id == "dashboard":
            _add_neuro_glass_badges(prs, slide_data, palette)
        elif layout_id == "comparison":
            _add_neuro_blob_image_right(prs, slide_data, palette)
        elif layout_id == "process":
            _add_neuro_liquid_process(prs, slide_data, palette)
        elif layout_id == "timeline":
            _add_neuro_blob_image_left(prs, slide_data, palette)
        elif layout_id == "architecture":
            _add_neuro_leaf_background(prs, slide_data, palette)
        else:
            if stype == "title":
                _add_neuro_blob_hero(prs, slide_data, palette)
            else:
                _add_neuro_photo_showcase(prs, slide_data, palette)
        return

    if layout_id:
        if layout_id == "1-column":
            _add_1_column_slide(prs, slide_data, palette)
        elif layout_id == "2-column":
            _add_2_column_slide(prs, slide_data, palette)
        elif layout_id == "3-column":
            _add_3_column_slide(prs, slide_data, palette)
        elif layout_id == "4-grid":
            _add_4_grid_slide(prs, slide_data, palette)
        elif layout_id == "hero":
            _add_hero_slide(prs, slide_data, palette)
        elif layout_id == "dashboard":
            _add_dashboard_slide(prs, slide_data, palette)
        elif layout_id == "architecture":
            _add_architecture_slide(prs, slide_data, palette)
        elif layout_id == "timeline":
            _add_timeline_slide(prs, slide_data, palette)
        elif layout_id == "comparison":
            _add_comparison_slide_layout(prs, slide_data, palette)
        elif layout_id == "process":
            _add_process_slide_layout(prs, slide_data, palette)
        else:
            _add_fallback_slide_from_scratch(prs, slide_data, stype, palette)
    else:
        _add_fallback_slide_from_scratch(prs, slide_data, stype, palette)
    
    # E-commerce post-processing: override BG with gradient + add decorative shapes
    # (hero slide already handles its own gradient/decorations)
    if is_ecommerce and layout_id != "hero":
        slide = prs.slides[-1]  # Last added slide
        _fill_bg(slide, palette["bg"], palette=palette)
        _add_ecommerce_decorative_shapes(slide, palette)


def _add_1_column_slide(prs, data, palette):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, palette["bg"])
    
    # Title
    title_text = data.get("title") or ""
    title_color = RGBColor(255, 255, 255) if (palette.get("grad_start") is not None) else palette["text"]
    _add_textbox(slide, title_text, l=Inches(0.8), t=Inches(0.6), w=Inches(11.73), h=Inches(1.0),
                 font_size=FONT_SIZES['slide_title'], bold=FONT_WEIGHTS['slide_title'], color=title_color)
                 
    # ── Visual area: try illustration or hero image ──
    illustration_embedded = False
    if _VISUAL_ENGINE_AVAILABLE:
        try:
            img_source = (data.get("visual_items")[0] if data.get("visual_items") else None) or data.get("image_keyword") or data.get("title") or ""
            illust_path = get_illustration(img_source) if not (img_source.startswith("http://") or img_source.startswith("https://") or os.path.exists(img_source)) else None
            if not illust_path:
                illust_path = _resolve_image_to_local_path(img_source)
            if illust_path and os.path.exists(illust_path):
                slide.shapes.add_picture(
                    illust_path,
                    Inches(0.8), Inches(1.8), Inches(11.73), Inches(2.4)
                )
                illustration_embedded = True
        except Exception as e:
            print(f"[PPTX] 1-column illustration embed failed: {e}")

    if not illustration_embedded:
        # Fallback: card surface with label
        has_image = data.get("visual_type") == "image" and data.get("visual_items")
        img_text = data.get("visual_items")[0] if has_image else "Featured Image Area"
        _add_rect_styled(slide, l=Inches(0.8), t=Inches(1.8), w=Inches(11.73), h=Inches(2.4),
                         fill_color=palette["surface"], line_color=palette["surface"], rounded=True)
        _add_textbox(slide, img_text, l=Inches(0.8), t=Inches(2.7), w=Inches(11.73), h=Inches(0.6),
                     font_size=FONT_SIZES['bullet_point'], color=palette["text_muted"], align=PP_ALIGN.CENTER)
                 
    # Headline and body lines below
    card_h = Inches(2.6)
    if palette.get("grad_start") is not None:
        _add_rect_styled(slide, l=Inches(0.8), t=Inches(4.5), w=Inches(11.73), h=card_h,
                         fill_color=palette["surface"], line_color=palette["surface"], rounded=True)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.73), card_h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    
    font_body = FONT_NAME
    
    headline = data.get("headline") or ""
    if headline == title_text:
        headline = ""
    bullets = data.get("bullet_points", [])
    
    is_first = True
    if headline:
        p = tf.paragraphs[0]
        p.text = headline
        p.space_after = Pt(10)
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = font_body
        run.font.size = FONT_SIZES['section_heading']
        run.font.bold = FONT_WEIGHTS['section_heading']
        run.font.color.rgb = palette["text"]
        is_first = False
        
    _populate_bullets_in_text_frame(tf, bullets, font_body, palette, is_first=is_first)
        
    _set_speaker_notes(slide, data)


def _add_2_column_slide(prs, data, palette):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, palette["bg"])
    
    title_text = data.get("title") or ""
    title_color = RGBColor(255, 255, 255) if (palette.get("grad_start") is not None) else palette["text"]
    _add_textbox(slide, title_text, l=Inches(0.8), t=Inches(0.6), w=Inches(11.73), h=Inches(1.0),
                 font_size=FONT_SIZES['slide_title'], bold=FONT_WEIGHTS['slide_title'], color=title_color)
                 
    left_headline = data.get("left_headline") or "Key Points"
    left_text = data.get("left_text") or ""
    bullets = data.get("bullet_points", [])
    
    # Left column: text/bullets styled hierarchy
    _add_rect_styled(slide, l=Inches(0.8), t=Inches(1.8), w=Inches(5.6), h=Inches(4.9),
                     fill_color=palette["surface"], line_color=palette["surface"], rounded=True)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    
    font_body = FONT_NAME
    
    is_first = True
    if left_headline:
        p = tf.paragraphs[0]
        p.text = left_headline
        p.space_after = Pt(10)
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = font_body
        run.font.size = FONT_SIZES['section_heading']
        run.font.bold = FONT_WEIGHTS['section_heading']
        run.font.color.rgb = palette["text"]
        is_first = False
        
    if left_text:
        p = tf.paragraphs[0] if is_first else tf.add_paragraph()
        p.text = left_text
        p.space_after = Pt(10)
        p.line_spacing = 1.25
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = font_body
        run.font.size = FONT_SIZES['hero_body']
        run.font.bold = FONT_WEIGHTS['hero_body']
        run.font.color.rgb = palette["text_muted"]
        is_first = False
        
    _populate_bullets_in_text_frame(tf, bullets, font_body, palette, is_first=is_first)
                 
    # Right column: visual content
    visual_type = data.get("visual_type", "none")
    visual_items = data.get("visual_items", [])
    
    # Right card rect
    _add_rect_styled(slide, l=Inches(6.8), t=Inches(1.8), w=Inches(5.7), h=Inches(4.9),
                     fill_color=palette["surface"], line_color=palette["surface"], rounded=True)
                     
    _add_visual_in_rect(slide, visual_type, visual_items, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.9), palette, data)
    _set_speaker_notes(slide, data)


def _add_3_column_slide(prs, data, palette):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, palette["bg"])
    
    title_text = data.get("title") or ""
    title_color = RGBColor(255, 255, 255) if (palette.get("grad_start") is not None) else palette["text"]
    _add_textbox(slide, title_text, l=Inches(0.8), t=Inches(0.6), w=Inches(11.73), h=Inches(1.0),
                 font_size=FONT_SIZES['slide_title'], bold=FONT_WEIGHTS['slide_title'], color=title_color)
                 
    bullets = data.get("bullet_points", [])
    icon_keywords = data.get("icon_keywords", [])
    grid_items = data.get("grid_items")
    cols = []
    if grid_items:
        for idx, item in enumerate(grid_items[:3]):
            cols.append((item.get("title") or f"Focus Area {idx + 1}", item.get("description") or ""))
    else:
        for idx, bullet in enumerate(bullets[:3]):
            if ":" in bullet:
                h, t = bullet.split(":", 1)
                cols.append((h.strip(), t.strip()))
            elif " - " in bullet:
                h, t = bullet.split(" - ", 1)
                cols.append((h.strip(), t.strip()))
            else:
                cols.append((f"Focus Area {idx + 1}", bullet.strip()))
            
    col_w = Inches(4.07)
    gap = Inches(0.15)
    y = Inches(1.3)
    h = Inches(5.3)
    
    zone_content = data.get("zone_content") or {}
    
    for idx, (heading, text) in enumerate(cols):
        x = Inches(0.4) + idx * (col_w + gap)
        _add_rect_styled(slide, l=x, t=y, w=col_w, h=h,
                         fill_color=palette["surface"], line_color=palette["surface"], rounded=True)

        # ── Try SVG icon from visual_engine ──
        icon_embedded = False
        if _VISUAL_ENGINE_AVAILABLE and icon_keywords:
            try:
                kw = icon_keywords[idx] if idx < len(icon_keywords) else heading
                svg_str = get_icon(kw)
                if svg_str:
                    svg_path = os.path.join(tempfile.gettempdir(), f"icon_{idx}_{id(slide)}.svg")
                    with open(svg_path, "w", encoding="utf-8") as f:
                        f.write(svg_str)
                    try:
                        slide.shapes.add_picture(
                            svg_path,
                            x + Inches(0.15), y + Inches(0.15), Inches(0.4), Inches(0.4)
                        )
                        icon_embedded = True
                    except Exception:
                        pass  # SVG not supported as picture in python-pptx, fallback below
            except Exception as e:
                print(f"[PPTX] 3-col icon embed failed: {e}")

        if not icon_embedded:
            bubble = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.15), Inches(0.4), Inches(0.4))
            _style_shape(bubble, palette["primary"], palette["primary"])
            _set_center_text(bubble, f"0{idx + 1}", 11, RGBColor(255, 255, 255), font_name=FONT_NAME)
        
        _add_textbox(slide, heading, l=x + Inches(0.15), t=y + Inches(0.6), w=col_w - Inches(0.3), h=Inches(0.25),
                                   font_size=FONT_SIZES['section_heading'], bold=FONT_WEIGHTS['section_heading'], color=palette["text"], font_name=FONT_NAME)
        
        # Limit text height to prevent overlap with the circular image at the bottom
        _add_textbox(slide, text, l=x + Inches(0.15), t=y + Inches(0.9), w=col_w - Inches(0.3), h=Inches(0.3),
                                 font_size=FONT_SIZES['card_body'], color=palette["text_muted"], font_name=FONT_NAME)

        # ── Circle Image at the bottom of the card ──
        img_url = zone_content.get(f"col_{idx + 1}_image") or zone_content.get(f"col_{idx + 1}_img")
        img_path = None
        if img_url:
            img_path = _resolve_image_to_local_path(img_url)
        if not img_path:
            # Fallback: search using column heading or fallback keyword
            img_path = _resolve_image_to_local_path(heading)

        img_w = Inches(2.2)
        img_h = Inches(2.2)
        img_x = x + (col_w - img_w) / 2
        img_y = y + h - img_h - Inches(0.4)
        
        circle_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, img_x, img_y, img_w, img_h)
        if img_path and os.path.exists(img_path):
            success = _set_shape_picture_fill(circle_shape, img_path)
            if not success:
                _style_shape(circle_shape, palette.get("secondary", palette["primary"]), palette.get("secondary", palette["primary"]))
        else:
            _style_shape(circle_shape, palette.get("secondary", palette["primary"]), palette.get("secondary", palette["primary"]))
                                 
    _set_speaker_notes(slide, data)


def _add_4_grid_slide(prs, data, palette):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, palette["bg"])
    
    title_text = data.get("title") or ""
    title_color = RGBColor(255, 255, 255) if (palette.get("grad_start") is not None) else palette["text"]
    _add_textbox(slide, title_text, l=Inches(0.8), t=Inches(0.6), w=Inches(11.73), h=Inches(1.0),
                 font_size=FONT_SIZES['slide_title'], bold=FONT_WEIGHTS['slide_title'], color=title_color)
                 
    bullets = data.get("bullet_points", [])
    grid_items = data.get("grid_items")
    cells = []
    if grid_items:
        for idx, item in enumerate(grid_items[:4]):
            cells.append({
                "label": item.get("title") or f"Metric {idx + 1}",
                "value": item.get("description") or "—",
                "change": item.get("change") or "",
                "trend": item.get("trend") or "up"
            })
    else:
        for idx, bullet in enumerate(bullets[:4]):
            label = f"Metric {idx + 1}"
            value = "—"
            change = ""
            trend = "up"
            if ":" in bullet:
                parts = bullet.split(":", 1)
                label = parts[0].strip()
                rest = parts[1].strip()
                import re
                match = re.search(r'\(([-+]\d+(?:\.\d+)?%?)\)', rest)
                if match:
                    change = match.group(1)
                    value = rest.replace(match.group(0), "").strip()
                else:
                    bare_match = re.search(r'\b([-+]\d+(?:\.\d+)?%?)\b', rest)
                    if bare_match:
                        change = bare_match.group(1)
                        value = rest.replace(bare_match.group(0), "").strip()
                    else:
                        value = rest
                if change and change.startswith("-"):
                     trend = "down"
            else:
                value = bullet.strip()
            cells.append({"label": label, "value": value, "change": change, "trend": trend})
            
    col_w = Inches(5.61)
    row_h = Inches(2.25)
    gap_x = Inches(0.5)
    gap_y = Inches(0.4)
    
    start_x = Inches(0.8)
    start_y = Inches(1.9)
    
    for idx, cell in enumerate(cells):
        col = idx % 2
        row = idx // 2
        x = start_x + col * (col_w + gap_x)
        y = start_y + row * (row_h + gap_y)
        
        _add_rect_styled(slide, l=x, t=y, w=col_w, h=row_h,
                         fill_color=palette["surface"], line_color=palette["surface"], rounded=True)
                         
        val_str = str(cell["value"])
        is_long = len(val_str.split()) > 3 or len(val_str) > 15
        
        if is_long:
            # Case A: Title + Description
            # Title on top
            _add_textbox(slide, cell["label"], l=x + Inches(0.25), t=y + Inches(0.2), w=col_w - Inches(0.5), h=Inches(0.5),
                         font_size=16, bold=True, color=palette["text"], font_name=FONT_NAME)
            # Description below
            _add_textbox(slide, val_str, l=x + Inches(0.25), t=y + Inches(0.75), w=col_w - Inches(0.5), h=Inches(1.3),
                         font_size=13, bold=False, color=palette["text_muted"], font_name=FONT_NAME, fit_text=True)
        else:
            # Case B: Metric Value + Muted Label
            # Muted label on top
            _add_textbox(slide, cell["label"], l=x + Inches(0.25), t=y + Inches(0.2), w=col_w - Inches(2.0), h=Inches(0.5),
                         font_size=11, bold=False, color=palette["text_muted"], font_name=FONT_NAME)
            # Large metric below
            _add_textbox(slide, val_str, l=x + Inches(0.25), t=y + Inches(0.9), w=col_w - Inches(0.5), h=Inches(1.0),
                         font_size=32, bold=True, color=palette["text"], font_name=FONT_NAME)
                         
        if cell["change"] and not is_long:
            badge_text = f"{'▲' if cell['trend'] == 'up' else '▼'} {cell['change']}"
            if cell['trend'] == 'up':
                fill_color = RGBColor(220, 252, 231)
                text_color = RGBColor(22, 163, 74)
            else:
                fill_color = RGBColor(254, 226, 226)
                text_color = RGBColor(220, 38, 38)
            badge_rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + col_w - Inches(1.5), y + Inches(0.2), Inches(1.25), Inches(0.45))
            _style_shape(badge_rect, fill_color, fill_color)
            _set_center_text(badge_rect, badge_text, 11, text_color, font_name=FONT_NAME)
            
    _set_speaker_notes(slide, data)


def _add_hero_slide(prs, data, palette):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    
    is_ecommerce = palette.get("grad_start") is not None
    
    if is_ecommerce:
        # E-commerce: gradient background + decorative shapes + white card
        _fill_bg(slide, palette["bg"], palette=palette)
        _add_ecommerce_decorative_shapes(slide, palette)
        
        # Large white rounded floating card (center)
        card_l, card_t = Inches(2.0), Inches(1.2)
        card_w, card_h = Inches(9.33), Inches(5.1)
        _add_rect_styled(slide, l=card_l, t=card_t, w=card_w, h=card_h,
                         fill_color=palette["surface"], line_color=palette["surface"], rounded=True)
        
        title_text = data.get("title") or data.get("deck_title") or "Untitled Section"
        subtitle_text = ""
        if data.get("bullet_points"):
            subtitle_text = data["bullet_points"][0]
        
        # Title on white card (dark text)
        _add_textbox(slide, title_text, l=Inches(2.8), t=Inches(2.4), w=Inches(7.73), h=Inches(1.8),
                     font_size=_get_hero_title_font_size(title_text), bold=True,
                     color=palette["text"], align=PP_ALIGN.CENTER, line_spacing=1.1, font_name=FONT_NAME)
        
        # Accent divider bar (coral/purple gradient feel)
        _add_rect_styled(slide, l=Inches(5.66), t=Inches(4.4), w=Inches(2.0), h=Inches(0.06),
                         fill_color=palette["primary"], line_color=palette["primary"])
        
        if subtitle_text:
            _add_textbox(slide, subtitle_text, l=Inches(3.2), t=Inches(4.7), w=Inches(6.93), h=Inches(1.2),
                         font_size=FONT_SIZES['hero_body'], color=palette["text_muted"], align=PP_ALIGN.CENTER, line_spacing=1.25, font_name=FONT_NAME)
    else:
        # Standard hero slide for all other tones
        _fill_bg(slide, palette["bg"])
        
        # ── Try embedding a hero image from visual_engine ──
        if _VISUAL_ENGINE_AVAILABLE:
            try:
                img_source = (data.get("visual_items")[0] if data.get("visual_items") else None) or data.get("image_keyword") or data.get("title") or "presentation"
                hero_path = _resolve_image_to_local_path(img_source)
                if hero_path and os.path.exists(hero_path):
                    slide.shapes.add_picture(
                        hero_path,
                        Inches(0), Inches(0),
                        prs.slide_width, prs.slide_height
                    )
                    print(f"[PPTX] Hero image embedded: {hero_path}")
            except Exception as e:
                print(f"[PPTX] Hero image embed failed: {e}")

        title_text = data.get("title") or data.get("deck_title") or "Untitled Section"
        subtitle_text = ""
        if data.get("bullet_points"):
            subtitle_text = data["bullet_points"][0]
            
        _add_textbox(slide, title_text, l=Inches(1.2), t=Inches(2.0), w=Inches(10.93), h=Inches(1.8),
                     font_size=_get_hero_title_font_size(title_text), bold=True, color=palette["text"], align=PP_ALIGN.CENTER, line_spacing=1.1, font_name=FONT_NAME)
                     
        _add_rect_styled(slide, l=Inches(5.66), t=Inches(4.2), w=Inches(2.0), h=Inches(0.04),
                         fill_color=palette["primary"], line_color=palette["primary"])
                         
        if subtitle_text:
            _add_textbox(slide, subtitle_text, l=Inches(2.5), t=Inches(4.6), w=Inches(8.33), h=Inches(1.5),
                         font_size=FONT_SIZES['hero_body'], color=palette["text_muted"], align=PP_ALIGN.CENTER, line_spacing=1.25, font_name=FONT_NAME)
                     
    _set_speaker_notes(slide, data)



def _add_dashboard_slide(prs, data, palette):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, palette["bg"])
    
    title_text = data.get("title") or ""
    title_color = RGBColor(255, 255, 255) if (palette.get("grad_start") is not None) else palette["text"]
    _add_textbox(slide, title_text, l=Inches(0.8), t=Inches(0.6), w=Inches(11.73), h=Inches(1.0),
                 font_size=FONT_SIZES['slide_title'], bold=FONT_WEIGHTS['slide_title'], color=title_color)
                 
    bullets = data.get("bullet_points", [])
    
    _add_rect_styled(slide, l=Inches(0.8), t=Inches(1.8), w=Inches(6.8), h=Inches(4.9),
                     fill_color=palette["surface"], line_color=palette["surface"], rounded=True)
                     
    dashboard_metrics = data.get("dashboard_metrics")
    stats = []
    if dashboard_metrics:
        for item in dashboard_metrics[:3]:
            stats.append({"label": item.get("label", ""), "value": item.get("value", "")})
    else:
        for idx, bullet in enumerate(bullets[:3]):
            stats.append(_parse_bullet_to_metric_pptx(bullet, idx))
    
    mini_w = Inches(1.9)
    mini_gap = Inches(0.2)
    # Determine stat card background and border colors
    is_ecommerce = palette.get("grad_start") is not None
    is_academic = palette.get("font_heading") == "Georgia" and palette["bg"] == RGBColor(250, 250, 249)
    is_light_bg = is_ecommerce or is_academic
    
    card_fill = palette["surface"] if is_light_bg else palette["bg"]
    card_line = palette["surface"] if is_light_bg else palette["primary"]

    for idx, stat in enumerate(stats):
        x = Inches(1.1) + idx * (mini_w + mini_gap)
        y = Inches(2.1)
        _add_rect_styled(slide, l=x, t=y, w=mini_w, h=Inches(1.2),
                         fill_color=card_fill, line_color=card_line, rounded=True)
        _add_textbox(slide, stat["label"].upper()[:18], l=x + Inches(0.1), t=y + Inches(0.1), w=mini_w - Inches(0.2), h=Inches(0.4),
                     font_size=11, color=palette["text_muted"], align=PP_ALIGN.CENTER, font_name=FONT_NAME)
        
        val_str = str(stat["value"])
        val_fs = 22 if len(val_str) > 10 else 32
        _add_textbox(slide, val_str, l=x + Inches(0.1), t=y + Inches(0.55), w=mini_w - Inches(0.2), h=Inches(0.6),
                     font_size=val_fs, bold=True, color=palette["text"], align=PP_ALIGN.CENTER, font_name=FONT_NAME)

    # ── Try matplotlib chart from visual_engine ──
    chart_embedded = False
    chart_data = data.get("chart_data")
    if _VISUAL_ENGINE_AVAILABLE and chart_data:
        try:
            chart_path = render_chart(chart_data)
            if chart_path and os.path.exists(chart_path):
                slide.shapes.add_picture(
                    chart_path,
                    Inches(1.1), Inches(3.5), Inches(6.2), Inches(3.0)
                )
                chart_embedded = True
                print(f"[PPTX] Dashboard chart embedded: {chart_path}")
        except Exception as e:
            print(f"[PPTX] Dashboard chart embed failed: {e}")

    if not chart_embedded:
        chart_y = Inches(3.6)
        chart_title = data.get("chart_title") or (f"{data.get('title')} Analysis" if data.get("title") else "Trends Analysis")
        _add_textbox(slide, chart_title, l=Inches(1.1), t=chart_y, w=Inches(6.2), h=Inches(0.4),
                     font_size=FONT_SIZES['card_body'], bold=True, color=palette["text_muted"], font_name=FONT_NAME)
        chart_items = [s["label"] for s in stats] if stats else [b.split(":")[0].strip() for b in bullets if ":" in b]
        if not chart_items:
            chart_items = ["Trend 1", "Trend 2", "Trend 3"]
        _add_chart_visual_styled(slide, chart_items[:3], Inches(1.1), chart_y + Inches(0.4), Inches(6.2), Inches(2.3), palette)
    
    _add_rect_styled(slide, l=Inches(8.0), t=Inches(1.8), w=Inches(4.53), h=Inches(4.9),
                     fill_color=palette["surface"], line_color=palette["surface"], rounded=True)
                     
    insight_text = data.get("dashboard_insight") or (bullets[1] if len(bullets) > 1 else (bullets[0] if bullets else ""))
    
    sparkle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.3), Inches(2.1), Inches(0.65), Inches(0.65))
    _style_shape(sparkle, palette["primary"], palette["primary"])
    _set_center_text(sparkle, "✦", 20, RGBColor(255, 255, 255), font_name=FONT_NAME)
    
    _add_textbox(slide, "Executive Insights", l=Inches(9.1), t=Inches(2.1), w=Inches(3.2), h=Inches(0.6),
                 font_size=FONT_SIZES['section_heading'], bold=FONT_WEIGHTS['section_heading'], color=palette["text"], font_name=FONT_NAME)
                 
    _add_textbox(slide, insight_text, l=Inches(8.3), t=Inches(2.9), w=Inches(3.9), h=Inches(2.6),
                 font_size=FONT_SIZES['card_body'], color=palette["text_muted"], font_name=FONT_NAME)
                 
    badge_rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(5.8), Inches(3.9), Inches(0.5))
    _style_shape(badge_rect, palette["bg"], palette["primary"])
    
    badge_text = data.get("dashboard_insight_summary") or (f"★ Primary: {stats[0]['label']} - {stats[0]['value']}" if stats else "Target Achieved")
    _set_center_text(badge_rect, badge_text, 11, palette["primary"], font_name=FONT_NAME)
    
    _set_speaker_notes(slide, data)


def _add_architecture_slide(prs, data, palette):
    print(f"[PPTX] _add_architecture_slide called")
    print(f"   drawio_xml present: {'drawio_xml' in data}")
    print(f"   diagram_png present: {'diagram_png' in data}")
    drawio_xml = data.get("drawio_xml")
    diagram_png = data.get("diagram_png")
    
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, palette["bg"])
    
    # Try to use pre-generated PNG first (from /generate-json)
    if diagram_png:
        try:
            # Convert URL to file path correctly using E: temp directory
            import tempfile
            filename = diagram_png.split("/")[-1]
            if "diagram_cache" in diagram_png:
                png_path = os.path.join(tempfile.gettempdir(), "diagram_cache", filename)
            else:
                png_path = os.path.join(tempfile.gettempdir(), filename)
                if not os.path.exists(png_path):
                    png_path = os.path.join(tempfile.gettempdir(), "diagram_cache", filename)
            
            print(f"[PPTX] Using pre-generated PNG: {png_path}")
            if os.path.exists(png_path):
                print(f"[PPTX] Adding PNG to slide with aspect ratio preservation: {png_path}")
                from PIL import Image
                with Image.open(png_path) as img:
                    img_w, img_h = img.size
                
                img_ratio = img_w / img_h
                max_w_in = 12.33
                max_h_in = 6.5
                left_margin_in = 0.5
                top_margin_in = 0.5
                
                if img_ratio > (max_w_in / max_h_in):
                    w_in = max_w_in
                    h_in = max_w_in / img_ratio
                    l_in = left_margin_in
                    t_in = top_margin_in + (max_h_in - h_in) / 2.0
                else:
                    h_in = max_h_in
                    w_in = max_h_in * img_ratio
                    t_in = top_margin_in
                    l_in = left_margin_in + (max_w_in - w_in) / 2.0
                
                print(f"[PPTX] Calculated dimensions: l={l_in:.2f}\", t={t_in:.2f}\", w={w_in:.2f}\", h={h_in:.2f}\"")
                slide.shapes.add_picture(
                    png_path,
                    Inches(l_in),
                    Inches(t_in),
                    Inches(w_in),
                    Inches(h_in)
                )
                _set_speaker_notes(slide, data)
                return
            else:
                print(f"[PPTX] Pre-generated PNG not found: {png_path}")
        except Exception as e:
            print(f"[PPTX] Failed to use pre-generated PNG: {e}")
    
    # Fallback: generate PNG from drawio_xml
    if drawio_xml:
        try:
            from services.visual_engine import generate_drawio_diagram
            print(f"[PPTX] Calling generate_drawio_diagram with drawio_xml...")
            svg_path, png_path = generate_drawio_diagram(
                drawio_xml,
                topic=data.get("title") or data.get("deck_title") or "",
                disable_fallback=False  # Enable matplotlib fallback
            )
            print(f"[PPTX] generate_drawio_diagram returned: svg={svg_path}, png={png_path}")
            if png_path and os.path.exists(png_path):
                print(f"[PPTX] Adding PNG to slide with aspect ratio preservation: {png_path}")
                from PIL import Image
                with Image.open(png_path) as img:
                    img_w, img_h = img.size
                
                img_ratio = img_w / img_h
                max_w_in = 12.33
                max_h_in = 6.5
                left_margin_in = 0.5
                top_margin_in = 0.5
                
                if img_ratio > (max_w_in / max_h_in):
                    w_in = max_w_in
                    h_in = max_w_in / img_ratio
                    l_in = left_margin_in
                    t_in = top_margin_in + (max_h_in - h_in) / 2.0
                else:
                    h_in = max_h_in
                    w_in = max_h_in * img_ratio
                    t_in = top_margin_in
                    l_in = left_margin_in + (max_w_in - w_in) / 2.0
                
                print(f"[PPTX] Calculated dimensions: l={l_in:.2f}\", t={t_in:.2f}\", w={w_in:.2f}\", h={h_in:.2f}\"")
                slide.shapes.add_picture(
                    png_path,
                    Inches(l_in),
                    Inches(t_in),
                    Inches(w_in),
                    Inches(h_in)
                )
                _set_speaker_notes(slide, data)
                return
            else:
                print(f"[PPTX] Draw.io diagram render failed (PNG not found): {png_path}")
                print(f"[PPTX] Using matplotlib fallback instead of card layout")
                # Try matplotlib fallback
                try:
                    from services.visual_engine import _render_arch_matplotlib
                    matplotlib_png = _render_arch_matplotlib(
                        drawio_xml,
                        topic=data.get("title") or data.get("deck_title") or "",
                        slide_title=data.get("title") or "",
                        topology=data.get("topology", "")
                    )
                    if matplotlib_png and os.path.exists(matplotlib_png):
                        print(f"[PPTX] Adding matplotlib PNG to slide: {matplotlib_png}")
                        from PIL import Image
                        with Image.open(matplotlib_png) as img:
                            img_w, img_h = img.size
                        
                        img_ratio = img_w / img_h
                        max_w_in = 12.33
                        max_h_in = 6.5
                        left_margin_in = 0.5
                        top_margin_in = 0.5
                        
                        if img_ratio > (max_w_in / max_h_in):
                            w_in = max_w_in
                            h_in = max_w_in / img_ratio
                            l_in = left_margin_in
                            t_in = top_margin_in + (max_h_in - h_in) / 2.0
                        else:
                            h_in = max_h_in
                            w_in = max_h_in * img_ratio
                            t_in = top_margin_in
                            l_in = left_margin_in + (max_w_in - w_in) / 2.0
                        
                        slide.shapes.add_picture(
                            matplotlib_png,
                            Inches(l_in),
                            Inches(t_in),
                            Inches(w_in),
                            Inches(h_in)
                        )
                        _set_speaker_notes(slide, data)
                        return
                except Exception as matplotlib_err:
                    print(f"[PPTX] Matplotlib fallback also failed: {matplotlib_err}")
        except Exception as e:
            print(f"[PPTX] Draw.io diagram inline render failed: {e}")
            import traceback
            traceback.print_exc()
            print(f"[PPTX] Trying matplotlib fallback instead of card layout")
            # Try matplotlib fallback
            try:
                from services.visual_engine import _render_arch_matplotlib
                matplotlib_png = _render_arch_matplotlib(
                    drawio_xml,
                    topic=data.get("title") or data.get("deck_title") or "",
                    slide_title=data.get("title") or "",
                    topology=data.get("topology", "")
                )
                if matplotlib_png and os.path.exists(matplotlib_png):
                    print(f"[PPTX] Adding matplotlib PNG to slide: {matplotlib_png}")
                    from PIL import Image
                    with Image.open(matplotlib_png) as img:
                        img_w, img_h = img.size
                    
                    img_ratio = img_w / img_h
                    max_w_in = 12.33
                    max_h_in = 6.5
                    left_margin_in = 0.5
                    top_margin_in = 0.5
                    
                    if img_ratio > (max_w_in / max_h_in):
                        w_in = max_w_in
                        h_in = max_w_in / img_ratio
                        l_in = left_margin_in
                        t_in = top_margin_in + (max_h_in - h_in) / 2.0
                    else:
                        h_in = max_h_in
                        w_in = max_h_in * img_ratio
                        t_in = top_margin_in
                        l_in = left_margin_in + (max_w_in - w_in) / 2.0
                    
                    slide.shapes.add_picture(
                        matplotlib_png,
                        Inches(l_in),
                        Inches(t_in),
                        Inches(w_in),
                        Inches(h_in)
                    )
                    _set_speaker_notes(slide, data)
                    return
            except Exception as matplotlib_err:
                print(f"[PPTX] Matplotlib fallback also failed: {matplotlib_err}")
                # Only fall back to cards if both DrawIO and matplotlib fail
                print(f"[PPTX] All diagram rendering failed, using card-based layout as last resort")
                 
    # Fallback to full-width horizontal architecture cards
    bullets = data.get("bullet_points", [])
    architecture_nodes = data.get("architecture_nodes")
    nodes = []
    if architecture_nodes:
        for idx, item in enumerate(architecture_nodes[:4]):
            nodes.append({
                "title": item.get("title") or item.get("label") or f"Component {idx + 1}",
                "description": item.get("description") or item.get("body") or ""
            })
    else:
        for idx, bullet in enumerate(bullets[:4]):
            title = f"Component {idx + 1}"
            description = bullet
            if ":" in bullet:
                parts = bullet.split(":", 1)
                title = parts[0].strip()
                description = parts[1].strip()
            nodes.append({"title": title, "description": description})
            
    if not nodes:
        nodes = [
            {"title": "API Gateway", "description": "Routes incoming web traffic"},
            {"title": "App Service", "description": "Processes business logic"},
            {"title": "Auth Server", "description": "Authenticates user requests"},
            {"title": "Database", "description": "Stores secure application state"}
        ]
        
    canvas_w = Inches(11.73)
    canvas_h = Inches(5.5)
    canvas_t = Inches(1.0)
    # Background full-width surface canvas centered vertically
    _add_rect_styled(slide, l=Inches(0.8), t=canvas_t, w=canvas_w, h=canvas_h,
                     fill_color=palette["surface"], line_color=palette["surface"], rounded=True)
                     
    # Delegate rendering of nodes horizontally inside the full-width canvas
    _add_architecture_visual_styled(slide, nodes, Inches(0.8), canvas_t, canvas_w, canvas_h, palette)
    
    _set_speaker_notes(slide, data)


def _add_timeline_slide(prs, data, palette):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, palette["bg"])
    
    title_text = data.get("title") or ""
    title_color = RGBColor(255, 255, 255) if (palette.get("grad_start") is not None) else palette["text"]
    _add_textbox(slide, title_text, l=Inches(0.8), t=Inches(0.6), w=Inches(11.73), h=Inches(1.0),
                 font_size=FONT_SIZES['slide_title'], bold=FONT_WEIGHTS['slide_title'], color=title_color)
                 
    bullets = data.get("bullet_points", [])
    timeline_steps = data.get("timeline_steps")
    events = []
    if timeline_steps:
        for idx, item in enumerate(timeline_steps[:4]):
            events.append({
                "date": item.get("title") or item.get("date") or f"Phase {idx + 1}",
                "title": item.get("title") if item.get("date") else "Event",
                "description": item.get("description") or item.get("body") or ""
            })
    else:
        for idx, bullet in enumerate(bullets[:4]):
            date = f"202{4 + idx} Q{(idx % 4) + 1}"
            title = "Event"
            description = bullet
            if ":" in bullet:
                parts = bullet.split(":", 1)
                date = parts[0].strip()
                rest = parts[1].strip()
                if " - " in rest:
                    t_parts = rest.split(" - ", 1)
                    title = t_parts[0].strip()
                    description = t_parts[1].strip()
                else:
                    title = rest
            events.append({"date": date, "title": title, "description": description})
        
    track_y = Inches(4.2)
    _add_rect_styled(slide, l=Inches(0.8), t=track_y, w=Inches(11.73), h=Inches(0.06),
                     fill_color=palette["primary"], line_color=palette["primary"])
                     
    col_w = Inches(2.6)
    gap = Inches(0.4)
    start_x = Inches(0.8)
    
    for idx, evt in enumerate(events):
        x = start_x + idx * (col_w + gap)
        
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + col_w / 2 - Inches(0.2), track_y - Inches(0.17), Inches(0.4), Inches(0.4))
        _style_shape(dot, palette["secondary"], palette["secondary"])
        
        if idx % 2 == 0:
            _add_textbox(slide, evt["date"], l=x, t=track_y - Inches(0.7), w=col_w, h=Inches(0.4),
                         font_size=FONT_SIZES['section_heading'], bold=True, color=palette["primary"], align=PP_ALIGN.CENTER, font_name=FONT_NAME)
            
            card_y = track_y + Inches(0.4)
            _add_rect_styled(slide, l=x, t=card_y, w=col_w, h=Inches(2.0),
                             fill_color=palette["surface"], line_color=palette["surface"], rounded=True)
            _add_textbox(slide, evt["title"], l=x + Inches(0.1), t=card_y + Inches(0.1), w=col_w - Inches(0.2), h=Inches(0.45),
                         font_size=FONT_SIZES['section_heading'], bold=True, color=palette["text"], align=PP_ALIGN.CENTER, font_name=FONT_NAME)
            _add_textbox(slide, evt["description"], l=x + Inches(0.15), t=card_y + Inches(0.6), w=col_w - Inches(0.3), h=Inches(1.3),
                         font_size=FONT_SIZES['card_body'], color=palette["text_muted"], align=PP_ALIGN.CENTER, font_name=FONT_NAME)
        else:
            card_y = track_y - Inches(2.4)
            _add_rect_styled(slide, l=x, t=card_y, w=col_w, h=Inches(2.0),
                             fill_color=palette["surface"], line_color=palette["surface"], rounded=True)
            _add_textbox(slide, evt["title"], l=x + Inches(0.1), t=card_y + Inches(0.1), w=col_w - Inches(0.2), h=Inches(0.45),
                         font_size=FONT_SIZES['section_heading'], bold=True, color=palette["text"], align=PP_ALIGN.CENTER, font_name=FONT_NAME)
            _add_textbox(slide, evt["description"], l=x + Inches(0.15), t=card_y + Inches(0.6), w=col_w - Inches(0.3), h=Inches(1.3),
                         font_size=FONT_SIZES['card_body'], color=palette["text_muted"], align=PP_ALIGN.CENTER, font_name=FONT_NAME)
                         
            _add_textbox(slide, evt["date"], l=x, t=track_y + Inches(0.4), w=col_w, h=Inches(0.4),
                         font_size=FONT_SIZES['section_heading'], bold=True, color=palette["primary"], align=PP_ALIGN.CENTER, font_name=FONT_NAME)
                         
    _set_speaker_notes(slide, data)


def _add_comparison_slide_layout(prs, data, palette):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, palette["bg"])
    
    title_text = data.get("title") or ""
    title_color = RGBColor(255, 255, 255) if (palette.get("grad_start") is not None) else palette["text"]
    _add_textbox(slide, title_text, l=Inches(0.8), t=Inches(0.6), w=Inches(11.73), h=Inches(1.0),
                 font_size=FONT_SIZES['slide_title'], bold=FONT_WEIGHTS['slide_title'], color=title_color)

    # ── Try structured table from visual_engine ──
    table_data = data.get("table_data")
    if _VISUAL_ENGINE_AVAILABLE and table_data and table_data.get("headers"):
        try:
            table_shape = build_table(
                slide, table_data,
                left=Inches(0.8), top=Inches(1.8),
                width=Inches(11.73), height=Inches(4.9)
            )
            if table_shape:
                print("[PPTX] Comparison table embedded via visual_engine")
                _set_speaker_notes(slide, data)
                return
        except Exception as e:
            print(f"[PPTX] Comparison table embed failed: {e}")

    # Fallback: original two-column comparison cards
    comparison_items = data.get("comparison_items")
    left_title = data.get("comparison_left_title") or ""
    right_title = data.get("comparison_right_title") or ""
    
    left_points = data.get("left_points", [])
    right_points = data.get("right_points", [])
    
    if comparison_items:
        left_points = [item.get("left", "") for item in comparison_items]
        right_points = [item.get("right", "") for item in comparison_items]
        
    if not left_title or not right_title:
        if data.get("visual_items") and len(data.get("visual_items")) >= 2:
            left_title = left_title or data.get("visual_items")[0]
            right_title = right_title or data.get("visual_items")[1]
        elif not left_points and data.get("bullet_points"):
            bullets = data.get("bullet_points")
            if bullets and ":" in bullets[0]:
                left_title = left_title or bullets[0].split(":", 1)[0].strip()
            # If not split yet, split bullets in half
            half = len(bullets) // 2 if len(bullets) > 1 else len(bullets)
            left_points = bullets[:half]
            right_points = bullets[half:]
            
    if not left_title:
        left_title = "Overview"
    if not right_title:
        right_title = "Details"
        
    col_w = Inches(5.6)
    y = Inches(1.8)
    h = Inches(4.9)
    
    _add_rect_styled(slide, l=Inches(0.8), t=y, w=col_w, h=h,
                     fill_color=palette["surface"], line_color=palette["surface"], rounded=True)
    _add_textbox(slide, left_title, l=Inches(1.0), t=y + Inches(0.3), w=col_w - Inches(0.4), h=Inches(0.6),
                 font_size=FONT_SIZES['section_heading'], bold=True, color=palette["primary"], font_name=FONT_NAME)
    left_body = "\n".join(f"•  {pt}" for pt in left_points)
    _add_textbox(slide, left_body, l=Inches(1.0), t=y + Inches(1.0), w=col_w - Inches(0.4), h=h - Inches(1.2),
                 font_size=FONT_SIZES['card_body'], color=palette["text_muted"], font_name=FONT_NAME)
                 
    _add_rect_styled(slide, l=Inches(6.8), t=y, w=col_w, h=h,
                     fill_color=palette["surface"], line_color=palette["surface"], rounded=True)
    _add_textbox(slide, right_title, l=Inches(7.0), t=y + Inches(0.3), w=col_w - Inches(0.4), h=Inches(0.6),
                 font_size=FONT_SIZES['section_heading'], bold=True, color=palette["secondary"], font_name=FONT_NAME)
    right_body = "\n".join(f"•  {pt}" for pt in right_points)
    _add_textbox(slide, right_body, l=Inches(7.0), t=y + Inches(1.0), w=col_w - Inches(0.4), h=h - Inches(1.2),
                 font_size=FONT_SIZES['card_body'], color=palette["text_muted"], font_name=FONT_NAME)
                 
    _set_speaker_notes(slide, data)


def _add_process_slide_layout(prs, data, palette):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, palette["bg"])
    
    title_text = data.get("title") or ""
    title_color = RGBColor(255, 255, 255) if (palette.get("grad_start") is not None) else palette["text"]
    _add_textbox(slide, title_text, l=Inches(0.8), t=Inches(0.6), w=Inches(11.73), h=Inches(1.0),
                 font_size=FONT_SIZES['slide_title'], bold=FONT_WEIGHTS['slide_title'], color=title_color)
                 
    bullets = data.get("bullet_points", [])
    icon_keywords = data.get("icon_keywords", [])
    process_steps = data.get("process_steps")
    steps = []
    if process_steps:
        for idx, item in enumerate(process_steps[:4]):
            steps.append({"title": item.get("title") or f"Phase {idx + 1}", "body": item.get("body") or item.get("description") or ""})
    else:
        for idx, bullet in enumerate(bullets[:4]):
            title = f"Phase {idx + 1}"
            description = bullet
            if ":" in bullet:
                parts = bullet.split(":", 1)
                title = parts[0].strip()
                description = parts[1].strip()
            elif " - " in bullet:
                parts = bullet.split(" - ", 1)
                title = parts[0].strip()
                description = parts[1].strip()
            steps.append({"title": title, "body": description})
        
    col_w = Inches(2.6)
    gap = Inches(0.4)
    start_x = Inches(0.8)
    y = Inches(2.2)
    h = Inches(4.2)
    
    for idx, step in enumerate(steps):
        x = start_x + idx * (col_w + gap)
        
        _add_rect_styled(slide, l=x, t=y, w=col_w, h=h,
                         fill_color=palette["surface"], line_color=palette["surface"], rounded=True)

        # ── Try SVG icon from visual_engine ──
        icon_embedded = False
        if _VISUAL_ENGINE_AVAILABLE and icon_keywords:
            try:
                kw = icon_keywords[idx] if idx < len(icon_keywords) else step["title"]
                svg_str = get_icon(kw)
                if svg_str:
                    svg_path = os.path.join(tempfile.gettempdir(), f"proc_icon_{idx}_{id(slide)}.svg")
                    with open(svg_path, "w", encoding="utf-8") as f:
                        f.write(svg_str)
                    try:
                        slide.shapes.add_picture(
                            svg_path,
                            x + col_w / 2 - Inches(0.35), y + Inches(0.3), Inches(0.7), Inches(0.7)
                        )
                        icon_embedded = True
                    except Exception:
                        pass  # SVG not supported as picture in python-pptx, fallback below
            except Exception as e:
                print(f"[PPTX] Process icon embed failed: {e}")

        if not icon_embedded:
            bubble = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + col_w / 2 - Inches(0.35), y + Inches(0.3), Inches(0.7), Inches(0.7))
            _style_shape(bubble, palette["primary"], palette["primary"])
            _set_center_text(bubble, f"0{idx + 1}", 11, RGBColor(255, 255, 255), font_name=FONT_NAME)
        
        _add_textbox(slide, step["title"], l=x + Inches(0.1), t=y + Inches(1.2), w=col_w - Inches(0.2), h=Inches(0.6),
                     font_size=FONT_SIZES['section_heading'], bold=True, color=palette["text"], align=PP_ALIGN.CENTER, font_name=FONT_NAME)
                     
        _add_textbox(slide, step["body"], l=x + Inches(0.15), t=y + Inches(1.9), w=col_w - Inches(0.3), h=h - Inches(2.1),
                     font_size=FONT_SIZES['card_body'], color=palette["text_muted"], align=PP_ALIGN.CENTER, font_name=FONT_NAME)
                     
        if idx < len(steps) - 1:
            cx = x + col_w
            cy = y + Inches(2.0)
            _add_rect_styled(slide, l=cx, t=cy, w=gap, h=Inches(0.04),
                             fill_color=palette["primary"], line_color=palette["primary"])
                             
    _set_speaker_notes(slide, data)


def _add_fallback_slide_from_scratch(prs, slide_data, stype, palette):
    if stype == "title":
        _add_title_slide(prs, slide_data, palette)
    elif stype == "conclusion":
        _add_conclusion_slide(prs, slide_data, palette)
    else:
        _add_content_slide(prs, slide_data, palette)


def _populate_template_slide(slide, data, palette):
    title_shape, body_shape, subtitle_shape = _select_template_slots(slide)
    if not body_shape and subtitle_shape:
        body_shape = subtitle_shape

    keep_shapes = [title_shape, body_shape]
    _clear_template_text(slide, keep_shapes)
    _remove_large_content_visuals(slide, keep_shapes)

    title_text = data.get("title") or data.get("deck_title") or ""
    if not title_shape:
        title_shape = _add_fallback_text(slide, title_text, palette, is_title=True)
    else:
        _set_shape_text(title_shape, title_text)

    bullets = data.get("bullet_points", [])
    visual_type = data.get("visual_type") or _visual_for_slide_type(data.get("slide_type"))
    visual_items = data.get("visual_items") or bullets

    if not body_shape:
        body_shape = _add_fallback_text(slide, "", palette, is_title=False)

    if visual_type and visual_type != "none" and data.get("slide_type") != "title":
        _fit_body_to_left(body_shape)
        _set_bullets(body_shape, bullets[:4], font_name=palette.get("font_body", "Segoe UI"))
        _add_visual_styled(slide, visual_type, visual_items, Inches(7.0), Inches(1.55), Inches(5.6), Inches(4.9), palette, data)
    else:
        _set_bullets(body_shape, bullets, font_name=palette.get("font_body", "Segoe UI"))


def _select_template_slots(slide):
    title_shape = None
    body_shape = None
    subtitle_shape = None

    for shape in slide.placeholders:
        ph = shape.placeholder_format
        if ph.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.VERTICAL_TITLE):
            title_shape = shape
        elif ph.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.VERTICAL_BODY):
            body_shape = shape
        elif ph.type == PP_PLACEHOLDER.SUBTITLE:
            subtitle_shape = shape

    if not title_shape:
        title_shape = _find_best_text_shape(slide, prefer_title=True)
    if not body_shape:
        body_shape = _find_best_text_shape(slide, prefer_title=False, exclude=[title_shape])

    return title_shape, body_shape, subtitle_shape


def _clear_template_text(slide, keep_shapes):
    keep_ids = {id(shape) for shape in keep_shapes if shape is not None}
    for shape in slide.shapes:
        if shape.has_text_frame and id(shape) not in keep_ids:
            shape.text_frame.clear()


def _remove_large_content_visuals(slide, keep_shapes):
    keep_ids = {id(shape) for shape in keep_shapes if shape is not None}
    slide_area = int(Inches(13.33)) * int(Inches(7.5))
    for shape in list(slide.shapes):
        if id(shape) in keep_ids:
            continue
        area = _shape_area(shape)
        if getattr(shape, "shape_type", None) == 13 and area > slide_area * 0.08:
            _delete_shape(shape)


def _delete_shape(shape):
    element = shape._element
    element.getparent().remove(element)


def _fit_body_to_left(shape):
    if not shape:
        return
    shape.left = Inches(0.7)
    shape.top = max(shape.top, Inches(1.45))
    shape.width = Inches(5.7)
    shape.height = min(shape.height, Inches(4.9))


def _visual_for_slide_type(slide_type):
    return {
        "abstract": "icons",
        "problem": "icons",
        "solution": "process",
        "tech_stack": "icons",
        "architecture": "architecture",
        "process": "flowchart",
        "timeline": "process",
        "comparison": "comparison",
        "data": "chart",
        "conclusion": "icons",
    }.get(slide_type, "icons")


def _add_visual_styled(slide, visual_type, items, l, t, w, h, palette, data):
    labels = _visual_labels(items, data)
    if visual_type in ("architecture", "flowchart"):
        _add_architecture_visual_styled(slide, labels, l, t, w, h, palette)
    elif visual_type in ("process", "mindmap"):
        _add_process_visual_styled(slide, labels, l, t, w, h, palette)
    elif visual_type in ("chart", "infographic"):
        _add_chart_visual_styled(slide, labels, l, t, w, h, palette)
    elif visual_type == "table":
        _add_table_visual_styled(slide, labels, l, t, w, h, palette)
    elif visual_type == "comparison":
        _add_comparison_visual_styled(slide, labels, l, t, w, h, palette)
    elif visual_type == "image":
        image_embedded = False
        if _VISUAL_ENGINE_AVAILABLE:
            try:
                img_source = (items[0] if items else None) or data.get("image_keyword") or data.get("title") or "presentation"
                img_path = _resolve_image_to_local_path(img_source)
                if img_path and os.path.exists(img_path):
                    slide.shapes.add_picture(img_path, l, t, w, h)
                    image_embedded = True
            except Exception as e:
                print(f"[PPTX] TwoColumn image embed failed: {e}")
        if not image_embedded:
            _add_rect_styled(slide, l=l, t=t, w=w, h=h, fill_color=palette["surface"], line_color=palette["surface"], rounded=True)
            img_text = items[0] if items else "Featured Image Area"
            _add_textbox(slide, img_text, l=l, t=t+h/2-Inches(0.3), w=w, h=Inches(0.6), font_size=20, color=palette["text_muted"], align=PP_ALIGN.CENTER)
    else:
        _add_icon_visual_styled(slide, labels, l, t, w, h, palette)


def _add_visual_in_rect(slide, visual_type, items, l, t, w, h, palette, data):
    # Padding within the card container
    pl = l + Inches(0.2)
    pt = t + Inches(0.2)
    pw = w - Inches(0.4)
    ph = h - Inches(0.4)
    _add_visual_styled(slide, visual_type, items, pl, pt, pw, ph, palette, data)


def _visual_labels(items, data):
    labels = [str(item).strip() for item in (items or []) if str(item).strip()]
    if not labels:
        labels = [str(item).strip() for item in data.get("bullet_points", []) if str(item).strip()]
    return labels[:6] or [data.get("title", "Key point")]


def _add_icon_visual_styled(slide, labels, l, t, w, h, palette):
    cols = 2
    cell_w = w / cols
    cell_h = h / 3
    for idx, label in enumerate(labels[:6]):
        col = idx % cols
        row = idx // cols
        x = l + cell_w * col + Inches(0.15)
        y = t + cell_h * row + Inches(0.15)
        icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.65), Inches(0.65))
        _style_shape(icon, palette["primary"], palette["primary"])
        _set_center_text(icon, "✓", 16, RGBColor(255, 255, 255))
        
        _add_textbox(slide, label, x + Inches(0.8), y, cell_w - Inches(0.95), Inches(0.65), font_size=18, color=palette["text"], line_spacing=1.5)


def _add_process_visual_styled(slide, labels, l, t, w, h, palette):
    count = min(len(labels), 4)
    if count == 0:
        return
    gap = w / count
    y = t + h / 2
    for idx, label in enumerate(labels[:count]):
        x = l + gap * idx + Inches(0.15)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y - Inches(0.5), Inches(1.0), Inches(1.0))
        _style_shape(box, palette["primary"], palette["primary"])
        _set_center_text(box, str(idx + 1), 18, RGBColor(255, 255, 255))
        _add_textbox(slide, label, x - Inches(0.2), y + Inches(0.6), Inches(1.4), Inches(0.9), font_size=16, color=palette["text_muted"], align=PP_ALIGN.CENTER, line_spacing=1.4)
        if idx < count - 1:
            cx = x + Inches(1.05)
            cw = gap - Inches(1.2)
            _add_rect_styled(slide, l=cx, t=y, w=cw, h=Inches(0.04), fill_color=palette["primary"], line_color=palette["primary"])


def _add_architecture_visual_styled(slide, nodes, l, t, w, h, palette):
    # Normalize nodes to list of dicts
    normalized_nodes = []
    for idx, item in enumerate(nodes):
        if isinstance(item, dict):
            normalized_nodes.append({
                "title": item.get("title") or f"Component {idx + 1}",
                "description": item.get("description") or ""
            })
        else:
            # It's a string (label)
            normalized_nodes.append({
                "title": f"Component {idx + 1}",
                "description": str(item)
            })
            
    while len(normalized_nodes) < 4:
        normalized_nodes.append({
            "title": f"Component {len(normalized_nodes) + 1}",
            "description": ""
        })
        
    nodes_count = min(len(normalized_nodes), 4)
    
    card_w = Inches(2.6)
    card_h = Inches(3.2)
    
    # Calculate spacing centered horizontally in width w
    if nodes_count > 1:
        gap = (w - (card_w * nodes_count)) / (nodes_count - 1)
        if gap > Inches(1.5):
            gap = Inches(1.5)
        total_w = card_w * nodes_count + gap * (nodes_count - 1)
        start_x = l + (w - total_w) / 2
    else:
        start_x = l + (w - card_w) / 2
        gap = Inches(0)
        
    y_center = t + (h - card_h) / 2
    
    # Draw horizontal connector lines first
    line_color = palette.get("secondary", palette["primary"])
    for i in range(nodes_count - 1):
        cx = start_x + card_w + i * (card_w + gap)
        cy = y_center + card_h / 2
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy - Inches(0.02), gap, Inches(0.04))
        _style_shape(line, line_color, line_color)
        
    # Draw cards on top
    colors = [palette["primary"], palette.get("secondary", palette["primary"]), palette.get("text_muted", palette["primary"]), palette["primary"]]
    for i, node in enumerate(normalized_nodes[:4]):
        x = start_x + i * (card_w + gap)
        card = _add_rect_styled(slide, l=x, t=y_center, w=card_w, h=card_h,
                                 fill_color=palette["surface"], line_color=colors[i % len(colors)], rounded=True)
        card.line.width = Pt(1.5)
        
        # Title
        _add_textbox(slide, node["title"], l=x + Inches(0.15), t=y_center + Inches(0.3), w=card_w - Inches(0.3), h=Inches(0.6),
                     font_size=15, bold=True, color=palette["text"], align=PP_ALIGN.CENTER, font_name=FONT_NAME)
                     
        # Description
        _add_textbox(slide, node["description"], l=x + Inches(0.15), t=y_center + Inches(1.0), w=card_w - Inches(0.3), h=card_h - Inches(1.2),
                     font_size=11, color=palette["text_muted"], align=PP_ALIGN.CENTER, font_name=FONT_NAME)


def _add_chart_visual_styled(slide, labels, l, t, w, h, palette):
    max_bar_h = h - Inches(1.4)
    bar_w = w / max(len(labels), 1) - Inches(0.25)
    for idx, label in enumerate(labels[:4]):
        value = (idx + 2) / (len(labels[:4]) + 2)
        bar_h = int(max_bar_h * value)
        x = l + (w / max(len(labels[:4]), 1)) * idx + Inches(0.15)
        y = t + max_bar_h - bar_h + Inches(0.25)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, bar_w, bar_h)
        _style_shape(bar, palette["primary"], palette["primary"])
        _add_textbox(slide, label, x, t + max_bar_h + Inches(0.35), bar_w, Inches(0.7), font_size=16, color=palette["text_muted"], align=PP_ALIGN.CENTER, line_spacing=1.4)


def _add_table_visual_styled(slide, labels, l, t, w, h, palette):
    active_palette = getattr(thread_local, 'palette', DEFAULT_PALETTE)
    f_heading = active_palette.get("font_heading", "Segoe UI")
    f_body = active_palette.get("font_body", "Segoe UI")
    rows = min(len(labels), 4) + 1
    table_shape = slide.shapes.add_table(rows, 2, l + Inches(0.25), t + Inches(0.25), w - Inches(0.5), h - Inches(0.5))
    table = table_shape.table
    
    # Configure column widths
    table.columns[0].width = Inches(1.1)
    table.columns[1].width = w - Inches(1.6)
    
    # Style Header
    cell_lh = table.cell(0, 0)
    cell_lh.text = "Item"
    cell_lh.fill.solid()
    cell_lh.fill.fore_color.rgb = palette["primary"]
    run_lh = cell_lh.text_frame.paragraphs[0].runs[0]
    run_lh.font.name = f_heading
    run_lh.font.size = Pt(20)
    run_lh.font.color.rgb = RGBColor(255, 255, 255)
    
    cell_rh = table.cell(0, 1)
    cell_rh.text = "Details"
    cell_rh.fill.solid()
    cell_rh.fill.fore_color.rgb = palette["primary"]
    run_rh = cell_rh.text_frame.paragraphs[0].runs[0]
    run_rh.font.name = f_heading
    run_rh.font.size = Pt(20)
    run_rh.font.color.rgb = RGBColor(255, 255, 255)
    
    for idx, label in enumerate(labels[:rows - 1], start=1):
        cell_idx = table.cell(idx, 0)
        cell_idx.text = f"0{idx}"
        cell_idx.fill.solid()
        cell_idx.fill.fore_color.rgb = palette["surface"]
        p_idx = cell_idx.text_frame.paragraphs[0]
        if p_idx.runs:
            p_idx.runs[0].font.name = f_body
            p_idx.runs[0].font.size = Pt(18)
            p_idx.runs[0].font.color.rgb = palette["primary"]
            
        cell_lbl = table.cell(idx, 1)
        cell_lbl.text = label
        cell_lbl.fill.solid()
        cell_lbl.fill.fore_color.rgb = palette["surface"]
        p_lbl = cell_lbl.text_frame.paragraphs[0]
        if p_lbl.runs:
            p_lbl.runs[0].font.name = f_body
            p_lbl.runs[0].font.size = Pt(18)
            p_lbl.runs[0].font.color.rgb = palette["text"]


def _add_comparison_visual_styled(slide, labels, l, t, w, h, palette):
    left = labels[0] if labels else "Option A"
    right = labels[1] if len(labels) > 1 else "Option B"
    cw = w / 2 - Inches(0.25)
    for idx, label in enumerate((left, right)):
        x = l + (w / 2) * idx + Inches(0.15)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, t, cw, h)
        _style_shape(box, palette["surface"], palette["secondary"])
        _set_center_text(box, label, 20, palette["text"])


def _style_shape(shape, fill_color, line_color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color


def _set_center_text(shape, text, font_size, color, font_name=None):
    active_palette = getattr(thread_local, 'palette', DEFAULT_PALETTE)
    if font_name is None:
        font_name = active_palette.get("font_body", "Segoe UI")
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color


def _connect_centers(slide, start_shape, end_shape, color):
    x1 = start_shape.left + start_shape.width / 2
    y1 = start_shape.top + start_shape.height / 2
    x2 = end_shape.left + end_shape.width / 2
    y2 = end_shape.top + end_shape.height / 2
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)


def _set_speaker_notes(slide, data):
    if data.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = data["speaker_notes"]


def _delete_slide(prs, index):
    r_id = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(r_id)
    del prs.slides._sldIdLst[index]


def _find_best_text_shape(slide, prefer_title=False, exclude=None):
    exclude_ids = {id(shape) for shape in (exclude or set()) if shape is not None}
    shapes = [
        shape for shape in slide.shapes
        if id(shape) not in exclude_ids and shape.has_text_frame and not shape.is_placeholder
    ]
    if not shapes:
        return None

    if prefer_title:
        shapes.sort(key=lambda shape: (shape.top, -_shape_area(shape)))
    else:
        shapes.sort(key=lambda shape: (-_shape_area(shape), shape.top))
    return shapes[0]


def _shape_area(shape):
    return int(shape.width) * int(shape.height)


def _set_shape_text(shape, text):
    tf = shape.text_frame
    tf.clear()
    tf.paragraphs[0].text = text


def _set_bullets(shape, bullets, font_name=None):
    active_palette = getattr(thread_local, 'palette', DEFAULT_PALETTE)
    if font_name is None:
        font_name = FONT_NAME
    tf = shape.text_frame
    tf.clear()
    
    # 1. Join all bullet points with newlines to detect language and check fitting
    full_text = "\n".join(str(b) for b in bullets)
    is_multi, lang, char_factor = detect_multilingual_mode(full_text)
    
    # Enable word wrap
    tf.word_wrap = True
    
    # Usable font size estimation
    current_font_size_val = 14.0  # default bullet point size
    try:
        current_font_size_val = float(FONT_SIZES['bullet_point'].pt)
    except Exception:
        pass
        
    if is_multi:
        fitted_val = fit_multilingual_text_to_shape(full_text, shape.width, shape.height, current_font_size_val, lang)
        bullet_font_size = Pt(fitted_val)
        
        # Override line spacing for Hindi, Tamil, Telugu
        if lang in ("hindi", "tamil", "telugu"):
            line_spacing = 1.3
        elif lang in ("arabic", "cjk"):
            line_spacing = 1.2
        else:
            line_spacing = 1.1
    else:
        bullet_font_size = FONT_SIZES['bullet_point']
        line_spacing = 1.25
        
    for idx, point in enumerate(bullets):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        
        # Determine indentation level
        level = 0
        cleaned = point
        if isinstance(point, str):
            stripped = point.lstrip()
            indent = len(point) - len(stripped)
            if indent >= 4:
                level = 2
            elif indent >= 2:
                level = 1
            cleaned = stripped
            
        para.level = level
        para.space_after = Pt(8)
        para.line_spacing = line_spacing
        
        # Set RTL properties if Arabic
        if is_multi and lang == "arabic":
            try:
                pPr = para._p.get_or_add_pPr()
                pPr.set('rtl', '1')
                para.alignment = PP_ALIGN.RIGHT
            except Exception as e:
                print(f"[RTL] Warning: failed to set rtl on bullet paragraph: {e}")
                
        # Clean prefix
        for prefix in ["•", "-", "*", "✓"]:
            if cleaned.startswith(prefix):
                cleaned = cleaned[len(prefix):].strip()
                
        if level == 0:
            if lang == "arabic":
                run_text = para.add_run()
                run_text.text = cleaned
                run_text.font.name = font_name
                run_text.font.size = bullet_font_size
                run_text.font.bold = FONT_WEIGHTS['bullet_point']
                run_text.font.color.rgb = active_palette["text_muted"]
            else:
                run_check = para.add_run()
                run_check.text = "✓  "
                run_check.font.name = font_name
                run_check.font.size = bullet_font_size
                run_check.font.bold = True
                run_check.font.color.rgb = RGBColor(34, 197, 94)
                
                run_text = para.add_run()
                run_text.text = cleaned
                run_text.font.name = font_name
                run_text.font.size = bullet_font_size
                run_text.font.bold = FONT_WEIGHTS['bullet_point']
                run_text.font.color.rgb = active_palette["text_muted"]
        else:
            # Sub-bullets
            run_text = para.add_run()
            run_text.text = cleaned
            run_text.font.name = font_name
            run_text.font.size = max(Pt(9), bullet_font_size - Pt(2 * level))
            run_text.font.bold = False
            run_text.font.color.rgb = active_palette["text_muted"]


def _add_fallback_text(slide, text, palette, is_title=False):
    # Search if there is already a textbox we can write to
    for shape in slide.shapes:
        if shape.has_text_frame and not shape.is_placeholder:
            tf = shape.text_frame
            if not tf.text.strip():
                tf.text = text
                return shape
    
    # Otherwise add a new textbox - Gamma-style typography
    if is_title:
        return _add_textbox(slide, text, Inches(1.2), Inches(1.2), Inches(10.93), Inches(1.2), font_size=FONT_SIZES['slide_title'], bold=FONT_WEIGHTS['slide_title'], color=palette["text"], line_spacing=1.2, space_after=10, font_name=FONT_NAME)
    else:
        return _add_textbox(slide, text, Inches(1.2), Inches(2.6), Inches(10.93), Inches(4.5), font_size=FONT_SIZES['card_body'], color=palette["text_muted"], line_spacing=1.25, space_after=8, font_name=FONT_NAME)


def _add_title_slide(prs, data, palette):
    layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, palette["bg"])
    
    title_text = data.get("deck_title") or data.get("title", "Presentation")
    title_color = RGBColor(255, 255, 255) if (palette.get("grad_start") is not None) else palette["text"]
    subtitle_color = RGBColor(240, 240, 240) if (palette.get("grad_start") is not None) else palette["text_muted"]
    _add_textbox(slide, title_text,
                 l=Inches(1.2), t=Inches(2.6), w=Inches(10.93), h=Inches(1.8),
                 font_size=_get_hero_title_font_size(title_text), bold=True, color=title_color, align=PP_ALIGN.CENTER, line_spacing=1.1, space_after=10, font_name=FONT_NAME)
                 
    if data.get("bullet_points"):
        subtitle = data["bullet_points"][0]
        _add_textbox(slide, subtitle,
                     l=Inches(2.5), t=Inches(4.7), w=Inches(8.33), h=Inches(1.0),
                     font_size=FONT_SIZES['hero_body'], color=subtitle_color, align=PP_ALIGN.CENTER, line_spacing=1.25, font_name=FONT_NAME)


def _add_content_slide(prs, data, palette):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, palette["bg"])
    
    # Title - Gamma-style large title
    title_text = data.get("title", "")
    title_color = RGBColor(255, 255, 255) if (palette.get("grad_start") is not None) else palette["text"]
    _add_textbox(slide, title_text,
                 l=Inches(0.8), t=Inches(0.6), w=Inches(11.73), h=Inches(1.0),
                 font_size=FONT_SIZES['slide_title'], bold=FONT_WEIGHTS['slide_title'], color=title_color, line_spacing=1.2, space_after=10, font_name=FONT_NAME)
                 
    has_visual = data.get("visual_type") and data.get("visual_type") != "none"
    body_width = Inches(6.0) if has_visual else Inches(11.73)

    # Bullet points - limited to 5 with safer spacing
    bullets = data.get("bullet_points", [])[:5]
    if palette.get("grad_start") is not None:
        _add_rect_styled(slide, l=Inches(0.8), t=Inches(1.9), w=body_width, h=Inches(4.9),
                         fill_color=palette["surface"], line_color=palette["surface"], rounded=True)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), body_width, Inches(4.9))
    frame = txBox.text_frame
    frame.word_wrap = True
    frame.clear()
    
    font_body = FONT_NAME
    for i, point in enumerate(bullets):
        para = frame.add_paragraph() if i > 0 else frame.paragraphs[0]
        para.space_after = Pt(8)
        para.line_spacing = 1.25
        
        # Clean prefix
        cleaned = point.strip()
        for prefix in ["•", "-", "*", "✓"]:
            if cleaned.startswith(prefix):
                cleaned = cleaned[len(prefix):].strip()
                
        # Check if challenge
        is_challenge = any(cleaned.lower().startswith(pfx) for pfx in ["challenge:", "key challenge:", "problem:"])
        
        # Run 1: Green checkmark ✓
        run_check = para.add_run()
        run_check.text = "✓  "
        run_check.font.name = font_body
        run_check.font.size = FONT_SIZES['bullet_point']
        run_check.font.bold = True
        run_check.font.color.rgb = RGBColor(34, 197, 94) # Green check
        
        # Run 2: Bullet text
        if is_challenge:
            colon_idx = cleaned.find(":")
            label = cleaned[:colon_idx+1] if colon_idx != -1 else "Challenge:"
            desc = cleaned[colon_idx+1:].strip() if colon_idx != -1 else cleaned
            
            run_label = para.add_run()
            run_label.text = label + " "
            run_label.font.name = font_body
            run_label.font.size = FONT_SIZES['bullet_point']
            run_label.font.bold = True
            run_label.font.color.rgb = RGBColor(16, 185, 129) # Standout green for Challenge label
            
            run_text = para.add_run()
            run_text.text = desc
            run_text.font.name = font_body
            run_text.font.size = FONT_SIZES['bullet_point']
            run_text.font.bold = False
            run_text.font.color.rgb = palette["text"]
        else:
            run_text = para.add_run()
            run_text.text = cleaned
            run_text.font.name = font_body
            run_text.font.size = FONT_SIZES['bullet_point']
            run_text.font.bold = FONT_WEIGHTS['bullet_point']
            run_text.font.color.rgb = palette["text_muted"]

    if has_visual:
        # Draw visual card container - rounded corners
        _add_rect_styled(slide, l=Inches(7.1), t=Inches(2.1), w=Inches(5.43), h=Inches(4.8), fill_color=palette["surface"], line_color=palette["surface"], rounded=True)
        _add_visual_in_rect(slide, data.get("visual_type"), data.get("visual_items") or bullets, Inches(7.1), Inches(2.1), Inches(5.43), Inches(4.8), palette, data)
        
    # Speaker notes
    if data.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = data["speaker_notes"]


def _add_conclusion_slide(prs, data, palette):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, palette["bg"])
    
    # Title - Gamma-style large title
    title_text = data.get("title", "")
    title_color = RGBColor(255, 255, 255) if (palette.get("grad_start") is not None) else palette["text"]
    _add_textbox(slide, title_text,
                 l=Inches(1.5), t=Inches(2.0), w=Inches(10.33), h=Inches(1.2),
                 font_size=FONT_SIZES['slide_title'], bold=FONT_WEIGHTS['slide_title'], color=title_color, align=PP_ALIGN.CENTER, line_spacing=1.2, space_after=10, font_name=FONT_NAME)
                 
    bullets = data.get("bullet_points", [])[:5]  # Limit to 5 bullets
    if palette.get("grad_start") is not None:
        _add_rect_styled(slide, l=Inches(2.5), t=Inches(3.5), w=Inches(8.33), h=Inches(3.2),
                         fill_color=palette["surface"], line_color=palette["surface"], rounded=True)
    txBox = slide.shapes.add_textbox(Inches(2.5), Inches(3.5), Inches(8.33), Inches(3.2))
    frame = txBox.text_frame
    frame.word_wrap = True
    frame.clear()
    
    font_body = FONT_NAME
    for idx, b in enumerate(bullets):
        para = frame.add_paragraph() if idx > 0 else frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.space_after = Pt(8)
        para.line_spacing = 1.25
        
        # Clean prefix
        cleaned = b.strip()
        for prefix in ["•", "-", "*", "✓"]:
            if cleaned.startswith(prefix):
                cleaned = cleaned[len(prefix):].strip()
                
        # Run 1: Green checkmark ✓
        run_check = para.add_run()
        run_check.text = "✓  "
        run_check.font.name = font_body
        run_check.font.size = FONT_SIZES['bullet_point']
        run_check.font.bold = True
        run_check.font.color.rgb = RGBColor(34, 197, 94) # Green check
        
        # Run 2: Bullet text
        run_text = para.add_run()
        run_text.text = cleaned
        run_text.font.name = font_body
        run_text.font.size = FONT_SIZES['bullet_point']
        run_text.font.bold = FONT_WEIGHTS['bullet_point']
        run_text.font.color.rgb = palette["text_muted"]


def _set_shape_picture_fill(shape, image_path):
    if not image_path or not os.path.exists(image_path):
        return False
    try:
        from pptx.oxml.xmlchemy import OxmlElement
        from pptx.oxml.ns import qn
        part = shape.part
        image_part, rId = part.get_or_add_image_part(image_path)
        spPr = shape._element.spPr
        for fill_tag in ['a:solidFill', 'a:gradFill', 'a:pattFill', 'a:noFill']:
            fill_el = spPr.find(qn(fill_tag))
            if fill_el is not None:
                spPr.remove(fill_el)
        blipFill = OxmlElement('a:blipFill')
        blip = OxmlElement('a:blip')
        blip.set(qn('r:embed'), rId)
        stretch = OxmlElement('a:stretch')
        fillRect = OxmlElement('a:fillRect')
        stretch.append(fillRect)
        blipFill.append(blip)
        blipFill.append(stretch)
        spPr.append(blipFill)
        if shape.line:
            shape.line.fill.background()
        return True
    except Exception as e:
        print(f"[PPTX] Failed to apply shape XML picture fill: {e}")
        return False


# ── helpers ──────────────────────────────────────────────────
def _fill_bg(slide, color, palette=None):
    """Fill slide background. Uses gradient if palette has grad_start/grad_mid/grad_end."""
    bg = slide.background
    fill = bg.fill
    if palette and palette.get("grad_start"):
        # E-commerce gradient background
        from pptx.oxml.ns import qn
        fill.gradient()
        fill.gradient_angle = 135  # Diagonal gradient
        # Clear existing stops and add custom ones using oxml element wrapper
        gs_lst = fill._fill._gradFill.get_or_add_gsLst()
        for child in list(gs_lst):
            gs_lst.remove(child)
        # Add 3 gradient stops
        for pos, stop_key in [(0, "grad_start"), (50000, "grad_mid"), (100000, "grad_end")]:
            gs = gs_lst.makeelement(qn('a:gs'), {'pos': str(pos)})
            # RGB Color value in hex format
            hex_color = "{:02X}{:02X}{:02X}".format(palette[stop_key][0], palette[stop_key][1], palette[stop_key][2])
            srgb = gs.makeelement(qn('a:srgbClr'), {'val': hex_color})
            gs.append(srgb)
            gs_lst.append(gs)
    else:
        fill.solid()
        fill.fore_color.rgb = color


def _add_ecommerce_decorative_shapes(slide, palette):
    """Add scattered decorative geometric shapes for e-commerce tone.
    Creates triangles, plus signs, squares, and small circles
    in semi-transparent style around the slide periphery."""
    from pptx.oxml.ns import qn
    
    decorations = [
        # (shape_type, left, top, width, height, rotation, color, opacity)
        # Small triangles
        (MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(0.4), Inches(0.5), Inches(0.5), Inches(0.5), 15, palette["primary"], 0.25),
        (MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(12.0), Inches(5.5), Inches(0.6), Inches(0.6), -20, palette["secondary"], 0.20),
        (MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.5), Inches(1.2), Inches(0.35), Inches(0.35), 45, palette.get("accent", palette["primary"]), 0.30),
        # Plus signs (cross shapes)
        (MSO_SHAPE.CROSS, Inches(1.0), Inches(5.8), Inches(0.45), Inches(0.45), 0, palette["secondary"], 0.25),
        (MSO_SHAPE.CROSS, Inches(12.4), Inches(3.0), Inches(0.35), Inches(0.35), 15, palette["primary"], 0.20),
        # Small squares
        (MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(3.2), Inches(0.3), Inches(0.3), 30, palette.get("accent", palette["primary"]), 0.22),
        (MSO_SHAPE.RECTANGLE, Inches(12.6), Inches(0.8), Inches(0.25), Inches(0.25), 45, palette["primary"], 0.18),
        # Circles (dot clusters)
        (MSO_SHAPE.OVAL, Inches(0.6), Inches(1.8), Inches(0.15), Inches(0.15), 0, palette["secondary"], 0.30),
        (MSO_SHAPE.OVAL, Inches(0.9), Inches(1.95), Inches(0.12), Inches(0.12), 0, palette["secondary"], 0.22),
        (MSO_SHAPE.OVAL, Inches(0.7), Inches(2.15), Inches(0.10), Inches(0.10), 0, palette["secondary"], 0.18),
        (MSO_SHAPE.OVAL, Inches(12.1), Inches(6.2), Inches(0.15), Inches(0.15), 0, palette["primary"], 0.28),
        (MSO_SHAPE.OVAL, Inches(12.35), Inches(6.35), Inches(0.10), Inches(0.10), 0, palette["primary"], 0.20),
    ]
    
    for shape_type, left, top, width, height, rotation, color, opacity in decorations:
        try:
            shape = slide.shapes.add_shape(shape_type, left, top, width, height)
            shape.rotation = rotation
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            # Set opacity via XML
            solid_fill = shape.fill._fill._solidFill
            if solid_fill is not None:
                srgb = solid_fill.find(qn('a:srgbClr'))
                if srgb is not None:
                    alpha = srgb.makeelement(qn('a:alpha'), {'val': str(int(opacity * 100000))})
                    srgb.append(alpha)
            shape.line.fill.background()  # No border
        except Exception as e:
            print(f"[PPTX] Decorative shape failed: {e}")




def _add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def _add_rect_styled(slide, l, t, w, h, fill_color, line_color, rounded=False):
    if rounded:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    else:
        shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    return shape


def _add_textbox(slide, text, l, t, w, h,
                 font_size=15, bold=False,
                 color=RGBColor(0,0,0), align=PP_ALIGN.LEFT, font_name=None,
                 line_spacing=1.2, space_after=6, fit_text=False):
    active_palette = getattr(thread_local, 'palette', DEFAULT_PALETTE)
    
    # 1. Detect language and enable multilingual mode if appropriate
    is_multi, lang, char_factor = detect_multilingual_mode(str(text))
    
    # 2. Handle font size conversion
    try:
        from pptx.util import Length
        if isinstance(font_size, Length):
            font_size_val = float(font_size.pt)
        else:
            font_size_val = float(font_size)
    except Exception:
        font_size_val = float(font_size) if isinstance(font_size, (int, float)) else 14.0
        
    # 3. Fit font size and set line spacing if in multilingual mode
    if is_multi:
        fitted_val = fit_multilingual_text_to_shape(str(text), w, h, font_size_val, lang)
        sz = Pt(fitted_val)
        
        # Increase line spacing for Indic scripts
        if lang in ("hindi", "tamil", "telugu"):
            line_spacing = 1.3
        elif lang in ("arabic", "cjk"):
            line_spacing = 1.2
    else:
        sz = Pt(font_size_val)

    if font_name is None:
        has_emoji = any(ord(char) > 0xffff for char in str(text)) if text else False
        if has_emoji:
            font_name = 'Segoe UI Emoji'
        else:
            font_name = FONT_NAME if FONT_NAME else (active_palette.get("font_heading", "Segoe UI") if bold else active_palette.get("font_body", "Segoe UI"))
        
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True # Word wrap always enabled
    
    from pptx.enum.text import MSO_AUTO_SIZE
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
        
    lines = text.split("\n") if isinstance(text, str) else [str(text)]
    
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        
        # Handle paragraph alignment for Arabic/RTL
        p_align = align
        if is_multi and lang == "arabic":
            if align == PP_ALIGN.LEFT or align is None:
                p_align = PP_ALIGN.RIGHT
        p.alignment = p_align
        
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        
        if is_multi and lang == "arabic":
            try:
                pPr = p._p.get_or_add_pPr()
                pPr.set('rtl', '1')
            except Exception as e:
                print(f"[RTL] Warning: failed to set rtl on paragraph: {e}")
        
        # Configure formatting on run level
        if p.runs:
            for run in p.runs:
                run.font.name = font_name
                run.font.size = sz
                run.font.bold = bold
                run.font.color.rgb = color
        else:
            run = p.add_run()
            run.text = line
            run.font.name = font_name
            run.font.size = sz
            run.font.bold = bold
            run.font.color.rgb = color
            
    return txBox


def _populate_bullets_in_text_frame(tf, bullets, font_name, palette, is_first=False):
    active_palette = palette or getattr(thread_local, 'palette', DEFAULT_PALETTE)
    
    # 1. Join all bullet points with newlines to detect language and check fitting
    full_text = "\n".join(str(b) for b in bullets)
    is_multi, lang, char_factor = detect_multilingual_mode(full_text)
    
    # Enable word wrap
    tf.word_wrap = True
    
    # Usable font size estimation
    current_font_size_val = 14.0  # default bullet point size
    try:
        current_font_size_val = float(FONT_SIZES['bullet_point'].pt)
    except Exception:
        pass
        
    if is_multi:
        # We need height/width of shape to fit text, try to extract them
        try:
            shape = tf._parent
            w, h = shape.width, shape.height
        except Exception:
            w, h = Inches(5.6), Inches(4.9)  # Standard column size fallback
            
        fitted_val = fit_multilingual_text_to_shape(full_text, w, h, current_font_size_val, lang)
        bullet_font_size = Pt(fitted_val)
        
        # Override line spacing for Hindi, Tamil, Telugu
        if lang in ("hindi", "tamil", "telugu"):
            line_spacing = 1.3
        elif lang in ("arabic", "cjk"):
            line_spacing = 1.2
        else:
            line_spacing = 1.1
    else:
        bullet_font_size = FONT_SIZES['bullet_point']
        line_spacing = 1.25
        
    for idx, b in enumerate(bullets):
        p = tf.paragraphs[0] if (idx == 0 and is_first) else tf.add_paragraph()
        p.space_after = Pt(8)
        p.line_spacing = line_spacing
        
        # Determine indentation level
        level = 0
        cleaned = b
        if isinstance(b, str):
            stripped = b.lstrip()
            indent = len(b) - len(stripped)
            if indent >= 4:
                level = 2
            elif indent >= 2:
                level = 1
            cleaned = stripped
            
        p.level = level
        
        # Set RTL properties if Arabic
        if is_multi and lang == "arabic":
            try:
                pPr = p._p.get_or_add_pPr()
                pPr.set('rtl', '1')
                p.alignment = PP_ALIGN.RIGHT
            except Exception as e:
                print(f"[RTL] Warning: failed to set rtl on bullet paragraph: {e}")
                
        # Clean prefix
        for prefix in ["•", "-", "*", "✓"]:
            if cleaned.startswith(prefix):
                cleaned = cleaned[len(prefix):].strip()
                
        # Check if challenge
        is_challenge = any(cleaned.lower().startswith(pfx) for pfx in ["challenge:", "key challenge:", "problem:"])
        
        if level == 0:
            if lang == "arabic":
                run_text = p.add_run()
                run_text.text = cleaned
                run_text.font.name = font_name
                run_text.font.size = bullet_font_size
                run_text.font.bold = FONT_WEIGHTS['bullet_point']
                run_text.font.color.rgb = active_palette["text_muted"]
            else:
                run_check = p.add_run()
                run_check.text = "✓  "
                run_check.font.name = font_name
                run_check.font.size = bullet_font_size
                run_check.font.bold = True
                run_check.font.color.rgb = RGBColor(34, 197, 94) # Green check
                
                if is_challenge:
                    colon_idx = cleaned.find(":")
                    label = cleaned[:colon_idx+1] if colon_idx != -1 else "Challenge:"
                    desc = cleaned[colon_idx+1:].strip() if colon_idx != -1 else cleaned
                    
                    run_label = p.add_run()
                    run_label.text = label + " "
                    run_label.font.name = font_name
                    run_label.font.size = bullet_font_size
                    run_label.font.bold = True
                    run_label.font.color.rgb = RGBColor(16, 185, 129)
                    
                    run_text = p.add_run()
                    run_text.text = desc
                    run_text.font.name = font_name
                    run_text.font.size = bullet_font_size
                    run_text.font.bold = False
                    run_text.font.color.rgb = active_palette["text"]
                else:
                    run_text = p.add_run()
                    run_text.text = cleaned
                    run_text.font.name = font_name
                    run_text.font.size = bullet_font_size
                    run_text.font.bold = FONT_WEIGHTS['bullet_point']
                    run_text.font.color.rgb = active_palette["text_muted"]
        else:
            # Sub-bullets: standard bullets using level indent, no manual checkmark prefix
            run_text = p.add_run()
            run_text.text = cleaned
            run_text.font.name = font_name
            run_text.font.size = max(Pt(9), bullet_font_size - Pt(2 * level))
            run_text.font.bold = False
            run_text.font.color.rgb = active_palette["text_muted"]


def detect_multilingual_mode(text: str):
    if not text:
        return False, "en", 0.4
    
    has_devanagari = False
    has_tamil = False
    has_telugu = False
    has_arabic = False
    has_cjk = False
    
    for char in text:
        o = ord(char)
        # Devanagari: 2304 to 2431
        if 2304 <= o <= 2431:
            has_devanagari = True
        # Tamil: 2944 to 3071
        elif 2944 <= o <= 3071:
            has_tamil = True
        # Telugu: 3072 to 3199
        elif 3072 <= o <= 3199:
            has_telugu = True
        # Arabic: 1536-1791, 1872-1919, 2208-2303
        elif (1536 <= o <= 1791) or (1872 <= o <= 1919) or (2208 <= o <= 2303):
            has_arabic = True
        # CJK Unified Ideographs, Hiragana, Katakana, Hangul, CJK Symbols:
        elif (19968 <= o <= 40959) or (12352 <= o <= 12543) or (44032 <= o <= 55215) or (12288 <= o <= 12351):
            has_cjk = True
            
    is_multi = has_devanagari or has_tamil or has_telugu or has_arabic or has_cjk
    
    if has_cjk:
        return is_multi, "cjk", 1.0
    elif has_devanagari or has_tamil or has_telugu:
        lang = "hindi" if has_devanagari else ("tamil" if has_tamil else "telugu")
        return is_multi, lang, 0.6
    elif has_arabic:
        return is_multi, "arabic", 0.5
    else:
        return is_multi, "other", 0.4


def fit_multilingual_text_to_shape(text: str, width, height, current_font_size: float, language: str) -> float:
    try:
        width_in = float(width.inches)
    except Exception:
        width_in = float(width) / 72.0 if isinstance(width, (int, float)) else 3.0
        
    try:
        height_in = float(height.inches)
    except Exception:
        height_in = float(height) / 72.0 if isinstance(height, (int, float)) else 2.0

    usable_width = max(0.5, width_in - 0.3)
    usable_height = max(0.5, height_in - 0.24)

    if language == "cjk":
        char_factor = 1.0
        line_spacing = 1.2
    elif language in ("hindi", "tamil", "telugu"):
        char_factor = 0.6
        line_spacing = 1.3
    elif language == "arabic":
        char_factor = 0.5
        line_spacing = 1.2
    else:
        char_factor = 0.4
        line_spacing = 1.1

    import math
    font_size = float(current_font_size)
    space_after = 6.0  # Default space after paragraph in Pt

    while font_size >= 9.0:
        char_width_in = (font_size / 72.0) * char_factor
        chars_per_line = usable_width / char_width_in
        
        paragraphs = text.split("\n") if isinstance(text, str) else [str(text)]
        total_lines = 0
        for para in paragraphs:
            para_len = len(para)
            if para_len == 0:
                total_lines += 1
            else:
                para_usable_width = max(0.5, usable_width - 0.25)
                total_lines += int(math.ceil(para_len / (para_usable_width / char_width_in)))
        
        line_height_in = line_spacing * (font_size / 72.0)
        space_after_in = (space_after / 72.0)
        
        total_text_height = (total_lines * line_height_in) + (len(paragraphs) - 1) * space_after_in
        
        if total_text_height <= usable_height:
            break
            
        font_size -= 1.0
        
    return max(9.0, font_size)


def fix_presentation_complex_scripts(prs):
    """Post-processor that ensures all runs containing complex script text (non-ASCII characters)
    have the complex script typeface (<a:cs>) properly set in their OOXML properties.
    This also handles RTL layout rules and auto-fitting text to prevent overflow.
    """
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
    import math

    def _get_all_shapes_recursive(shapes_collection):
        result = []
        for shape in shapes_collection:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    result.extend(_get_all_shapes_recursive(shape.shapes))
                else:
                    result.append(shape)
            except Exception:
                result.append(shape)
        return result

    def _process_text_frame(tf, width, height):
        full_text = ""
        paragraphs_text = []
        for para in tf.paragraphs:
            para_text = "".join(run.text for run in para.runs if run.text)
            paragraphs_text.append(para_text)
        full_text = "\n".join(paragraphs_text)
        
        if not full_text.strip():
            return
            
        is_multi, lang, char_factor = detect_multilingual_mode(full_text)
        if not is_multi:
            return
            
        if lang in ("hindi", "tamil", "telugu"):
            line_spacing = 1.3
        elif lang in ("arabic", "cjk"):
            line_spacing = 1.2
        else:
            line_spacing = 1.1

        tf.word_wrap = True
        
        # Determine current font size
        current_font_size = 14.0
        first_run_size = None
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.size:
                    first_run_size = float(run.font.size.pt)
                    break
            if first_run_size is not None:
                break
        if first_run_size is not None:
            current_font_size = first_run_size
            
        # Fit font size to shape
        fitted_font_size = fit_multilingual_text_to_shape(
            full_text, width, height, current_font_size, lang
        )
        
        scale_factor = fitted_font_size / current_font_size if current_font_size > 0 else 1.0
        
        # Truncate text if it still overflows at 9 pt to prevent vertical/ugly wrapping
        try:
            width_in = float(width.inches)
        except Exception:
            width_in = 3.0
        try:
            height_in = float(height.inches)
        except Exception:
            height_in = 2.0
            
        usable_width = max(0.5, width_in - 0.3)
        usable_height = max(0.5, height_in - 0.24)
        
        char_width_in = (fitted_font_size / 72.0) * char_factor
        line_height_in = line_spacing * (fitted_font_size / 72.0)
        space_after_in = (6.0 / 72.0)
        
        # Calculate total lines
        total_lines = 0
        for para in tf.paragraphs:
            para_text = "".join(run.text for run in para.runs if run.text)
            para_len = len(para_text)
            if para_len == 0:
                total_lines += 1
            else:
                para_usable_width = max(0.5, usable_width - 0.25)
                total_lines += int(math.ceil(para_len / (para_usable_width / char_width_in)))
                
        total_text_height = (total_lines * line_height_in) + (len(tf.paragraphs) - 1) * space_after_in
        
        if total_text_height > usable_height and fitted_font_size <= 9.5:
            # Truncation logic to prevent vertical text or severe clipping
            max_lines = max(1, int(usable_height // line_height_in))
            lines_written = 0
            paragraphs_to_keep = []
            
            for para in tf.paragraphs:
                para_text = "".join(run.text for run in para.runs if run.text)
                para_len = len(para_text)
                if para_len == 0:
                    para_lines = 1
                else:
                    para_usable_width = max(0.5, usable_width - 0.25)
                    para_lines = int(math.ceil(para_len / (para_usable_width / char_width_in)))
                    
                if lines_written + para_lines <= max_lines:
                    paragraphs_to_keep.append(para)
                    lines_written += para_lines
                else:
                    remaining_lines = max_lines - lines_written
                    if remaining_lines > 0:
                        chars_per_line = max(1, int(para_usable_width // char_width_in))
                        char_budget = remaining_lines * chars_per_line - 3
                        if char_budget > 5:
                            truncated_text = para_text[:char_budget] + "..."
                            if para.runs:
                                para.runs[0].text = truncated_text
                                for r in para.runs[1:]:
                                    r.text = ""
                            else:
                                r = para.add_run()
                                r.text = truncated_text
                            paragraphs_to_keep.append(para)
                        lines_written += remaining_lines
                    break
                    
            # Clear text in discarded paragraphs
            all_paras = list(tf.paragraphs)
            for p_to_del in all_paras[len(paragraphs_to_keep):]:
                for r in p_to_del.runs:
                    r.text = ""
                p_to_del.text = ""

        # Format remaining paragraphs
        for para in tf.paragraphs:
            para.line_spacing = line_spacing
            
            # Apply RTL if Arabic
            if lang == "arabic":
                try:
                    pPr = para._p.get_or_add_pPr()
                    pPr.set('rtl', '1')
                    if para.alignment is None or para.alignment == PP_ALIGN.LEFT:
                        para.alignment = PP_ALIGN.RIGHT
                except Exception as e:
                    print(f"[COMPLEX_SCRIPT] Warning: failed to apply RTL: {e}")
                    
            for run in para.runs:
                if run.text:
                    if run.font.size:
                        orig_pt = float(run.font.size.pt)
                        run.font.size = Pt(max(9.0, round(orig_pt * scale_factor)))
                    else:
                        run.font.size = Pt(fitted_font_size)
                        
                    # Apply complex script tag
                    try:
                        rPr = run.font._rPr
                        if rPr is not None:
                            font_name = run.font.name
                            if not font_name:
                                font_name = "Calibri"
                            
                            latin = rPr.find(qn('a:latin'))
                            if latin is None:
                                latin = OxmlElement('a:latin')
                                rPr.append(latin)
                            if not latin.get('typeface'):
                                latin.set('typeface', font_name)
                                
                            cs = rPr.find(qn('a:cs'))
                            if cs is None:
                                cs = OxmlElement('a:cs')
                                rPr.append(cs)
                            cs.set('typeface', font_name)
                            
                            ea = rPr.find(qn('a:ea'))
                            if ea is None:
                                ea = OxmlElement('a:ea')
                                rPr.append(ea)
                            if not ea.get('typeface'):
                                ea.set('typeface', font_name)
                    except Exception as e:
                        print(f"[COMPLEX_SCRIPT] Warning: failed to apply cs tag: {e}")
                        
        # Finally enable PowerPoint's native auto-sizing as secondary safeguard
        try:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception as e:
            print(f"[COMPLEX_SCRIPT] Warning: failed to apply auto_size: {e}")

    for slide in prs.slides:
        all_shapes = _get_all_shapes_recursive(slide.shapes)
        for shape in all_shapes:
            try:
                w, h = shape.width, shape.height
            except Exception:
                w, h = Inches(3), Inches(2)
                
            if shape.has_text_frame:
                _process_text_frame(shape.text_frame, w, h)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        try:
                            cell_w = cell.width
                            cell_h = row.height
                        except Exception:
                            cell_w, cell_h = Inches(2), Inches(1)
                        _process_text_frame(cell.text_frame, cell_w, cell_h)
