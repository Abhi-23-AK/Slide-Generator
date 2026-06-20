import os
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Design System Tokens
PRIMARY = RGBColor(255, 214, 0)       # #FFD600 (Yellow accent blocks)
BACKGROUND = RGBColor(255, 255, 255) # #FFFFFF (White background)
CARD_BG = RGBColor(245, 245, 245)    # #F5F5F5 (Light gray cards)
TITLE_COLOR = RGBColor(20, 20, 20)   # #141414 (Black typography)
BODY_COLOR = RGBColor(100, 100, 100) # #646464 (Small muted body text)

def _get_prof_hero_title_font_size(title_text: str) -> int:
    t_len = len(str(title_text))
    if t_len > 50:
        return 28
    elif t_len > 35:
        return 34
    elif t_len > 20:
        return 40
    else:
        return 48

def _parse_bullet_to_metric(bullet, idx):
    import re
    bullet = bullet.strip()
    
    # Clean prefixes
    cleaned = bullet
    for prefix in ["•", "-", "*", "✓"]:
        if cleaned.startswith(prefix):
            cleaned = cleaned[len(prefix):].strip()
            
    # Try split by colon or dash
    parts = None
    if ":" in cleaned:
        parts = cleaned.split(":", 1)
    elif " - " in cleaned:
        parts = cleaned.split(" - ", 1)
        
    if parts:
        left, right = parts[0].strip(), parts[1].strip()
        # A right part is a short metric if it has <= 3 words or is <= 15 chars
        is_short_metric = len(right.split()) <= 3 or len(right) <= 15
        if is_short_metric:
            return {"label": left, "value": right}
        else:
            # Right part is a long sentence description. We treat the whole bullet as the label.
            # Try to extract a metric from the text, otherwise use fallback.
            match = re.search(r'(\+?-?\d+\.?\d*%\s*)|(\$\d+,?\d*\s*)|(\b\d+\.?\d*[KMB]?\b\s*)', cleaned)
            if match:
                val = match.group(0).strip()
            else:
                default_values = ["84.5%", "$12,400", "98.2%", "+12.4%"]
                val = default_values[idx % len(default_values)]
            return {"label": cleaned, "value": val}
            
    # Try to extract a percentage, currency, or count from a flat bullet
    match = re.search(r'(\+?-?\d+\.?\d*%\s*)|(\$\d+,?\d*\s*)|(\b\d+\.?\d*[KMB]?\b\s*)', cleaned)
    if match:
        val = match.group(0).strip()
        label = cleaned.replace(val, "").strip()
        label = re.sub(r'\s+', ' ', label).strip().strip(',').strip('.').strip(':').strip(' - ')
        if not label:
            label = f"Metric {idx + 1}"
        return {"label": label, "value": val}
    else:
        default_values = ["84.5%", "$12,400", "98.2%", "+12.4%"]
        val = default_values[idx % len(default_values)]
        return {"label": cleaned, "value": val}

def _add_prof_hero_slide(prs, data, palette):
    from services.pptx_builder import (
        _add_textbox,
        _add_rect_styled,
        _set_speaker_notes,
        _resolve_image_to_local_path,
        _set_shape_picture_fill,
        _fill_bg,
    )

    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, BACKGROUND)

    # Left yellow vertical rectangle (25% width)
    # 13.33 * 0.25 = 3.33 inches
    col_w = Inches(3.33)
    _add_rect_styled(
        slide,
        l=Inches(0),
        t=Inches(0),
        w=col_w,
        h=Inches(7.5),
        fill_color=PRIMARY,
        line_color=PRIMARY,
    )

    # Image inside lower portion of yellow column
    img_source = (
        (data.get("visual_items")[0] if data.get("visual_items") else None)
        or data.get("image_keyword")
        or data.get("title")
        or "presentation"
    )
    img_path = _resolve_image_to_local_path(img_source)

    img_w = Inches(2.6)
    img_h = Inches(2.6)
    img_x = Inches(0.365)  # (3.33 - 2.6)/2
    img_y = Inches(4.3)

    img_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, img_x, img_y, img_w, img_h)
    img_card.fill.solid()
    img_card.fill.fore_color.rgb = CARD_BG
    img_card.line.color.rgb = RGBColor(220, 220, 220)
    img_card.line.width = Pt(1)

    if img_path and os.path.exists(img_path):
        _set_shape_picture_fill(img_card, img_path)

    # Title, subtitle, minimal body text on right
    title_text = data.get("title") or data.get("deck_title") or "Untitled Section"
    subtitle_text = ""
    bullets = data.get("bullet_points", [])

    if bullets:
        subtitle_text = bullets[0]
        body_bullets = bullets[1:3]
    else:
        body_bullets = []

    # Large bold title on right
    title_fs = _get_prof_hero_title_font_size(title_text)
    _add_textbox(
        slide,
        title_text,
        l=Inches(4.2),
        t=Inches(1.5),
        w=Inches(8.3),
        h=Inches(1.8),
        font_size=title_fs,
        bold=True,
        color=TITLE_COLOR,
        align=PP_ALIGN.LEFT,
        line_spacing=1.1,
    )

    # Subtitle below title
    if subtitle_text:
        _add_textbox(
            slide,
            subtitle_text,
            l=Inches(4.2),
            t=Inches(3.3),
            w=Inches(8.3),
            h=Inches(1.0),
            font_size=20,
            bold=True,
            color=BODY_COLOR,
            align=PP_ALIGN.LEFT,
            line_spacing=1.2,
        )

    # Minimal body text below subtitle
    if body_bullets:
        body_text = "\n".join(body_bullets)
        _add_textbox(
            slide,
            body_text,
            l=Inches(4.2),
            t=Inches(4.5),
            w=Inches(8.3),
            h=Inches(2.0),
            font_size=13,
            bold=False,
            color=BODY_COLOR,
            align=PP_ALIGN.LEFT,
            line_spacing=1.3,
            space_after=6,
        )

    _set_speaker_notes(slide, data)

def _add_prof_1_column_slide(prs, data, palette):
    from services.pptx_builder import (
        _add_textbox,
        _add_rect_styled,
        _set_speaker_notes,
        _resolve_image_to_local_path,
        _set_shape_picture_fill,
        _fill_bg,
    )

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, BACKGROUND)

    # Title
    title_text = data.get("title") or ""
    _add_textbox(
        slide,
        title_text,
        l=Inches(0.8),
        t=Inches(0.6),
        w=Inches(11.73),
        h=Inches(1.0),
        font_size=32,
        bold=True,
        color=TITLE_COLOR,
    )

    bullets = data.get("bullet_points", [])
    headline = data.get("headline") or ""
    if headline == title_text:
        headline = ""

    # Check for image
    img_source = (
        (data.get("visual_items")[0] if data.get("visual_items") else None)
        or data.get("image_keyword")
        or data.get("title")
        or ""
    )
    img_path = _resolve_image_to_local_path(img_source) if img_source else None
    has_image = img_path and os.path.exists(img_path)

    body_width = Inches(6.0) if has_image else Inches(11.73)

    # Left column: text/bullets
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), body_width, Inches(4.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)

    is_first = True
    if headline:
        p = tf.paragraphs[0]
        p.text = headline
        p.space_after = Pt(14)
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = "Calibri"
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = TITLE_COLOR
        is_first = False

    # Limit bullet points to 4 for minimal/non-crowded feel
    for b in bullets[:4]:
        p = tf.paragraphs[0] if is_first else tf.add_paragraph()
        p.space_after = Pt(10)
        p.line_spacing = 1.3

        cleaned = b.strip()
        for prefix in ["•", "-", "*", "✓"]:
            if cleaned.startswith(prefix):
                cleaned = cleaned[len(prefix):].strip()

        # Run 1: Yellow checkmark ✓
        run_check = p.add_run()
        run_check.text = "✓  "
        run_check.font.name = "Calibri"
        run_check.font.size = Pt(14)
        run_check.font.bold = True
        run_check.font.color.rgb = PRIMARY

        # Run 2: Text
        run_text = p.add_run()
        run_text.text = cleaned
        run_text.font.name = "Calibri"
        run_text.font.size = Pt(14)
        run_text.font.bold = False
        run_text.font.color.rgb = BODY_COLOR

        is_first = False

    # Right column: Image
    if has_image:
        # Gray background card
        _add_rect_styled(
            slide,
            l=Inches(7.3),
            t=Inches(1.8),
            w=Inches(5.2),
            h=Inches(4.9),
            fill_color=CARD_BG,
            line_color=CARD_BG,
            rounded=True,
        )

        # Image framed inside card
        img_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.45), Inches(1.95), Inches(4.9), Inches(4.6)
        )
        img_card.fill.solid()
        img_card.fill.fore_color.rgb = CARD_BG
        # Thin border around image
        img_card.line.color.rgb = RGBColor(220, 220, 220)
        img_card.line.width = Pt(1)

        _set_shape_picture_fill(img_card, img_path)

    _set_speaker_notes(slide, data)

def _add_prof_2_column_slide(prs, data, palette):
    from services.pptx_builder import (
        _add_textbox,
        _add_rect_styled,
        _set_speaker_notes,
        _resolve_image_to_local_path,
        _set_shape_picture_fill,
        _fill_bg,
        _add_visual_in_rect,
    )

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, BACKGROUND)

    # Title
    title_text = data.get("title") or ""
    _add_textbox(
        slide,
        title_text,
        l=Inches(0.8),
        t=Inches(0.6),
        w=Inches(11.73),
        h=Inches(1.0),
        font_size=32,
        bold=True,
        color=TITLE_COLOR,
    )

    left_headline = data.get("left_headline") or "Key Points"
    left_text = data.get("left_text") or ""
    bullets = data.get("bullet_points", [])

    # Left: Yellow content card
    _add_rect_styled(
        slide,
        l=Inches(0.8),
        t=Inches(1.8),
        w=Inches(5.6),
        h=Inches(4.9),
        fill_color=PRIMARY,
        line_color=PRIMARY,
        rounded=True,
    )

    # Text inside Yellow card
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    is_first = True
    if left_headline:
        p = tf.paragraphs[0]
        p.text = left_headline
        p.space_after = Pt(12)
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = "Calibri"
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = TITLE_COLOR
        is_first = False

    if left_text:
        p = tf.paragraphs[0] if is_first else tf.add_paragraph()
        p.text = left_text
        p.space_after = Pt(10)
        p.line_spacing = 1.25
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = "Calibri"
        run.font.size = Pt(16)
        run.font.bold = False
        run.font.color.rgb = TITLE_COLOR
        is_first = False

    for b in bullets[:4]:
        p = tf.paragraphs[0] if is_first else tf.add_paragraph()
        p.space_after = Pt(8)
        p.line_spacing = 1.25

        cleaned = b.strip()
        for prefix in ["•", "-", "*", "✓"]:
            if cleaned.startswith(prefix):
                cleaned = cleaned[len(prefix):].strip()

        # Dark checkmark on yellow card
        run_check = p.add_run()
        run_check.text = "✓  "
        run_check.font.name = "Calibri"
        run_check.font.size = Pt(14)
        run_check.font.bold = True
        run_check.font.color.rgb = TITLE_COLOR

        run_text = p.add_run()
        run_text.text = cleaned
        run_text.font.name = "Calibri"
        run_text.font.size = Pt(14)
        run_text.font.bold = False
        run_text.font.color.rgb = TITLE_COLOR

        is_first = False

    # Right: Image card
    _add_rect_styled(
        slide,
        l=Inches(6.8),
        t=Inches(1.8),
        w=Inches(5.7),
        h=Inches(4.9),
        fill_color=CARD_BG,
        line_color=CARD_BG,
        rounded=True,
    )

    visual_type = data.get("visual_type", "none")
    visual_items = data.get("visual_items", [])
    img_source = visual_items[0] if visual_items else (data.get("image_keyword") or "")
    img_path = _resolve_image_to_local_path(img_source) if img_source else None

    if img_path and os.path.exists(img_path):
        img_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.95), Inches(1.95), Inches(5.4), Inches(4.6)
        )
        img_card.fill.solid()
        img_card.fill.fore_color.rgb = CARD_BG
        img_card.line.color.rgb = RGBColor(220, 220, 220)
        img_card.line.width = Pt(1)
        _set_shape_picture_fill(img_card, img_path)
    else:
        # Fallback using PPTX helper if visual_type is set
        _add_visual_in_rect(
            slide,
            visual_type,
            visual_items,
            Inches(6.95),
            Inches(1.95),
            Inches(5.4),
            Inches(4.6),
            palette,
            data,
        )

    _set_speaker_notes(slide, data)

def _add_prof_3_column_slide(prs, data, palette):
    from services.pptx_builder import (
        _add_textbox,
        _add_rect_styled,
        _style_shape,
        _set_center_text,
        _set_speaker_notes,
        _fill_bg,
    )

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, BACKGROUND)

    # Title
    title_text = data.get("title") or ""
    _add_textbox(
        slide,
        title_text,
        l=Inches(0.8),
        t=Inches(0.6),
        w=Inches(11.73),
        h=Inches(1.0),
        font_size=32,
        bold=True,
        color=TITLE_COLOR,
    )

    grid_items = data.get("grid_items")
    cols = []
    if grid_items:
        for item in grid_items[:3]:
            cols.append((item.get("title", ""), item.get("description", "")))
    else:
        bullets = data.get("bullet_points", [])
        for bullet in bullets[:3]:
            if ":" in bullet:
                h, t = bullet.split(":", 1)
                cols.append((h.strip(), t.strip()))
            elif " - " in bullet:
                h, t = bullet.split(" - ", 1)
                cols.append((h.strip(), t.strip()))
            else:
                cols.append((bullet, ""))

    col_w = Inches(3.7)
    gap = Inches(0.3)
    y = Inches(1.8)
    h = Inches(4.9)

    for idx, (heading, text) in enumerate(cols):
        x = Inches(0.8) + idx * (col_w + gap)

        # Card shape
        card = _add_rect_styled(
            slide, l=x, t=y, w=col_w, h=h, fill_color=CARD_BG, line_color=CARD_BG, rounded=True
        )
        card.line.color.rgb = RGBColor(220, 220, 220)
        card.line.width = Pt(1)

        # Modern icon circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.3), Inches(0.6), Inches(0.6))
        _style_shape(circle, PRIMARY, PRIMARY)
        _set_center_text(circle, f"0{idx + 1}", 12, TITLE_COLOR, font_name="Calibri")

        # Heading
        _add_textbox(
            slide,
            heading,
            l=x + Inches(0.3),
            t=y + Inches(1.1),
            w=col_w - Inches(0.6),
            h=Inches(0.5),
            font_size=18,
            bold=True,
            color=TITLE_COLOR,
        )

        # Description
        _add_textbox(
            slide,
            text,
            l=x + Inches(0.3),
            t=y + Inches(1.7),
            w=col_w - Inches(0.6),
            h=Inches(2.9),
            font_size=13,
            color=BODY_COLOR,
        )

    _set_speaker_notes(slide, data)

def _add_prof_4_grid_slide(prs, data, palette):
    from services.pptx_builder import (
        _add_textbox,
        _add_rect_styled,
        _style_shape,
        _set_center_text,
        _set_speaker_notes,
        _fill_bg,
    )

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, BACKGROUND)

    # Title - Centered section title
    title_text = data.get("title") or ""
    _add_textbox(
        slide,
        title_text,
        l=Inches(0.8),
        t=Inches(0.6),
        w=Inches(11.73),
        h=Inches(1.0),
        font_size=32,
        bold=True,
        color=TITLE_COLOR,
        align=PP_ALIGN.CENTER,
    )

    grid_items = data.get("grid_items")
    cells = []
    if grid_items:
        for item in grid_items[:4]:
            cells.append({"label": item.get("title", ""), "desc": item.get("description", "")})
    else:
        bullets = data.get("bullet_points", [])
        for idx, bullet in enumerate(bullets[:4]):
            label = f"Metric {idx + 1}"
            desc = bullet
            if ":" in bullet:
                parts = bullet.split(":", 1)
                label = parts[0].strip()
                desc = parts[1].strip()
            elif " - " in bullet:
                parts = bullet.split(" - ", 1)
                label = parts[0].strip()
                desc = parts[1].strip()
            cells.append({"label": label, "desc": desc})

    col_w = Inches(5.6)
    row_h = Inches(2.2)
    gap_x = Inches(0.5)
    gap_y = Inches(0.4)
    start_x = Inches(0.8)
    start_y = Inches(1.8)

    for idx, cell in enumerate(cells):
        col = idx % 2
        row = idx // 2
        x = start_x + col * (col_w + gap_x)
        y = start_y + row * (row_h + gap_y)

        # Card
        card = _add_rect_styled(
            slide, l=x, t=y, w=col_w, h=row_h, fill_color=CARD_BG, line_color=CARD_BG, rounded=True
        )
        card.line.color.rgb = RGBColor(220, 220, 220)
        card.line.width = Pt(1)

        # Circular icon on left
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.3), Inches(0.6), Inches(0.6))
        _style_shape(circle, PRIMARY, PRIMARY)
        _set_center_text(circle, f"0{idx + 1}", 12, TITLE_COLOR, font_name="Calibri")

        # Title on the right of the icon
        _add_textbox(
            slide,
            cell["label"],
            l=x + Inches(1.1),
            t=y + Inches(0.25),
            w=col_w - Inches(1.3),
            h=Inches(0.4),
            font_size=16,
            bold=True,
            color=TITLE_COLOR,
        )

        # Description
        _add_textbox(
            slide,
            cell["desc"],
            l=x + Inches(1.1),
            t=y + Inches(0.7),
            w=col_w - Inches(1.3),
            h=Inches(1.3),
            font_size=13,
            color=BODY_COLOR,
            fit_text=True,
        )

    _set_speaker_notes(slide, data)

def _add_prof_process_slide(prs, data, palette):
    from services.pptx_builder import (
        _add_textbox,
        _add_rect_styled,
        _style_shape,
        _set_center_text,
        _set_speaker_notes,
        _fill_bg,
    )

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, BACKGROUND)

    # Title
    title_text = data.get("title") or ""
    _add_textbox(
        slide,
        title_text,
        l=Inches(0.8),
        t=Inches(0.6),
        w=Inches(11.73),
        h=Inches(1.0),
        font_size=32,
        bold=True,
        color=TITLE_COLOR,
    )

    process_steps = data.get("process_steps")
    steps = []
    if process_steps:
        for item in process_steps[:4]:
            steps.append({"title": item.get("title", ""), "body": item.get("body", "")})
    else:
        bullets = data.get("bullet_points", [])
        for idx, bullet in enumerate(bullets[:4]):
            title = f"Phase {idx + 1}"
            description = bullet
            if ":" in bullet:
                parts = bullet.split(":", 1)
                title = parts[0].strip()
                description = parts[1].strip()
            elif " - " in bullet:
                parts = bullet.split(" - ", 1)
                title = parts[0].strip()
                description = parts[1].strip()
            steps.append({"title": title, "body": description})

    col_w = Inches(2.7)
    gap = Inches(0.3)
    start_x = Inches(0.8)
    y = Inches(2.0)
    h = Inches(4.5)

    for idx, step in enumerate(steps):
        x = start_x + idx * (col_w + gap)

        # Card
        card = _add_rect_styled(
            slide, l=x, t=y, w=col_w, h=h, fill_color=CARD_BG, line_color=CARD_BG, rounded=True
        )
        card.line.color.rgb = RGBColor(220, 220, 220)
        card.line.width = Pt(1)

        # Step number in yellow circle on top
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + (col_w - Inches(0.7)) / 2, y + Inches(0.4), Inches(0.7), Inches(0.7))
        _style_shape(circle, PRIMARY, PRIMARY)
        _set_center_text(circle, f"0{idx + 1}", 12, TITLE_COLOR, font_name="Calibri")

        # Title
        _add_textbox(
            slide,
            step["title"],
            l=x + Inches(0.1),
            t=y + Inches(1.3),
            w=col_w - Inches(0.2),
            h=Inches(0.5),
            font_size=16,
            bold=True,
            color=TITLE_COLOR,
            align=PP_ALIGN.CENTER,
        )

        # Description
        _add_textbox(
            slide,
            step["body"],
            l=x + Inches(0.15),
            t=y + Inches(1.9),
            w=col_w - Inches(0.3),
            h=h - Inches(2.1),
            font_size=13,
            color=BODY_COLOR,
            align=PP_ALIGN.CENTER,
        )

        # Connector line
        if idx < len(steps) - 1:
            cx = x + col_w
            cy = y + Inches(0.75)
            _add_rect_styled(
                slide, l=cx, t=cy, w=gap, h=Inches(0.02), fill_color=PRIMARY, line_color=PRIMARY
            )

    _set_speaker_notes(slide, data)

def _add_prof_dashboard_slide(prs, data, palette):
    from services.pptx_builder import (
        _add_textbox,
        _add_rect_styled,
        _style_shape,
        _set_center_text,
        _set_speaker_notes,
        _fill_bg,
        _VISUAL_ENGINE_AVAILABLE,
        render_chart,
    )

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, BACKGROUND)

    # Title
    title_text = data.get("title") or ""
    _add_textbox(
        slide,
        title_text,
        l=Inches(0.8),
        t=Inches(0.6),
        w=Inches(11.73),
        h=Inches(1.0),
        font_size=32,
        bold=True,
        color=TITLE_COLOR,
    )

    bullets = data.get("bullet_points", [])

    # Left Container Card (metrics + chart area)
    left_container = _add_rect_styled(
        slide,
        l=Inches(0.8),
        t=Inches(1.8),
        w=Inches(6.8),
        h=Inches(4.9),
        fill_color=BACKGROUND,
        line_color=BACKGROUND,
    )

    # Metrics
    dashboard_metrics = data.get("dashboard_metrics")
    stats = []
    if dashboard_metrics:
        for item in dashboard_metrics[:3]:
            stats.append({"label": item.get("label", ""), "value": item.get("value", "")})
    else:
        for idx, bullet in enumerate(bullets[:3]):
            stats.append(_parse_bullet_to_metric(bullet, idx))

    mini_w = Inches(1.9)
    mini_gap = Inches(0.2)

    for idx, stat in enumerate(stats):
        x = Inches(1.1) + idx * (mini_w + mini_gap)
        y = Inches(2.0)

        # Metric card
        m_card = _add_rect_styled(
            slide, l=x, t=y, w=mini_w, h=Inches(1.3), fill_color=CARD_BG, line_color=CARD_BG, rounded=True
        )
        m_card.line.color.rgb = RGBColor(220, 220, 220)
        m_card.line.width = Pt(1)

        # Value on top
        val_str = str(stat["value"])
        val_fs = 20 if len(val_str) > 10 else 24
        _add_textbox(
            slide,
            val_str,
            l=x + Inches(0.1),
            t=y + Inches(0.15),
            w=mini_w - Inches(0.2),
            h=Inches(0.4),
            font_size=val_fs,
            bold=True,
            color=TITLE_COLOR,
            align=PP_ALIGN.CENTER,
        )

        # Label on bottom
        _add_textbox(
            slide,
            stat["label"].upper(),
            l=x + Inches(0.05),
            t=y + Inches(0.6),
            w=mini_w - Inches(0.1),
            h=Inches(0.65),
            font_size=9,
            color=BODY_COLOR,
            align=PP_ALIGN.CENTER,
            fit_text=True,
        )

    # Chart Area below metrics
    chart_embedded = False
    chart_data = data.get("chart_data")
    if _VISUAL_ENGINE_AVAILABLE and chart_data:
        try:
            chart_path = render_chart(chart_data)
            if chart_path and os.path.exists(chart_path):
                slide.shapes.add_picture(
                    chart_path,
                    Inches(1.1), Inches(3.5), Inches(6.2), Inches(3.0)
                )
                chart_embedded = True
        except Exception as e:
            print(f"[PPTX] Dashboard chart embed failed: {e}")

    if not chart_embedded:
        chart_y = Inches(3.6)
        chart_title = data.get("chart_title", "Trends Analysis") if not stats else "Trends Analysis"
        _add_textbox(
            slide,
            chart_title,
            l=Inches(1.1),
            t=chart_y,
            w=Inches(6.2),
            h=Inches(0.4),
            font_size=14,
            bold=True,
            color=BODY_COLOR,
        )

        # Draw bar chart using metric labels if available, otherwise skip
        chart_items = [s["label"] for s in stats[:3]] if stats else []
        bar_colors = [PRIMARY, TITLE_COLOR, RGBColor(180, 180, 180)]
        for bar_idx, label in enumerate(chart_items):
            bx = Inches(1.5) + bar_idx * Inches(1.8)
            by = chart_y + Inches(0.8)
            bw = Inches(0.8)
            bh = Inches(1.5) - bar_idx * Inches(0.4)
            # Track
            _add_rect_styled(slide, l=bx, t=by, w=bw, h=Inches(1.5), fill_color=CARD_BG, line_color=CARD_BG, rounded=True)
            # Bar
            _add_rect_styled(slide, l=bx, t=by + Inches(1.5) - bh, w=bw, h=bh, fill_color=bar_colors[bar_idx], line_color=bar_colors[bar_idx], rounded=True)
            # Label
            _add_textbox(slide, label, l=bx - Inches(0.2), t=by + Inches(1.6), w=bw + Inches(0.4), h=Inches(0.3), font_size=11, color=BODY_COLOR, align=PP_ALIGN.CENTER)

    # Right: Executive Insights Card
    right_x = Inches(8.0)
    right_y = Inches(1.8)
    right_w = Inches(4.5)
    right_h = Inches(4.9)

    rx_card = _add_rect_styled(
        slide, l=right_x, t=right_y, w=right_w, h=right_h, fill_color=CARD_BG, line_color=CARD_BG, rounded=True
    )
    rx_card.line.color.rgb = RGBColor(220, 220, 220)
    rx_card.line.width = Pt(1)

    # Sparkle badge
    sparkle = slide.shapes.add_shape(MSO_SHAPE.OVAL, right_x + Inches(0.3), right_y + Inches(0.3), Inches(0.6), Inches(0.6))
    _style_shape(sparkle, PRIMARY, PRIMARY)
    _set_center_text(sparkle, "✦", 16, TITLE_COLOR, font_name="Calibri")

    # Title
    _add_textbox(
        slide,
        "Executive Insights",
        l=right_x + Inches(1.1),
        t=right_y + Inches(0.3),
        w=right_w - Inches(1.3),
        h=Inches(0.6),
        font_size=18,
        bold=True,
        color=TITLE_COLOR,
    )

    # Separator Line
    _add_rect_styled(
        slide,
        l=right_x + Inches(0.3),
        t=right_y + Inches(1.1),
        w=right_w - Inches(0.6),
        h=Inches(0.01),
        fill_color=RGBColor(210, 210, 210),
        line_color=RGBColor(210, 210, 210),
    )

    # Description
    dashboard_insight = data.get("dashboard_insight")
    if dashboard_insight:
        insight_text = dashboard_insight
    elif len(bullets) > 1:
        insight_text = bullets[1]
    else:
        insight_text = ""
    _add_textbox(
        slide,
        insight_text,
        l=right_x + Inches(0.3),
        t=right_y + Inches(1.3),
        w=right_w - Inches(0.6),
        h=Inches(2.6),
        font_size=13,
        color=BODY_COLOR,
    )

    # Badge/Highlight
    badge_text = stats[0]["value"] if stats else ""
    badge_rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x + Inches(0.3), right_y + Inches(4.1), right_w - Inches(0.6), Inches(0.5))
    _style_shape(badge_rect, PRIMARY, PRIMARY)
    _set_center_text(badge_rect, badge_text, 11, TITLE_COLOR, font_name="Calibri")

    _set_speaker_notes(slide, data)

def _add_prof_comparison_slide(prs, data, palette):
    from services.pptx_builder import (
        _add_textbox,
        _add_rect_styled,
        _set_speaker_notes,
        _fill_bg,
    )

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, BACKGROUND)

    # Title
    title_text = data.get("title") or ""
    _add_textbox(
        slide,
        title_text,
        l=Inches(0.8),
        t=Inches(0.6),
        w=Inches(11.73),
        h=Inches(1.0),
        font_size=32,
        bold=True,
        color=TITLE_COLOR,
    )

    comparison_items = data.get("comparison_items")
    left_title = data.get("comparison_left_title", "")
    right_title = data.get("comparison_right_title", "")
    left_points = data.get("left_points", [])
    right_points = data.get("right_points", [])

    if comparison_items:
        left_points = [item.get("left", "") for item in comparison_items]
        right_points = [item.get("right", "") for item in comparison_items]

    if not left_title or not right_title:
        if data.get("visual_items") and len(data.get("visual_items")) >= 2:
            left_title = left_title or data.get("visual_items")[0]
            right_title = right_title or data.get("visual_items")[1]
        elif not left_points and data.get("bullet_points"):
            bullets = data.get("bullet_points")
            if bullets and ":" in bullets[0]:
                left_title = left_title or bullets[0].split(":", 1)[0].strip()
            if not left_title:
                left_title = ""
            if not right_title:
                right_title = ""

    if not left_points and not comparison_items and data.get("bullet_points"):
        bullets = data.get("bullet_points")
        half = len(bullets) // 2 if len(bullets) > 1 else len(bullets)
        left_points = bullets[:half]
        right_points = bullets[half:]

    col_w = Inches(5.6)
    y = Inches(1.8)
    h = Inches(4.9)

    # Left Card
    card1 = _add_rect_styled(
        slide, l=Inches(0.8), t=y, w=col_w, h=h, fill_color=CARD_BG, line_color=CARD_BG, rounded=True
    )
    card1.line.color.rgb = RGBColor(220, 220, 220)
    card1.line.width = Pt(1)

    _add_textbox(
        slide,
        left_title,
        l=Inches(1.1),
        t=y + Inches(0.3),
        w=col_w - Inches(0.6),
        h=Inches(0.6),
        font_size=18,
        bold=True,
        color=TITLE_COLOR,
    )

    # Left Card Highlight Accent Line
    _add_rect_styled(
        slide,
        l=Inches(1.1),
        t=y + Inches(1.0),
        w=Inches(0.06),
        h=h - Inches(1.4),
        fill_color=PRIMARY,
        line_color=PRIMARY,
    )

    left_body = "\n\n".join(cleaned for cleaned in left_points)
    _add_textbox(
        slide,
        left_body,
        l=Inches(1.3),
        t=y + Inches(1.0),
        w=col_w - Inches(1.6),
        h=h - Inches(1.4),
        font_size=13,
        color=BODY_COLOR,
    )

    # Right Card
    card2 = _add_rect_styled(
        slide, l=Inches(6.8), t=y, w=col_w, h=h, fill_color=CARD_BG, line_color=CARD_BG, rounded=True
    )
    card2.line.color.rgb = RGBColor(220, 220, 220)
    card2.line.width = Pt(1)

    _add_textbox(
        slide,
        right_title,
        l=Inches(7.1),
        t=y + Inches(0.3),
        w=col_w - Inches(0.6),
        h=Inches(0.6),
        font_size=18,
        bold=True,
        color=TITLE_COLOR,
    )

    # Right Card Highlight Accent Line
    _add_rect_styled(
        slide,
        l=Inches(7.1),
        t=y + Inches(1.0),
        w=Inches(0.06),
        h=h - Inches(1.4),
        fill_color=PRIMARY,
        line_color=PRIMARY,
    )

    right_body = "\n\n".join(cleaned for cleaned in right_points)
    _add_textbox(
        slide,
        right_body,
        l=Inches(7.3),
        t=y + Inches(1.0),
        w=col_w - Inches(1.6),
        h=h - Inches(1.4),
        font_size=13,
        color=BODY_COLOR,
    )

    _set_speaker_notes(slide, data)


def _add_prof_timeline_slide(prs, data, palette):
    from services.pptx_builder import (
        _add_textbox, _add_rect_styled, _style_shape,
        _set_center_text, _set_speaker_notes, _fill_bg,
    )
    
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, BACKGROUND)
    
    title_text = data.get("title") or ""
    _add_textbox(
        slide,
        title_text,
        l=Inches(0.8),
        t=Inches(0.6),
        w=Inches(11.73),
        h=Inches(1.0),
        font_size=32,
        bold=True,
        color=TITLE_COLOR,
    )
    
    timeline_steps = data.get("timeline_steps")
    steps = []
    if timeline_steps:
        for idx, item in enumerate(timeline_steps[:5]):
            steps.append({
                "title": item.get("title") or f"Phase {idx + 1}",
                "description": item.get("description") or item.get("body") or ""
            })
    else:
        bullets = data.get("bullet_points", [])
        for idx, bullet in enumerate(bullets[:5]):
            title = f"Phase {idx + 1}"
            description = bullet
            if ":" in bullet:
                parts = bullet.split(":", 1)
                title = parts[0].strip()
                description = parts[1].strip()
            elif " - " in bullet:
                parts = bullet.split(" - ", 1)
                title = parts[0].strip()
                description = parts[1].strip()
            steps.append({"title": title, "description": description})
            
    steps_count = min(len(steps), 5)
    if steps_count == 0:
        steps_count = 1
        steps = [{"title": "Timeline", "description": ""}]
        
    total_w = Inches(11.73)
    start_x = Inches(0.8)
    gap = Inches(0.25)
    col_w = (total_w - (steps_count - 1) * gap) / steps_count
    
    line_start_x = start_x + col_w / 2
    line_end_x = start_x + (steps_count - 1) * (col_w + gap) + col_w / 2
    
    if steps_count > 1:
        _add_rect_styled(
            slide,
            l=line_start_x,
            t=Inches(3.0) - Inches(0.01),
            w=line_end_x - line_start_x,
            h=Inches(0.02),
            fill_color=PRIMARY,
            line_color=PRIMARY,
        )
        
    for idx, step in enumerate(steps[:steps_count]):
        x = start_x + idx * (col_w + gap)
        
        # Card background
        card = _add_rect_styled(slide, l=x, t=Inches(3.5), w=col_w, h=Inches(3.3), fill_color=CARD_BG, line_color=CARD_BG, rounded=True)
        card.line.color.rgb = RGBColor(220, 220, 220)
        card.line.width = Pt(1)
        
        # Yellow circle for timeline node
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + (col_w - Inches(0.6)) / 2, Inches(2.7), Inches(0.6), Inches(0.6))
        _style_shape(circle, PRIMARY, PRIMARY)
        _set_center_text(circle, f"0{idx + 1}", 11, TITLE_COLOR, font_name="Calibri")
        
        # Text inside Card
        _add_textbox(slide, step["title"], l=x + Inches(0.1), t=Inches(3.7), w=col_w - Inches(0.2), h=Inches(0.5), font_size=15, bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER)
        _add_textbox(slide, step["description"], l=x + Inches(0.1), t=Inches(4.3), w=col_w - Inches(0.2), h=Inches(2.3), font_size=12, color=BODY_COLOR, align=PP_ALIGN.CENTER)
        
    _set_speaker_notes(slide, data)


def _add_prof_architecture_slide(prs, data, palette):
    from services.pptx_builder import (
        _add_textbox, _add_rect_styled, _style_shape,
        _set_center_text, _set_speaker_notes, _fill_bg,
    )
    
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, BACKGROUND)
    
    drawio_xml = data.get("drawio_xml")
    if drawio_xml:
        try:
            from services.visual_engine import generate_drawio_diagram
            svg_path, png_path = generate_drawio_diagram(
                drawio_xml,
                topic=data.get("title") or data.get("deck_title") or "",
                disable_fallback=True
            )
            if png_path and os.path.exists(png_path):
                print(f"[PPTX] Adding PNG to professional slide with aspect ratio preservation: {png_path}")
                from PIL import Image
                with Image.open(png_path) as img:
                    img_w, img_h = img.size
                
                img_ratio = img_w / img_h
                max_w_in = 12.33
                max_h_in = 6.5
                left_margin_in = 0.5
                top_margin_in = 0.5
                
                if img_ratio > (max_w_in / max_h_in):
                    w_in = max_w_in
                    h_in = max_w_in / img_ratio
                    l_in = left_margin_in
                    t_in = top_margin_in + (max_h_in - h_in) / 2.0
                else:
                    h_in = max_h_in
                    w_in = max_h_in * img_ratio
                    t_in = top_margin_in
                    l_in = left_margin_in + (max_w_in - w_in) / 2.0
                
                print(f"[PPTX] Calculated professional dimensions: l={l_in:.2f}\", t={t_in:.2f}\", w={w_in:.2f}\", h={h_in:.2f}\"")
                slide.shapes.add_picture(
                    png_path,
                    Inches(l_in),
                    Inches(t_in),
                    Inches(w_in),
                    Inches(h_in)
                )
            else:
                print(f"[PPTX] Draw.io diagram render failed (PNG not found): {png_path}")
        except Exception as e:
            print(f"[PPTX] Draw.io diagram inline render failed: {e}")
        _set_speaker_notes(slide, data)
        return
    
    architecture_nodes = data.get("architecture_nodes")
    nodes = []
    if architecture_nodes:
        for idx, item in enumerate(architecture_nodes[:6]):
            nodes.append({
                "title": item.get("title") or f"Component {idx + 1}",
                "description": item.get("description") or item.get("body") or ""
            })
    else:
        bullets = data.get("bullet_points", [])
        for idx, bullet in enumerate(bullets[:6]):
            title = f"Component {idx + 1}"
            description = bullet
            if ":" in bullet:
                parts = bullet.split(":", 1)
                title = parts[0].strip()
                description = parts[1].strip()
            elif " - " in bullet:
                parts = bullet.split(" - ", 1)
                title = parts[0].strip()
                description = parts[1].strip()
            nodes.append({"title": title, "description": description})
            
    nodes_count = len(nodes)
    if nodes_count == 0:
        nodes_count = 1
        nodes = [{"title": "Architecture Node", "description": ""}]
        
    col_w = Inches(3.5)
    row_h = Inches(2.2)
    gap_x = Inches(0.61)
    gap_y = Inches(0.4)
    start_x = Inches(0.8)
    start_y = Inches(1.9)
    
    # Draw horizontal and vertical connectors first so they sit behind cards
    for idx, node in enumerate(nodes):
        row = idx // 3
        col = idx % 3
        x = start_x + col * (col_w + gap_x)
        y = start_y + row * (row_h + gap_y)
        
        # Horizontal connector line
        if col < 2 and idx + 1 < nodes_count:
            cx = x + col_w
            cy = y + row_h / 2
            _add_rect_styled(slide, l=cx, t=cy - Inches(0.01), w=gap_x, h=Inches(0.02), fill_color=PRIMARY, line_color=PRIMARY)
            # Yellow indicator dot on horizontal connector
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx + gap_x/2 - Inches(0.1), cy - Inches(0.1), Inches(0.2), Inches(0.2))
            _style_shape(dot, PRIMARY, PRIMARY)
            
        # Vertical connector line
        if row == 0 and idx + 3 < nodes_count:
            cx = x + col_w / 2
            cy = y + row_h
            _add_rect_styled(slide, l=cx - Inches(0.01), t=cy, w=Inches(0.02), h=gap_y, fill_color=PRIMARY, line_color=PRIMARY)
            # Yellow indicator dot on vertical connector
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.1), cy + gap_y/2 - Inches(0.1), Inches(0.2), Inches(0.2))
            _style_shape(dot, PRIMARY, PRIMARY)
            
    # Draw cards on top of connectors
    for idx, node in enumerate(nodes):
        row = idx // 3
        col = idx % 3
        x = start_x + col * (col_w + gap_x)
        y = start_y + row * (row_h + gap_y)
        
        # Card background
        card = _add_rect_styled(slide, l=x, t=y, w=col_w, h=row_h, fill_color=CARD_BG, line_color=CARD_BG, rounded=True)
        card.line.color.rgb = RGBColor(220, 220, 220)
        card.line.width = Pt(1)
        
        # Top yellow accent block
        _add_rect_styled(slide, l=x + Inches(0.1), t=y, w=col_w - Inches(0.2), h=Inches(0.08), fill_color=PRIMARY, line_color=PRIMARY)
        
        # Node title and description
        _add_textbox(slide, node["title"], l=x + Inches(0.3), t=y + Inches(0.2), w=col_w - Inches(0.6), h=Inches(0.4), font_size=15, bold=True, color=TITLE_COLOR)
        _add_textbox(slide, node["description"], l=x + Inches(0.3), t=y + Inches(0.7), w=col_w - Inches(0.6), h=row_h - Inches(0.9), font_size=12, color=BODY_COLOR)
        
    _set_speaker_notes(slide, data)

