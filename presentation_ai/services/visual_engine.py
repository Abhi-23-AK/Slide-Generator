"""
Visual Generation Engine
========================
Provides visual assets (images, icons, diagrams, charts, illustrations, tables)
for embedding into PPTX slides.  Every public function is crash-safe — failures
are logged and return None (except get_icon, which always returns an SVG string).
"""
import os
os.environ["TEMP"] = "E:\\Slide_Generator\\temp"
os.environ["TMP"] = "E:\\Slide_Generator\\temp"
os.makedirs("E:\\Slide_Generator\\temp", exist_ok=True)

import uuid
import tempfile
import subprocess
import difflib
import hashlib
import sys
import io
import re
from typing import Optional, Dict, List, Tuple
from concurrent.futures import ThreadPoolExecutor, as_completed

import requests
from dotenv import load_dotenv

# Import configuration
from config import DRAWIO_CLI_PATH

# Fix Windows cp1252 encoding issues with Unicode
if sys.stdout.encoding != 'utf-8':
    sys.stdout = io.TextIOWrapper(
        sys.stdout.buffer,
        encoding='utf-8',
        errors='replace'
    )

load_dotenv()

# ─── API keys ────────────────────────────────────────────────────────
UNSPLASH_ACCESS_KEY = os.getenv("UNSPLASH_ACCESS_KEY", "")
PEXELS_API_KEY = os.getenv("PEXELS_API_KEY", "")
PIXABAY_API_KEY = os.getenv("PIXABAY_API_KEY", "")

# ─── Temp directory ──────────────────────────────────────────────────
TMP_DIR = tempfile.gettempdir()

import threading
_CACHE_LOCK = threading.Lock()
DIAGRAM_CACHE = {}  # xml_hash -> (svg_path, png_path) - Thread-safe cache enabled


# =====================================================================
# 1. IMAGES  (Waterfall: Unsplash → Pexels → Pixabay)
# =====================================================================

def _download_image(url: str, headers: dict = None) -> Optional[str]:
    """Download an image URL to /tmp/ and return the local file path."""
    try:
        resp = requests.get(url, headers=headers, timeout=15, stream=True)
        resp.raise_for_status()
        ext = ".jpg"
        content_type = resp.headers.get("Content-Type", "")
        if "png" in content_type:
            ext = ".png"
        elif "webp" in content_type:
            ext = ".webp"
        path = os.path.join(TMP_DIR, f"hero_{uuid.uuid4().hex[:10]}{ext}")
        with open(path, "wb") as f:
            for chunk in resp.iter_content(8192):
                f.write(chunk)
        print(f"[VISUAL] Image downloaded -> {path}")
        return path
    except Exception as e:
        print(f"[VISUAL] Image download failed: {e}")
        return None


def _try_unsplash(keyword: str) -> Optional[str]:
    """Attempt to fetch a hero image from Unsplash."""
    if not UNSPLASH_ACCESS_KEY:
        print("[VISUAL] Unsplash: no API key configured, skipping")
        return None
    try:
        resp = requests.get(
            "https://api.unsplash.com/search/photos",
            params={"query": keyword, "per_page": 1, "orientation": "landscape"},
            headers={"Authorization": f"Client-ID {UNSPLASH_ACCESS_KEY}"},
            timeout=10,
        )
        # Check rate limit before processing
        remaining = resp.headers.get("X-Ratelimit-Remaining", "1")
        if remaining == "0":
            print("[VISUAL] Unsplash: rate limit exhausted, skipping")
            return None
        resp.raise_for_status()
        data = resp.json()
        results = data.get("results", [])
        if not results:
            print("[VISUAL] Unsplash: no results found")
            return None
        img_url = results[0].get("urls", {}).get("regular")
        if not img_url:
            return None
        print(f"[VISUAL] Unsplash: found image for '{keyword}'")
        return _download_image(img_url)
    except Exception as e:
        print(f"[VISUAL] Unsplash error: {e}")
        return None


def _try_pexels(keyword: str) -> Optional[str]:
    """Attempt to fetch a hero image from Pexels."""
    if not PEXELS_API_KEY:
        print("[VISUAL] Pexels: no API key configured, skipping")
        return None
    try:
        resp = requests.get(
            "https://api.pexels.com/v1/search",
            params={"query": keyword, "per_page": 1, "orientation": "landscape"},
            headers={"Authorization": PEXELS_API_KEY},
            timeout=10,
        )
        remaining = resp.headers.get("X-Ratelimit-Remaining", "1")
        if remaining == "0":
            print("[VISUAL] Pexels: rate limit exhausted, skipping")
            return None
        resp.raise_for_status()
        data = resp.json()
        photos = data.get("photos", [])
        if not photos:
            print("[VISUAL] Pexels: no results found")
            return None
        img_url = photos[0].get("src", {}).get("large2x") or photos[0].get("src", {}).get("large")
        if not img_url:
            return None
        print(f"[VISUAL] Pexels: found image for '{keyword}'")
        return _download_image(img_url)
    except Exception as e:
        print(f"[VISUAL] Pexels error: {e}")
        return None


def _try_pixabay(keyword: str) -> Optional[str]:
    """Attempt to fetch a hero image from Pixabay (final fallback)."""
    if not PIXABAY_API_KEY:
        print("[VISUAL] Pixabay: no API key configured, skipping")
        return None
    try:
        resp = requests.get(
            "https://pixabay.com/api/",
            params={
                "key": PIXABAY_API_KEY,
                "q": keyword,
                "per_page": 3,
                "orientation": "horizontal",
                "image_type": "photo",
                "safesearch": "true",
            },
            timeout=10,
        )
        remaining = resp.headers.get("X-Ratelimit-Remaining", "1")
        if remaining == "0":
            print("[VISUAL] Pixabay: rate limit exhausted, skipping")
            return None
        resp.raise_for_status()
        data = resp.json()
        hits = data.get("hits", [])
        if not hits:
            print("[VISUAL] Pixabay: no results found")
            return None
        img_url = hits[0].get("largeImageURL") or hits[0].get("webformatURL")
        if not img_url:
            return None
        print(f"[VISUAL] Pixabay: found image for '{keyword}'")
        return _download_image(img_url)
    except Exception as e:
        print(f"[VISUAL] Pixabay error: {e}")
        return None


def get_hero_image(keyword: str) -> Optional[str]:
    """
    Fetch a hero image for the given keyword using a waterfall strategy:
    Unsplash → Pexels → Pixabay.
    Returns a local file path or None.  Never raises an exception.
    """
    try:
        print(f"[VISUAL] Searching hero image for: '{keyword}'")
        # Waterfall: try each source in priority order
        path = _try_unsplash(keyword)
        if path:
            return path

        path = _try_pexels(keyword)
        if path:
            return path

        path = _try_pixabay(keyword)
        if path:
            return path

        print(f"[VISUAL] No hero image found for '{keyword}' from any source")
        return None
    except Exception as e:
        print(f"[VISUAL] get_hero_image unexpected error: {e}")
        return None


def download_all_images(slides: list) -> dict:
    """Download all slide images in parallel"""
    import time
    image_map = {}  # slide_index → local_path
    
    def download_one(index, keyword):
        try:
            path = get_hero_image(keyword)
            return index, path
        except:
            return index, None
    
    # Collect all keywords
    tasks = []
    for i, slide in enumerate(slides):
        keyword = slide.get('image_keyword', '')
        if keyword:
            tasks.append((i, keyword))
    
    if not tasks:
        return image_map
    
    # Download all at once in parallel
    print(f"[TIMING] Starting parallel download of {len(tasks)} images...")
    start_time = time.time()
    
    with ThreadPoolExecutor(max_workers=5) as executor:
        futures = {
            executor.submit(download_one, i, kw): i 
            for i, kw in tasks
        }
        for future in as_completed(futures):
            index, path = future.result()
            image_map[index] = path
            print(f"[VISUAL] OK Image {index+1} ready")
    
    elapsed = time.time() - start_time
    print(f"[TIMING] Images downloaded in {elapsed:.2f}s")
    return image_map


# =====================================================================
# 2. ICONS  (lucide-py — always returns SVG, never None)
# =====================================================================

# Keyword → Lucide icon name mapping
ICON_KEYWORD_MAP = {
    "technology":  "cpu",
    "data":        "database",
    "chart":       "bar-chart-3",
    "people":      "users",
    "money":       "dollar-sign",
    "time":        "clock",
    "document":    "file-text",
    "success":     "check-circle",
    "warning":     "alert-triangle",
    "idea":        "lightbulb",
    "cloud":       "cloud",
    "security":    "shield",
    "growth":      "trending-up",
    "ai":          "brain",
    "process":     "git-branch",
    "education":   "graduation-cap",
    "health":      "heart-pulse",
    "rocket":      "rocket",
    "globe":       "globe",
    "settings":    "settings",
    "mail":        "mail",
    "phone":       "phone",
    "home":        "home",
    "search":      "search",
    "star":        "star",
    "heart":       "heart",
    "code":        "code",
    "server":      "server",
    "database":    "database",
    "network":     "network",
    "analytics":   "line-chart",
    "team":        "users",
    "target":      "target",
    "award":       "award",
}

_FALLBACK_CIRCLE_SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" width="24" height="24" '
    'viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" '
    'stroke-linecap="round" stroke-linejoin="round">'
    '<circle cx="12" cy="12" r="10"/>'
    '</svg>'
)

# Inline SVG paths for standard Lucide icons as fallback if lucide-py is unavailable
_INLINE_ICONS = {
    "cpu": '<rect x="4" y="4" width="16" height="16" rx="2" ry="2"></rect><rect x="9" y="9" width="6" height="6"></rect><line x1="9" y1="1" x2="9" y2="4"></line><line x1="15" y1="1" x2="15" y2="4"></line><line x1="9" y1="20" x2="9" y2="23"></line><line x1="15" y1="20" x2="15" y2="23"></line><line x1="20" y1="9" x2="23" y2="9"></line><line x1="20" y1="15" x2="23" y2="15"></line><line x1="1" y1="9" x2="4" y2="9"></line><line x1="1" y1="15" x2="4" y2="15"></line>',
    "database": '<ellipse cx="12" cy="5" rx="9" ry="3"></ellipse><path d="M3 5v14c0 1.66 4 3 9 3s9-1.34 9-3V5"></path><path d="M3 12c0 1.66 4 3 9 3s9-1.34 9-3"></path>',
    "bar-chart-3": '<line x1="18" y1="20" x2="18" y2="10"></line><line x1="12" y1="20" x2="12" y2="4"></line><line x1="6" y1="20" x2="6" y2="14"></line><line x1="2" y1="20" x2="22" y2="20"></line>',
    "users": '<path d="M17 21v-2a4 4 0 0 0-4-4H5a4 4 0 0 0-4 4v2"></path><circle cx="9" cy="7" r="4"></circle><path d="M23 21v-2a4 4 0 0 0-3-3.87"></path><path d="M16 3.13a4 4 0 0 1 0 7.75"></path>',
    "dollar-sign": '<line x1="12" y1="1" x2="12" y2="23"></line><path d="M17 5H9.5a3.5 3.5 0 0 0 0 7h5a3.5 3.5 0 0 1 0 7H6"></path>',
    "clock": '<circle cx="12" cy="12" r="10"></circle><polyline points="12 6 12 12 16 14"></polyline>',
    "file-text": '<path d="M14 2H6a2 2 0 0 0-2 2v16a2 2 0 0 0 2 2h12a2 2 0 0 0 2-2V8z"></path><polyline points="14 2 14 8 20 8"></polyline><line x1="16" y1="13" x2="8" y2="13"></line><line x1="16" y1="17" x2="8" y2="17"></line><polyline points="10 9 9 9 8 9"></polyline>',
    "check-circle": '<path d="M22 11.08V12a10 10 0 1 1-5.93-9.14"></path><polyline points="22 4 12 14.01 9 11.01"></polyline>',
    "alert-triangle": '<path d="M10.29 3.86L1.82 18a2 2 0 0 0 1.71 3h16.94a2 2 0 0 0 1.71-3L13.71 3.86a2 2 0 0 0-3.42 0z"></path><line x1="12" y1="9" x2="12" y2="13"></line><line x1="12" y1="17" x2="12.01" y2="17"></line>',
    "lightbulb": '<path d="M15 14c.2-1 .7-1.7 1.5-2.5 1-.9 1.5-2.2 1.5-3.5A6 6 0 0 0 6 8c0 1 .6 2.2 1.5 3.5.7.7 1.3 1.5 1.5 2.5"></path><path d="M9 18h6"></path><path d="M10 22h4"></path>',
    "cloud": '<path d="M18 10h-1.26A8 8 0 1 0 9 20h9a5 5 0 0 0 0-10z"></path>',
    "shield": '<path d="M12 22s8-4 8-10V5l-8-3-8 3v7c0 6 8 10 8 10z"></path>',
    "trending-up": '<polyline points="23 6 13.5 15.5 8.5 10.5 1 18"></polyline><polyline points="17 6 23 6 23 12"></polyline>',
    "brain": '<path d="M9.5 2A2.5 2.5 0 0 1 12 4.5v15a2.5 2.5 0 0 1-4.96-.44 2.5 2.5 0 0 1 0-3.12 3 3 0 0 1 0-4.88 2.5 2.5 0 0 1 0-3.12A2.5 2.5 0 0 1 9.5 2zM14.5 2A2.5 2.5 0 0 0 12 4.5v15a2.5 2.5 0 0 0 4.96-.44 2.5 2.5 0 0 0 0-3.12 3 3 0 0 0 0-4.88 2.5 2.5 0 0 0 0-3.12A2.5 2.5 0 0 0 14.5 2z"></path>',
    "git-branch": '<line x1="6" y1="3" x2="6" y2="15"></line><circle cx="18" cy="6" r="3"></circle><circle cx="6" cy="18" r="3"></circle><path d="M18 9a9 9 0 0 1-9 9"></path>',
    "graduation-cap": '<path d="M22 10v6M2 10l10-5 10 5-10 5z"></path><path d="M6 12v5c0 2 2 3 6 3s6-1 6-3v-5"></path>',
    "heart-pulse": '<path d="M20.42 4.58a5.4 5.4 0 0 0-7.65 0l-.77.78-.77-.78a5.4 5.4 0 0 0-7.65 0 5.4 5.4 0 0 0 0 7.65L12 20l8.42-8.42a5.4 5.4 0 0 0 0-7.58z"></path><path d="M3.5 12h2.5l2.5-4 2.5 8 2.5-6 1.5 2h3"></path>',
    "rocket": '<path d="M4.5 16.5c-1.5 1.26-2 2.5-2 2.5s1.24-.5 2.5-2L9 13.5l-4.5-4.5-4.5 4.5 4.5 3zM21.5 2.5s-2 2.5-4.5 4.5L12.5 2.5l-4.5 4.5 4.5 4.5-4.5 4.5c-2.5 2.5-4.5 4.5-4.5 4.5s2-2 4.5-4.5L12.5 12.5l4.5 4.5-4.5-4.5 4.5-4.5z"></path>',
    "globe": '<circle cx="12" cy="12" r="10"></circle><line x1="2" y1="12" x2="22" y2="12"></line><path d="M12 2a15.3 15.3 0 0 1 4 10 15.3 15.3 0 0 1-4 10 15.3 15.3 0 0 1-4-10 15.3 15.3 0 0 1 4-10z"></path>',
    "settings": '<circle cx="12" cy="12" r="3"></circle><path d="M19.4 15a1.65 1.65 0 0 0 .33 1.82l.06.06a2 2 0 1 1-2.83 2.83l-.06-.06a1.65 1.65 0 0 0-1.82-.33 1.65 1.65 0 0 0-1 1.51V21a2 2 0 0 1-4 0v-.09A1.65 1.65 0 0 0 9 19.4a1.65 1.65 0 0 0-1.82.33l-.06.06a2 2 0 1 1-2.83-2.83l.06-.06a1.65 1.65 0 0 0 .33-1.82 1.65 1.65 0 0 0-1.51-1H3a2 2 0 0 1 0-4h.09A1.65 1.65 0 0 0 4.6 9a1.65 1.65 0 0 0-.33-1.82l-.06-.06a2 2 0 1 1 2.83-2.83l.06.06a1.65 1.65 0 0 0 1.82.33H9a1.65 1.65 0 0 0 1-1.51V3a2 2 0 0 1 4 0v.09a1.65 1.65 0 0 0 1 1.51 1.65 1.65 0 0 0 1.82-.33l.06-.06a2 2 0 1 1 2.83 2.83l-.06.06a1.65 1.65 0 0 0-.33 1.82V9a1.65 1.65 0 0 0 1.51 1H21a2 2 0 0 1 0 4h-.09a1.65 1.65 0 0 0-1.51 1z"></path>',
    "mail": '<path d="M4 4h16c1.1 0 2 .9 2 2v12c0 1.1-.9 2-2 2H4c-1.1 0-2-.9-2-2V6c0-1.1.9-2 2-2z"></path><polyline points="22,6 12,13 2,6"></polyline>',
    "phone": '<path d="M22 16.92v3a2 2 0 0 1-2.18 2 19.79 19.79 0 0 1-8.63-3.07 19.5 19.5 0 0 1-6-6 19.79 19.79 0 0 1-3.07-8.67A2 2 0 0 1 4.11 2h3a2 2 0 0 1 2 1.72 12.84 12.84 0 0 0 .7 2.81 2 2 0 0 1-.45 2.11L8.09 9.91a16 16 0 0 0 6 6l1.27-1.27a2 2 0 0 1 2.11-.45 12.84 12.84 0 0 0 2.81.7A2 2 0 0 1 22 16.92z"></path>',
    "home": '<path d="M3 9l9-7 9 7v11a2 2 0 0 1-2 2H5a2 2 0 0 1-2-2z"></path><polyline points="9 22 9 12 15 12 15 22"></polyline>',
    "search": '<circle cx="11" cy="11" r="8"></circle><line x1="21" y1="21" x2="16.65" y2="16.65"></line>',
    "star": '<polygon points="12 2 15.09 8.26 22 9.27 17 14.14 18.18 21.02 12 17.77 5.82 21.02 7 14.14 2 9.27 8.91 8.26 12 2"></polygon>',
    "heart": '<path d="M20.84 4.61a5.5 5.5 0 0 0-7.78 0L12 5.67l-1.06-1.06a5.5 5.5 0 0 0-7.78 7.78l1.06 1.06L12 21.23l7.78-7.78 1.06-1.06a5.5 5.5 0 0 0 0-7.78z"></path>',
    "code": '<polyline points="16 18 22 12 16 6"></polyline><polyline points="8 6 2 12 8 18"></polyline>',
    "server": '<rect x="2" y="2" width="20" height="8" rx="2" ry="2"></rect><rect x="2" y="14" width="20" height="8" rx="2" ry="2"></rect><line x1="6" y1="6" x2="6.01" y2="6"></line><line x1="6" y1="18" x2="6.01" y2="18"></line>',
    "network": '<rect x="16" y="16" width="6" height="6" rx="1"></rect><rect x="2" y="16" width="6" height="6" rx="1"></rect><rect x="9" y="2" width="6" height="6" rx="1"></rect><path d="M12 8v8M5 16v-3a1 1 0 0 1 1-1h12a1 1 0 0 1 1 1v3"></path>',
    "line-chart": '<path d="M3 3v18h18"></path><path d="M18.7 8l-5.1 5.2-2.8-2.7L7 14.3"></path>',
    "target": '<circle cx="12" cy="12" r="10"></circle><circle cx="12" cy="12" r="6"></circle><circle cx="12" cy="12" r="2"></circle>',
    "award": '<circle cx="12" cy="8" r="7"></circle><polyline points="8.21 13.89 7 23 12 20 17 23 15.79 13.88"></polyline>'
}


def _resolve_icon_name(keyword: str) -> str:
    """Resolve a keyword to a Lucide icon name via exact or fuzzy match."""
    kw = keyword.lower().strip()

    # Exact match
    if kw in ICON_KEYWORD_MAP:
        return ICON_KEYWORD_MAP[kw]

    # Fuzzy match against mapping keys
    matches = difflib.get_close_matches(kw, ICON_KEYWORD_MAP.keys(), n=1, cutoff=0.4)
    if matches:
        return ICON_KEYWORD_MAP[matches[0]]

    # Try checking if the keyword itself is a valid icon name (pass-through)
    return kw


def get_icon(keyword: str) -> str:
    """
    Return an SVG string for the given keyword using lucide-py.
    Falls back to an inline SVG database, or a simple circle SVG. **Never returns None.**
    """
    try:
        print(f"[VISUAL] Getting icon for: '{keyword}'")
        icon_name = _resolve_icon_name(keyword)

        try:
            from lucide_py import Icon
            svg_str = Icon(icon_name).svg()
            if svg_str:
                print(f"[VISUAL] Icon found: '{icon_name}' via lucide-py")
                return svg_str
        except Exception as e:
            # Fall back to inline icons database
            if icon_name in _INLINE_ICONS:
                print(f"[VISUAL] Icon found: '{icon_name}' via inline cache")
                return (
                    '<svg xmlns="http://www.w3.org/2000/svg" width="24" height="24" '
                    'viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" '
                    f'stroke-linecap="round" stroke-linejoin="round">{_INLINE_ICONS[icon_name]}</svg>'
                )
            print(f"[VISUAL] lucide-py lookup failed for '{icon_name}': {e}")

        # If everything failed, try matching custom SVG or return standard fallback
        print(f"[VISUAL] Returning fallback circle SVG for '{keyword}'")
        return _FALLBACK_CIRCLE_SVG

    except Exception as e:
        print(f"[VISUAL] get_icon unexpected error: {e}")
        return _FALLBACK_CIRCLE_SVG



# =====================================================================
# 3. DIAGRAMS  (mermaid-cli / mmdc)
# =====================================================================

def generate_diagram(mermaid_code: str) -> Optional[str]:
    """
    Render a Mermaid diagram to SVG using mmdc (mermaid-cli).
    Returns the SVG file path on success, None on failure.
    """
    uid = uuid.uuid4().hex[:10]
    mmd_path = os.path.join(TMP_DIR, f"diagram_{uid}.mmd")
    svg_path = os.path.join(TMP_DIR, f"diagram_{uid}.svg")

    try:
        print(f"[VISUAL] Generating mermaid diagram ({len(mermaid_code)} chars)")

        # Write mermaid source
        with open(mmd_path, "w", encoding="utf-8") as f:
            f.write(mermaid_code)

        # Render via mmdc
        result = subprocess.run(
            ["mmdc", "-i", mmd_path, "-o", svg_path, "-b", "transparent", "-w", "1200"],
            capture_output=True,
            text=True,
            timeout=15,
        )

        if result.returncode != 0:
            print(f"[VISUAL] mmdc failed (exit {result.returncode}): {result.stderr[:300]}")
            return None

        if not os.path.exists(svg_path):
            print("[VISUAL] mmdc ran but SVG file was not created")
            return None

        print(f"[VISUAL] Mermaid diagram rendered -> {svg_path}")
        return svg_path

    except FileNotFoundError:
        print("[VISUAL] mmdc (mermaid-cli) not installed — skipping diagram")
        return None
    except subprocess.TimeoutExpired:
        print("[VISUAL] mmdc timed out after 15 seconds")
        return None
    except Exception as e:
        print(f"[VISUAL] generate_diagram unexpected error: {e}")
        return None
    finally:
        # Clean up the .mmd temp file
        try:
            if os.path.exists(mmd_path):
                os.remove(mmd_path)
        except OSError:
            pass


def validate_drawio_complexity(xml: str, topology: str = None) -> dict:
    """
    Validate the complexity of a Draw.io diagram XML string.
    Returns a dictionary of counts and whether the diagram is sufficiently complex.
    Applies topology-aware validation rules.
    NOTE: This function is kept for validation but NOT used for V3 XML regeneration.
    """
    cell_count = 0
    arrow_count = 0
    nested_count = 0
    
    try:
        import xml.etree.ElementTree as ET
        xml_to_parse = xml.strip()
        if xml_to_parse:
            if not (xml_to_parse.startswith("<mxfile") or xml_to_parse.startswith("<mxGraphModel") or xml_to_parse.startswith("<?xml")):
                xml_to_parse = f"<root_wrapper>{xml_to_parse}</root_wrapper>"
            
            root = ET.fromstring(xml_to_parse)
            cells = root.findall(".//mxCell")
            cell_count = len(cells)
            
            for cell in cells:
                if cell.get("edge") == "1":
                    arrow_count += 1
                parent = cell.get("parent")
                if parent not in (None, "", "0", "1"):
                    nested_count += 1
    except Exception as e:
        print(f"[VISUAL] XML complexity validation parser error (falling back to regex): {e}")
        # Regex fallback
        cell_count = len(re.findall(r'<mxCell\b', xml))
        arrow_count = len(re.findall(r'\bedge=["\']1["\']', xml))
        parents = re.findall(r'\bparent=["\']([^"\']+)["\']', xml)
        nested_count = sum(1 for p in parents if p not in ("0", "1"))
        
    topo = str(topology).lower().strip() if topology else ""
    
    # Define topology-aware thresholds
    if "cnn" in topo or "transformer" in topo:
        min_cells = 8
        min_arrows = 1
        min_nested = 0
    elif "rag" in topo:
        min_cells = 10
        min_arrows = 2
        min_nested = 0
    else:
        min_cells = 15
        min_arrows = 3
        min_nested = 2
        
    is_complex = (
        cell_count >= min_cells and
        arrow_count >= min_arrows and
        nested_count >= min_nested
    )
    return {
        "cell_count": cell_count,
        "arrow_count": arrow_count,
        "nested_count": nested_count,
        "is_complex": is_complex
    }


def validate_diagram_png(png_path: str, drawio_xml: str, topology: str = None) -> Tuple[bool, str]:
    """
    Validates the rendered diagram PNG.
    Returns (is_valid, error_reason).
    """
    import os
    if not os.path.exists(png_path):
        return False, "PNG file does not exist"
        
    # Check 1: PNG file size < 20KB = blank/failed render
    try:
        sz = os.path.getsize(png_path)
        if sz < 20 * 1024:
            return False, f"PNG file size too small: {sz / 1024:.1f}KB (under 20KB)"
    except Exception as e:
        return False, f"Failed to check PNG size: {e}"
        
    # Check 2: Image is mostly one color = rendering error
    try:
        from PIL import Image
        with Image.open(png_path) as img:
            rgb_img = img.convert("RGB")
            w, h = rgb_img.size
            total_pixels = w * h
            
            colors = rgb_img.getcolors(maxcolors=total_pixels)
            if colors:
                colors_sorted = sorted(colors, key=lambda x: x[0], reverse=True)
                most_frequent_count = colors_sorted[0][0]
                ratio = most_frequent_count / total_pixels
                if ratio > 0.98:
                    return False, f"Image is mostly one color ({ratio*100:.1f}% dominated by single color)"
    except Exception as e:
        print(f"[VISUAL] PIL color check warning: {e}")
        
    # Check 4: Re-read XML and count cells - if not complex enough for topology, flag as too simple
    complexity = validate_drawio_complexity(drawio_xml, topology)
    if not complexity["is_complex"]:
        return False, f"XML complexity validation failed for topology '{topology}': cells={complexity['cell_count']}, arrows={complexity['arrow_count']}, nested={complexity['nested_count']}"
        
    # Check 3: Node count in XML vs visual density mismatch
    xml_cells = complexity["cell_count"]
    if xml_cells > 20 and os.path.getsize(png_path) < 30 * 1024:
         return False, f"Visual density mismatch: {xml_cells} XML cells but PNG is only {os.path.getsize(png_path)/1024:.1f}KB"

    return True, ""


def _render_xml_to_files(current_xml: str) -> Tuple[Optional[str], Optional[str]]:
    """
    Helper to compile a Draw.io XML string to SVG and PNG.
    """
    uid = uuid.uuid4().hex[:10]
    drawio_path = os.path.join(TMP_DIR, f"diagram_{uid}.drawio")
    svg_path = os.path.join(TMP_DIR, f"diagram_{uid}.svg")
    png_path = svg_path.replace(".svg", ".png")

    try:
        wrapped_xml = current_xml.strip()
        if not (wrapped_xml.startswith("<mxfile") or wrapped_xml.startswith("<?xml")):
            wrapped_xml = (
                f'<mxfile host="Electron" version="30.0.4">\n'
                f'  <diagram id="page_{uid}" name="Page-1">\n'
                f'    {current_xml}\n'
                f'  </diagram>\n'
                f'</mxfile>'
            )

        with open(drawio_path, "w", encoding="utf-8") as f:
            f.write(wrapped_xml)

        cli_path = DRAWIO_CLI_PATH
        if not os.path.exists(cli_path):
            print(f"[VISUAL] Draw.io CLI not found at: {cli_path}")
            return None, None

        # Render SVG
        res_svg = subprocess.run(
            [
                cli_path,
                "-x",
                "-f", "svg",
                "-o", svg_path,
                drawio_path,
                "--no-sandbox",
                "--disable-gpu",
                "--disable-software-rasterizer",
                "--headless"
            ],
            capture_output=True,
            text=True,
            timeout=120
        )
        print(f"[VISUAL] SVG Render Exit Code: {res_svg.returncode}")
        if res_svg.returncode != 0 or not os.path.exists(svg_path):
            print(f"[VISUAL] SVG Render Stdout: {res_svg.stdout.strip()}")
            print(f"[VISUAL] SVG Render Stderr: {res_svg.stderr.strip()}")

        if not os.path.exists(svg_path):
            print("[VISUAL] Draw.io CLI completed but SVG file was not created")
            return None, None

        # Render PNG (HD borderless cropping using scale=3, no fixed dimensions)
        cmd_png = [
            cli_path,
            '-x',
            '-f', 'png',
            '-s', '3',
            '--border', '20',
            '-o', png_path,
            drawio_path,
            '--no-sandbox',
            '--disable-gpu',
            '--disable-software-rasterizer',
            '--headless'
        ]
        
        res_png = subprocess.run(
            cmd_png,
            capture_output=True,
            timeout=120,
            text=True
        )
        print(f"[VISUAL] PNG Render Exit Code: {res_png.returncode}")
        if res_png.returncode != 0:
            print(f"[VISUAL] PNG Render Stdout: {res_png.stdout.strip()}")
            print(f"[VISUAL] PNG Render Stderr: {res_png.stderr.strip()}")
            print("[VISUAL] PNG Render with scale=3 failed. Retrying with scale=2...")
            cmd_png_s2 = list(cmd_png)
            for idx, arg in enumerate(cmd_png_s2):
                if arg == '-s':
                    cmd_png_s2[idx + 1] = '2'
                    break
            res_png = subprocess.run(
                cmd_png_s2,
                capture_output=True,
                timeout=120,
                text=True
            )
            print(f"[VISUAL] PNG Render (scale=2) Exit Code: {res_png.returncode}")
            
            if res_png.returncode != 0:
                print(f"[VISUAL] PNG Render (scale=2) Stdout: {res_png.stdout.strip()}")
                print(f"[VISUAL] PNG Render (scale=2) Stderr: {res_png.stderr.strip()}")
                print("[VISUAL] PNG Render with scale=2 failed. Retrying with scale=1...")
                cmd_png_s1 = list(cmd_png)
                for idx, arg in enumerate(cmd_png_s1):
                    if arg == '-s':
                        cmd_png_s1[idx + 1] = '1'
                        break
                res_png = subprocess.run(
                    cmd_png_s1,
                    capture_output=True,
                    timeout=120,
                    text=True
                )
                print(f"[VISUAL] PNG Render (scale=1) Exit Code: {res_png.returncode}")
                if res_png.returncode != 0:
                    print(f"[VISUAL] PNG Render (scale=1) Stdout: {res_png.stdout.strip()}")
                    print(f"[VISUAL] PNG Render (scale=1) Stderr: {res_png.stderr.strip()}")

        # SVG to PNG conversion check
        if not os.path.exists(png_path) and os.path.exists(svg_path):
            print("[VISUAL] PNG missing, converting SVG→PNG...")
            try:
                import cairosvg
                cairosvg.svg2png(
                    url=svg_path, 
                    write_to=png_path,
                    scale=3.0
                )
                print(f"[VISUAL] SVG→PNG conversion successful")
            except Exception as cairo_err:
                print(f"[VISUAL] cairosvg conversion failed: {cairo_err}")
                try:
                    # Try newer Inkscape v1.0+ syntax
                    subprocess.run(['inkscape', svg_path, 
                                  '-o', png_path,
                                  '--export-dpi', '300'],
                                  capture_output=True, timeout=60)
                except Exception as inkscape_err:
                    print(f"[VISUAL] Inkscape (v1+) conversion failed: {inkscape_err}")
                    try:
                        # Try older Inkscape syntax
                        subprocess.run(['inkscape', svg_path, 
                                      '--export-png', png_path,
                                      '--export-dpi', '300'],
                                      capture_output=True, timeout=60)
                    except Exception as old_inkscape_err:
                        print(f"[VISUAL] Inkscape (legacy) conversion failed: {old_inkscape_err}")

        if os.path.exists(png_path):
            return svg_path, png_path
            
        return None, None

    except Exception as e:
        print(f"[VISUAL] Error rendering XML: {e}")
        return None, None
    finally:
        try:
            if os.path.exists(drawio_path):
                os.remove(drawio_path)
        except OSError:
            pass


def print_diagram_diagnostics(xml: str, svg_path: str, png_path: str):
    """
    Computes and prints layout metrics (dimensions, node occupancy, and counts) for the diagram.
    """
    try:
        import re
        from PIL import Image
        
        # 1. Parse XML to count actual nodes and edges
        vertices = re.findall(r'<mxCell[^>]+vertex="1"[^>]*>', xml)
        edges = re.findall(r'<mxCell[^>]+edge="1"[^>]*>', xml)
        
        nodes = []
        containers = []
        for v in vertices:
            vid_m = re.search(r'id="([^"]+)"', v)
            if vid_m and vid_m.group(1) not in ("0", "1"):
                style_m = re.search(r'style="([^"]+)"', v)
                if style_m and ("is_container=1" in style_m.group(1) or "dashed=1" in style_m.group(1) or "fillColor=none" in style_m.group(1)):
                    containers.append(v)
                else:
                    nodes.append(v)
                    
        num_nodes = len(nodes)
        num_edges = len(edges)
        
        # 2. Extract canvas absolute geometry
        all_x = []
        all_y = []
        for m in re.finditer(r'<mxGeometry\s+x="(-?\d+)"\s+y="(-?\d+)"\s+width="(\d+)"\s+height="(\d+)"', xml):
            x = float(m.group(1))
            y = float(m.group(2))
            w = float(m.group(3))
            h = float(m.group(4))
            all_x.append(x)
            all_x.append(x + w)
            all_y.append(y)
            all_y.append(y + h)
            
        canvas_w = max(all_x) - min(all_x) if all_x else 0
        canvas_h = max(all_y) - min(all_y) if all_y else 0
        canvas_area = canvas_w * canvas_h
        
        # Sum of actual node areas
        sum_node_area = 0
        for v in nodes:
            vid_m = re.search(r'id="([^"]+)"', v)
            if vid_m:
                vid = vid_m.group(1)
                geom_p = rf'<mxCell\s+id="{vid}"[^>]+>.*?<mxGeometry\s+x="(-?\d+)"\s+y="(-?\d+)"\s+width="(\d+)"\s+height="(\d+)"'
                geom_m = re.search(geom_p, xml, re.DOTALL)
                if geom_m:
                    w = float(geom_m.group(3))
                    h = float(geom_m.group(4))
                    sum_node_area += (w * h)
                    
        occupancy = (sum_node_area / canvas_area * 100) if canvas_area > 0 else 0
        
        # 3. Read SVG properties
        svg_w, svg_h = "Unknown", "Unknown"
        if svg_path and os.path.exists(svg_path):
            with open(svg_path, 'r', encoding='utf-8', errors='ignore') as f:
                svg_content = f.read()
            svg_w_m = re.search(r'width="([^"]+)"', svg_content)
            svg_h_m = re.search(r'height="([^"]+)"', svg_content)
            if svg_w_m and svg_h_m:
                svg_w, svg_h = svg_w_m.group(1), svg_h_m.group(1)
                
        # 4. Read PNG dimensions and file size
        png_w, png_h = "Unknown", "Unknown"
        png_file_size = 0.0
        if png_path and os.path.exists(png_path):
            png_file_size = os.path.getsize(png_path) / 1024.0
            with Image.open(png_path) as im:
                png_w, png_h = im.width, im.height
                
        print("\n" + "=" * 60)
        print("DIAGRAM DIAGNOSTICS & METRICS")
        print("=" * 60)
        print(f"  Canvas Bounding Box : {canvas_w:.1f} x {canvas_h:.1f} points")
        print(f"  Number of Nodes     : {num_nodes}")
        print(f"  Number of Edges     : {num_edges}")
        print(f"  Node Occupancy %    : {occupancy:.1f}%")
        print(f"  SVG Dimensions      : {svg_w} x {svg_h}")
        print(f"  PNG Pixel Dimensions: {png_w} x {png_h} px")
        print(f"  PNG File Size       : {png_file_size:.2f} KB")
        print("=" * 60 + "\n")
    except Exception as diag_err:
        print(f"[DIAGNOSTICS] Failed to compute diagnostics: {diag_err}")


def generate_drawio_diagram(
    drawio_xml: str, 
    topic: str = "", 
    max_retries: int = 2, 
    disable_fallback: bool = False,
    topology: str = "",
    slide_title: str = "",
    slide_content: str = ""
) -> Tuple[Optional[str], Optional[str]]:
    """
    Render a Draw.io mxGraph XML diagram to SVG and PNG using draw.io CLI.
    Validates rendering quality (size, blankness, visual density, complexity).
    Architecture V4 is the ONLY source of XML - no V3 regeneration.
    Returns (svg_path, png_path) on success.
    Raises detailed exceptions on failure - never returns (None, None) silently.
    """
    import time
    import shutil
    diagram_start = time.time()
    
    print(f"[VISUAL] generate_drawio_diagram called!")
    print(f"   XML length: {len(drawio_xml) if drawio_xml else 0}")
    print(f"   XML preview: {drawio_xml[:200] if drawio_xml else 'EMPTY'}")

    if not drawio_xml:
        raise ValueError("[VISUAL] drawio_xml is empty - cannot render diagram")

    # Remove zero-width characters and other problematic Unicode
    drawio_xml = re.sub(r'[\u200B\u200C\u200D\uFEFF]', '', drawio_xml)  # Zero-width chars
    drawio_xml = re.sub(r'\u00A0', ' ', drawio_xml)  # nbsp to regular space
    drawio_xml = re.sub(r'[^\x09\x0A\x0D\x20-\x7E\x85\xA0-\uD7FF\uE000-\uFFFD]', '', drawio_xml)  # Invalid unicode ranges
    drawio_xml = drawio_xml.replace('\u2014', '-').replace('\u2013', '-').replace('\u201c', '"').replace('\u201d', '"').replace('\u2018', "'").replace('\u2019', "'")

    xml_hash = hashlib.md5(drawio_xml.encode()).hexdigest()[:12]

    # 1. Thread-safe in-memory cache lookup
    with _CACHE_LOCK:
        if xml_hash in DIAGRAM_CACHE:
            svg_cached, png_cached = DIAGRAM_CACHE[xml_hash]
            if os.path.exists(svg_cached) and os.path.exists(png_cached):
                print(f"[VISUAL] Returning in-memory cached diagram for hash '{xml_hash}'")
                print_diagram_diagnostics(drawio_xml, svg_cached, png_cached)
                return svg_cached, png_cached

    cache_dir = os.path.join(tempfile.gettempdir(), "diagram_cache")
    os.makedirs(cache_dir, exist_ok=True)
    cached_svg = os.path.join(cache_dir, f"{xml_hash}.svg")
    cached_png = os.path.join(cache_dir, f"{xml_hash}.png")
    
    if os.path.exists(cached_svg) and os.path.exists(cached_png):
        if os.path.getsize(cached_png) >= 20 * 1024:
            print(f"[VISUAL] Returning persistent disk-cached diagram for hash '{xml_hash}'")
            print_diagram_diagnostics(drawio_xml, cached_svg, cached_png)
            with _CACHE_LOCK:
                DIAGRAM_CACHE[xml_hash] = (cached_svg, cached_png)
            return cached_svg, cached_png

    # 2. Render XML to files (single attempt - no V3 retry loop)
    print(f"[VISUAL] Rendering XML...")
    svg_path, png_path = _render_xml_to_files(drawio_xml)
    
    if not svg_path or not png_path:
        raise RuntimeError(f"[VISUAL] Failed to render diagram - DrawIO CLI returned no output files")
    
    if not os.path.exists(svg_path):
        raise RuntimeError(f"[VISUAL] SVG file not created at {svg_path}")
    
    if not os.path.exists(png_path):
        raise RuntimeError(f"[VISUAL] PNG file not created at {png_path}")
    
    # 3. Validate PNG quality
    is_valid, error_reason = validate_diagram_png(png_path, drawio_xml, topology=topology)
    if not is_valid:
        print(f"[VISUAL] Render validation failed: {error_reason}")
        # Even if validation fails, return the files if they exist (fallback to matplotlib later if needed)
        print(f"[VISUAL] Returning rendered files despite validation failure")
    else:
        print(f"[VISUAL] Render validation passed")
    
    # 4. Cache the results
    try:
        shutil.copy(svg_path, cached_svg)
        shutil.copy(png_path, cached_png)
        print(f"[VISUAL] Cached rendered files to disk: {cached_png}")
        svg_path, png_path = cached_svg, cached_png
    except Exception as cache_write_err:
        print(f"[VISUAL] Warning: Failed to write to disk cache: {cache_write_err}")
    
    # Print diagnostics
    print_diagram_diagnostics(drawio_xml, svg_path, png_path)
    
    with _CACHE_LOCK:
        DIAGRAM_CACHE[xml_hash] = (svg_path, png_path)
    
    elapsed = time.time() - diagram_start
    print(f"[TIMING] Diagram rendered and validated in {elapsed:.2f}s")
    
    return svg_path, png_path


# =====================================================================
# 4. ARCHITECTURE DIAGRAMS  (matplotlib)
# =====================================================================

def _render_arch_matplotlib(xml: str, topic: str = "", slide_title: str = "", topology: str = "") -> str:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches
    from matplotlib.patches import FancyBboxPatch
    from matplotlib.lines import Line2D
    from matplotlib.offsetbox import OffsetImage, AnnotationBbox
    import matplotlib.image as mpimg
    import tempfile, uuid, re, os, base64
    from PIL import Image
    import services.architecture_v4.style_engine_v4 as style_engine_v4
    
    # ── Resolve Active Style ───────────────────────
    active_style = style_engine_v4.get_current_style()
    is_dark = (active_style == "ai_dark_neon")
    
    bg_color = '#0b0b16' if is_dark else '#FFFFFF'
    face_color = '#080810' if is_dark else '#F8FAFD'
    text_color = '#ffffff' if is_dark else '#1A202C'
    title_color = '#00f0ff' if is_dark else '#1A1A2E'
    edge_label_bg = '#181825' if is_dark else '#FFFFFF'
    edge_label_border = '#303046' if is_dark else '#CCCCCC'
    
    # ── Parse nodes & edges ────────────────────────
    nodes = []
    edges = []
    
    # Helper to parse Draw.io style string
    def parse_style_string(style: str) -> dict:
        res = {}
        if not style:
            return res
        parts = style.split(";")
        for p in parts:
            p = p.strip()
            if not p:
                continue
            if "=" in p:
                k, v = p.split("=", 1)
                res[k.strip()] = v.strip()
            else:
                res[p.strip()] = ""
        return res

    import xml.etree.ElementTree as ET
    xml_to_parse = xml.strip()
    if not (xml_to_parse.startswith("<mxfile") or xml_to_parse.startswith("<?xml")):
        xml_to_parse = f"<root_wrapper>{xml_to_parse}</root_wrapper>"
        
    try:
        root = ET.fromstring(xml_to_parse)
        for cell in root.findall(".//mxCell"):
            if cell.get("vertex") == "1":
                nid = cell.get("id")
                if nid in ("0", "1"):
                    continue
                label = cell.get("value") or ""
                
                geom = cell.find("mxGeometry")
                if geom is None:
                    continue
                    
                x = float(geom.get("x") or 0.0)
                y = float(geom.get("y") or 0.0)
                w = float(geom.get("width") or 120.0)
                h = float(geom.get("height") or 60.0)
                
                style_str = cell.get("style") or ""
                style_dict = parse_style_string(style_str)
                
                # Use style_engine_v4.get_palette() instead of hardcoded colors
                category = style_dict.get("category", "compute")
                palette = style_engine_v4.get_palette(category, active_style)
                fill = style_dict.get("fillColor", palette.get('fill', '#dae8fc'))
                stroke = style_dict.get("strokeColor", palette.get('stroke', '#6c8ebf'))
                fc = style_dict.get("fontColor", palette.get('font', '#000000'))
                dashed = style_dict.get("dashed") == "1"
                parent = cell.get("parent")  # Do NOT default to '1' - preserve actual parent
                
                # Metadata from XML style (written by drawio_xml_builder_v4)
                cluster_id = style_dict.get("cluster_id", "")
                namespace = style_dict.get("namespace", "")
                zone = style_dict.get("zone", "")
                lane = style_dict.get("lane", "")
                section = style_dict.get("section", "")
                phase = style_dict.get("phase", "")
                depth = int(style_dict.get("depth", 0))
                level = int(style_dict.get("level", 0))
                shape_hint = style_dict.get("shape_hint", "")
                confidence = float(style_dict.get("confidence", 1.0))
                is_container = style_dict.get("is_container") == "1"
                
                nodes.append({
                    'id': nid,
                    'label': label.replace('&#xa;', '\n')
                                 .replace('&amp;', '&')
                                 .replace('&lt;', '<')
                                 .replace('&gt;', '>'),
                    'x': x, 'y': y,
                    'w': w, 'h': h,
                    'fill': fill, 'stroke': stroke,
                    'fc': fc, 'dashed': dashed,
                    'area': w * h,
                    'parent': parent if parent else None,  # Preserve None if no parent
                    'category': category,
                    'shape_hint': shape_hint,
                    'confidence': confidence,
                    'is_container': is_container,
                    'depth': depth,
                    'level': level,
                    'cluster_id': cluster_id,
                    'namespace': namespace,
                    'zone': zone,
                    'lane': lane,
                    'section': section,
                    'phase': phase
                })
            elif cell.get("edge") == "1":
                src = cell.get("source")
                tgt = cell.get("target")
                label = cell.get("value") or ""
                if src and tgt:
                    edges.append({
                        'src': src, 'tgt': tgt,
                        'label': label.replace('&#xa;', '\n')
                                      .replace('&amp;', '&')
                                      .replace('&lt;', '<')
                                      .replace('&gt;', '>')
                    })
    except Exception as parse_err:
        print(f"[MATPLOTLIB] ElementTree parsing failed: {parse_err}. Falling back to regex.")
        nodes = []
        edges = []
        for m in re.finditer(
            r'<mxCell[^>]+id=["\']([^"\']+)["\']'
            r'[^>]+value=["\']([^"\']*)["\']'
            r'[^>]+vertex=["\']1["\'][^>]*>'
            r'[^<]*<mxGeometry[^>]*'
            r'x=["\']([^"\']+)["\'][^>]*'
            r'y=["\']([^"\']+)["\'][^>]*'
            r'width=["\']([^"\']+)["\'][^>]*'
            r'height=["\']([^"\']+)["\']',
            xml, re.DOTALL
        ):
            nid,label,x,y,w,h = m.groups()
            if nid in ('0','1'): continue
            
            style_str = ''
            style_m = re.search(r'style=["\']([^"\']+)["\']', m.group(0))
            if style_m: style_str = style_m.group(1)
            style_dict = parse_style_string(style_str)
            
            # Use style_engine_v4.get_palette() instead of hardcoded colors
            category = style_dict.get("category", "compute")
            palette = style_engine_v4.get_palette(category, active_style)
            fill = style_dict.get("fillColor", palette.get('fill', '#dae8fc'))
            stroke = style_dict.get("strokeColor", palette.get('stroke', '#6c8ebf'))
            fc = style_dict.get("fontColor", palette.get('font', '#000000'))
            dashed = style_dict.get("dashed") == "1"
            parent = None  # Do NOT default to '1'
            pm = re.search(r'parent=["\']([^"\']+)["\']', m.group(0))
            if pm: parent = pm.group(1)
            
            shape_hint = style_dict.get("shape_hint", "")
            confidence = float(style_dict.get("confidence", 1.0))
            is_container = style_dict.get("is_container") == "1" or float(w)*float(h) > 30000
            depth = int(style_dict.get("depth", 0))
            level = int(style_dict.get("level", 0))
            cluster_id = style_dict.get("cluster_id", "")
            namespace = style_dict.get("namespace", "")
            zone = style_dict.get("zone", "")
            lane = style_dict.get("lane", "")
            section = style_dict.get("section", "")
            phase = style_dict.get("phase", "")
            
            nodes.append({
                'id': nid,
                'label': label.replace('&#xa;','\n')
                             .replace('&amp;','&')
                             .replace('&lt;','<')
                             .replace('&gt;','>'),
                'x': float(x), 'y': float(y),
                'w': float(w), 'h': float(h),
                'fill': fill, 'stroke': stroke,
                'fc': fc, 'dashed': dashed,
                'area': float(w)*float(h),
                'parent': parent,
                'category': category,
                'shape_hint': shape_hint,
                'confidence': confidence,
                'is_container': is_container,
                'depth': depth,
                'level': level,
                'cluster_id': cluster_id,
                'namespace': namespace,
                'zone': zone,
                'lane': lane,
                'section': section,
                'phase': phase
            })
            
        for m in re.finditer(
            r'<mxCell[^>]+edge=["\']1["\'][^>]+'
            r'source=["\']([^"\']+)["\'][^>]+'
            r'target=["\']([^"\']+)["\']'
            r'(?:[^>]+value=["\']([^"\']*)["\'])?',
            xml, re.DOTALL
        ):
            src,tgt = m.group(1),m.group(2)
            label = m.group(3) or ''
            edges.append({
                'src': src, 'tgt': tgt,
                'label': label
            })

    if not nodes:
        return None

    # ── Resolve absolute positions ─────────────────
    def _resolve_absolute_positions(nodes: list) -> list:
        node_map = {n['id']: n for n in nodes}
        def get_abs_pos(node, visited=None):
            if visited is None:
                visited = set()
            if node['id'] in visited:
                return node['x'], node['y']
            visited.add(node['id'])
            parent_id = node.get('parent', '1')
            if parent_id in ('0', '1') or parent_id not in node_map:
                return node['x'], node['y']
            parent = node_map[parent_id]
            px, py = get_abs_pos(parent, visited)
            return node['x'] + px, node['y'] + py
        
        for node in nodes:
            abs_x, abs_y = get_abs_pos(node)
            node['abs_x'] = abs_x
            node['abs_y'] = abs_y
        return nodes

    nodes = _resolve_absolute_positions(nodes)

    # ── Split external and internal nodes ───────────
    external_nodes = [n for n in nodes 
                      if n.get('parent','1') == '1'
                      and not n.get('is_container', False)
                      and n['area'] < 10000]
    
    internal_nodes = [n for n in nodes if n not in external_nodes]

    # Position external nodes in wrapped columns on the left (Problem 6)
    num_ext = len(external_nodes)
    if num_ext > 0:
        max_rows = 4
        col_w = 120
        row_h = 100
        for i, node in enumerate(external_nodes):
            col_idx = i // max_rows
            row_idx = i % max_rows
            node['abs_x'] = -150 - (col_idx * col_w)
            node['abs_y'] = row_idx * row_h + 50

    print(f"[ARCH] Rendering {len(nodes)} nodes, {len(edges)} edges")

    # ── Aspect Ratio & Canvas Size ────────────────
    max_x = max(n['abs_x']+n['w'] for n in nodes)+80
    max_y = max(n['abs_y']+n['h'] for n in nodes)+80
    min_x = min(n['abs_x'] for n in nodes)-50

    # Force figure figsize=(15,8) and dpi=300 (Problem 14)
    fig, ax = plt.subplots(figsize=(15, 8), dpi=300)
    ax.set_xlim(min_x, max_x + 30)
    ax.set_ylim(-max_y-35, 30)
    ax.axis('off')
    fig.patch.set_facecolor(bg_color)
    ax.set_facecolor(bg_color)

    # ── Dynamic Sections and Separators (Topology-aware fallback) ───────────
    topo = style_engine_v4.normalize_topology(topology) if topology else ""
    if not topo:
        xml_lower = xml.lower()
        if "transformer" in xml_lower or "attention" in xml_lower:
            topo = "transformer_pipeline"
        elif "cnn" in xml_lower or "conv2d" in xml_lower:
            topo = "cnn_pipeline"
        elif "rag" in xml_lower or "retriever" in xml_lower or "pinecone" in xml_lower:
            topo = "rag_pipeline"
        elif "kubernetes" in xml_lower or "k8s" in xml_lower:
            topo = "kubernetes"
        elif "event" in xml_lower or "kafka" in xml_lower or "stream" in xml_lower:
            topo = "event-driven"
        elif "cloud" in xml_lower:
            topo = "cloud"
        elif "layered" in xml_lower:
            topo = "layered"
            
    internal_xs = [n['abs_x'] for n in internal_nodes if n['abs_x'] >= 0]
    min_internal_x = min(internal_xs) if internal_xs else 0
    max_internal_x = max(n['abs_x'] + n['w'] for n in internal_nodes if n['abs_x'] >= 0) if internal_xs else max_x
    span = max_internal_x - min_internal_x
    
    num_sections = 3 if span > 300 else 2
    # Topology-aware section labels
    if "transformer" in topo:
        labels = ["Input & Embeddings", "Encoder / Attention Block", "Decoder / Projection"]
    elif "cnn" in topo:
        labels = ["Input Layer", "Convolution & Pooling", "Dense / Output"]
    elif "rag" in topo:
        labels = ["Ingestion & Retrieval", "Embedding & Vector DB", "Generation & LLM Response"]
    elif "event" in topo:
        labels = ["Event Producers", "Message Broker / Stream", "Event Consumers / DB"]
    elif "k8s" in topo or "kubernetes" in topo:
        labels = ["Ingress Routing", "Compute (Pods / Services)", "State / Persistent Storage"]
    elif "cloud" in topo:
        labels = ["Edge / CDN", "Application Layer", "Data & Storage Layer"]
    elif "layered" in topo:
        labels = ["Presentation Layer", "Business Logic Layer", "Data Access Layer"]
    elif "react" in topo or "frontend" in xml.lower():
        labels = ["User Client (UI)", "Frontend Components", "Backend Services & Database"]
    else:
        labels = ["Public Access / Gateway", "Application Infrastructure", "Data & Cache Storage"]
        
    boundaries = []
    if num_sections == 2:
        labels = [labels[0], labels[2]]
        boundaries = [min_internal_x + span / 2]
    else:
        boundaries = [min_internal_x + span / 3, min_internal_x + 2 * span / 3]
        
    for bx in boundaries:
        ax.axvline(x=bx, color='#303046' if is_dark else '#E0E0E0', 
                   linewidth=1.5, linestyle='--', alpha=0.5, zorder=0)
                   
    for i in range(num_sections):
        sec_min = min_internal_x if i == 0 else boundaries[i-1]
        sec_max = boundaries[i] if i < len(boundaries) else max_internal_x
        sec_center = (sec_min + sec_max) / 2
        ax.text(sec_center, 15, labels[i],
                fontsize=11, color=title_color, 
                ha='center', fontweight='bold', alpha=0.9)

    if external_nodes:
        ext_max_x = max(n['abs_x'] + n['w'] for n in external_nodes)
        sep_x = (ext_max_x + min_internal_x) / 2
        ax.axvline(x=sep_x, color='#303046' if is_dark else '#E0E0E0', 
                   linewidth=2, linestyle='--', alpha=0.7, zorder=0)
        ax.text((ext_max_x + min_x)/2, 15, 'External Clients', 
                fontsize=11, color=title_color,
                ha='center', fontweight='bold', alpha=0.9)

    node_map = {n['id']: n for n in nodes}
    present_categories = set()
    temp_logo_paths = []
    
    # Separate nodes by type for z-order preservation
    container_nodes = [n for n in nodes if n.get('is_container', False)]
    regular_nodes = [n for n in nodes if not n.get('is_container', False)]
    
    def get_edge_color(style):
        colors = {
            "classic": "#4a5568",
            "drawio_vivid": "#4361ee",
            "ai_dark_neon": "#ff2a5f",
            "aws": "#ff9900",
            "azure": "#0078d4",
            "gcp": "#4285f4",
            "kubernetes": "#326ce5",
            "minimal": "#555555"
        }
        return colors.get(style, "#4a5568")

    # ── Draw nodes with z-order preservation ────────────
    # First: containers (zorder=0)
    for node in sorted(container_nodes, 
                       key=lambda n: n['area'],
                       reverse=True):
        x = node['abs_x']
        y = -(node['abs_y'] + node['h'])
        w, h = node['w'], node['h']
        is_container = node.get('is_container', False)
        category = node.get('category', 'compute')
        
        if not is_container:
            present_categories.add(category)
            
        palette = style_engine_v4.get_palette(category, active_style)
        
        if is_container:
            fill = '#0a0a12' if is_dark else '#FBFCFC'
            stroke_c = '#bd00ff' if is_dark else '#4a5568'
            ls = '--'
            lw = 2.0
            
            depth = node.get('depth', 0)
            if depth == 0:
                alpha = 0.08 if is_dark else 0.12
            elif depth == 1:
                alpha = 0.12 if is_dark else 0.18
            else:
                alpha = 0.18 if is_dark else 0.24
            zorder = 0  # Containers at bottom
        else:
            # Use style_engine_v4.get_palette() for consistent colors
            node_palette = style_engine_v4.get_palette(category, active_style)
            fill = node_palette.get('fill', '#dae8fc')
            stroke_c = node_palette.get('stroke', '#6c8ebf')
            ls = '-'
            lw = 1.5
            alpha = 0.95
            zorder = 5  # Regular nodes above containers

        if not is_container and active_style != "minimal":
            shadow = FancyBboxPatch(
                (x+3, y-3), w, h,
                boxstyle="round,pad=4",
                facecolor='#000000',
                edgecolor='none',
                alpha=0.15,
                zorder=zorder-1
            )
            ax.add_patch(shadow)

        rect = FancyBboxPatch(
            (x, y), w, h,
            boxstyle="round,pad=4",
            facecolor=fill,
            edgecolor=stroke_c,
            linewidth=lw,
            linestyle=ls,
            alpha=alpha,
            zorder=zorder
        )
        ax.add_patch(rect)

        label = node['label']
        lines = label.split('\n')
        lines = [l if len(l) <= 30 else l[:28] + '..' for l in lines[:4]]
        display = '\n'.join(lines)
        
        fc = node_palette.get('font', '#ffffff' if is_dark else '#000000')
        fw = 'bold'
        
        if is_container:
            fs = max(11, min(14, w / 45))
            tx, ty = x + 8, y + h - 12
            ha, va = 'left', 'top'
        else:
            fs = max(10, min(13, min(w, h) / 9))
            tx, ty = x + w / 2, y + h / 2
            ha, va = 'center', 'center'

        has_logo = False
        if not is_container and active_style != "minimal":
            try:
                from services.architecture_v4.ai_icon_engine import get_ai_icon
                res = get_ai_icon(node['label'], kind=node.get('category'), style=active_style)
                if res and res.get('style'):
                    import base64
                    m = re.search(r'image=data:image/svg\+xml;base64,([A-Za-z0-9+/=]+)', res['style'])
                    if m:
                        svg_bytes = base64.b64decode(m.group(1))
                        with tempfile.NamedTemporaryFile(suffix=".svg", delete=False) as temp_svg:
                            temp_svg.write(svg_bytes)
                            temp_svg_path = temp_svg.name
                        
                        temp_png_path = temp_svg_path.replace(".svg", ".png")
                        import cairosvg
                        cairosvg.svg2png(
                            url=temp_svg_path,
                            write_to=temp_png_path,
                            scale=2.0
                        )
                        
                        if os.path.exists(temp_png_path) and os.path.getsize(temp_png_path) > 0:
                            with Image.open(temp_png_path) as logo_img_pil:
                                logo_h = 32
                                logo_w = int(logo_img_pil.width * (logo_h / logo_img_pil.height))
                                img_resized = logo_img_pil.resize((logo_w, logo_h), Image.Resampling.LANCZOS)
                                resized_logo_path = temp_png_path.replace(".png", "_resized.png")
                                img_resized.save(resized_logo_path)
                                
                                logo_img = mpimg.imread(resized_logo_path)
                                imagebox = OffsetImage(logo_img, zoom=0.7)
                                logo_cx = x + w / 2
                                logo_cy = y + h / 2 + 10
                                text_cx = x + w / 2
                                text_cy = y + h / 2 - 14
                                
                                ab = AnnotationBbox(imagebox, (logo_cx, logo_cy), frameon=False, zorder=zorder+1)
                                ax.add_artist(ab)
                                has_logo = True
                                
                                temp_logo_paths.extend([temp_svg_path, temp_png_path, resized_logo_path])
            except Exception as logo_err:
                print(f"[MATPLOTLIB] Logo rendering skipped/failed: {logo_err}")

        if has_logo:
            ax.text(text_cx, text_cy, display,
                    ha='center', va='center',
                    fontsize=fs, fontweight=fw,
                    color=fc,
                    zorder=10,
                    multialignment='center',
                    clip_on=True)
        else:
            ax.text(tx, ty, display,
                    ha=ha, va=va,
                    fontsize=fs, fontweight=fw,
                    color=fc,
                    zorder=10,
                    multialignment='center',
                    clip_on=True)

    # Second: regular nodes (zorder=5) - render after containers
    for node in sorted(regular_nodes, 
                       key=lambda n: n['area'],
                       reverse=True):
        x = node['abs_x']
        y = -(node['abs_y'] + node['h'])
        w, h = node['w'], node['h']
        is_container = node.get('is_container', False)
        category = node.get('category', 'compute')
        
        present_categories.add(category)
        
        # Use style_engine_v4.get_palette() for consistent colors
        node_palette = style_engine_v4.get_palette(category, active_style)
        fill = node_palette.get('fill', '#dae8fc')
        stroke_c = node_palette.get('stroke', '#6c8ebf')
        ls = '-'
        lw = 1.5
        alpha = 0.95
        zorder = 5  # Regular nodes above containers

        if active_style != "minimal":
            shadow = FancyBboxPatch(
                (x+3, y-3), w, h,
                boxstyle="round,pad=4",
                facecolor='#000000',
                edgecolor='none',
                alpha=0.15,
                zorder=zorder-1
            )
            ax.add_patch(shadow)

        rect = FancyBboxPatch(
            (x, y), w, h,
            boxstyle="round,pad=4",
            facecolor=fill,
            edgecolor=stroke_c,
            linewidth=lw,
            linestyle=ls,
            alpha=alpha,
            zorder=zorder
        )
        ax.add_patch(rect)

        label = node['label']
        lines = label.split('\n')
        lines = [l if len(l) <= 30 else l[:28] + '..' for l in lines[:4]]
        display = '\n'.join(lines)
        
        fc = node_palette.get('font', '#ffffff' if is_dark else '#000000')
        fw = 'bold'
        
        fs = max(10, min(13, min(w, h) / 9))
        tx, ty = x + w / 2, y + h / 2
        ha, va = 'center', 'center'

        has_logo = False
        if active_style != "minimal":
            try:
                from services.architecture_v4.ai_icon_engine import get_ai_icon
                res = get_ai_icon(node['label'], kind=node.get('category'), style=active_style)
                if res and res.get('style'):
                    import base64
                    m = re.search(r'image=data:image/svg\+xml;base64,([A-Za-z0-9+/=]+)', res['style'])
                    if m:
                        svg_bytes = base64.b64decode(m.group(1))
                        with tempfile.NamedTemporaryFile(suffix=".svg", delete=False) as temp_svg:
                            temp_svg.write(svg_bytes)
                            temp_svg_path = temp_svg.name
                        
                        temp_png_path = temp_svg_path.replace(".svg", ".png")
                        import cairosvg
                        cairosvg.svg2png(
                            url=temp_svg_path,
                            write_to=temp_png_path,
                            scale=2.0
                        )
                        
                        if os.path.exists(temp_png_path) and os.path.getsize(temp_png_path) > 0:
                            with Image.open(temp_png_path) as logo_img_pil:
                                logo_h = 32
                                logo_w = int(logo_img_pil.width * (logo_h / logo_img_pil.height))
                                img_resized = logo_img_pil.resize((logo_w, logo_h), Image.Resampling.LANCZOS)
                                resized_logo_path = temp_png_path.replace(".png", "_resized.png")
                                img_resized.save(resized_logo_path)
                                
                                logo_img = mpimg.imread(resized_logo_path)
                                imagebox = OffsetImage(logo_img, zoom=0.7)
                                logo_cx = x + w / 2
                                logo_cy = y + h / 2 + 10
                                text_cx = x + w / 2
                                text_cy = y + h / 2 - 14
                                
                                ab = AnnotationBbox(imagebox, (logo_cx, logo_cy), frameon=False, zorder=zorder+1)
                                ax.add_artist(ab)
                                has_logo = True
                                
                                temp_logo_paths.extend([temp_svg_path, temp_png_path, resized_logo_path])
            except Exception as logo_err:
                print(f"[MATPLOTLIB] Logo rendering skipped/failed: {logo_err}")

        if has_logo:
            ax.text(text_cx, text_cy, display,
                    ha='center', va='center',
                    fontsize=fs, fontweight=fw,
                    color=fc,
                    zorder=10,
                    multialignment='center',
                    clip_on=True)
        else:
            ax.text(tx, ty, display,
                    ha=ha, va=va,
                    fontsize=fs, fontweight=fw,
                    color=fc,
                    zorder=10,
                    multialignment='center',
                    clip_on=True)

    # Third: edges (zorder=3-4) - draw after containers, before nodes
    # ── Draw edges (orthogonal routing) ────────────
    drawn = set()
    canvas_span = max(100.0, max_x - min_x)
    edge_offset_amount = max(8.0, min(25.0, canvas_span * 0.015))
    edge_label_fs = max(8, min(11, int(canvas_span / 150) + 7))

    for edge in edges:
        src = node_map.get(edge['src'])
        tgt = node_map.get(edge['tgt'])
        if not src or not tgt:
            continue
        
        key = tuple(sorted([edge['src'], edge['tgt']]))
        offset_val = 0
        if key in drawn:
            offset_val = edge_offset_amount
        drawn.add(key)
        
        src_cx = src['abs_x'] + src['w'] / 2
        src_cy = src['abs_y'] + src['h'] / 2
        tgt_cx = tgt['abs_x'] + tgt['w'] / 2
        tgt_cy = tgt['abs_y'] + tgt['h'] / 2
        
        dx = tgt_cx - src_cx
        dy = tgt_cy - src_cy
        
        if abs(dx) > abs(dy):
            if dx > 0:
                sx = src['abs_x'] + src['w']
                sy = src_cy
                tx = tgt['abs_x']
                ty = tgt_cy
            else:
                sx = src['abs_x']
                sy = src_cy
                tx = tgt['abs_x'] + tgt['w']
                ty = tgt_cy
        else:
            if dy > 0:
                sx = src_cx
                sy = src['abs_y'] + src['h']
                tx = tgt_cx
                ty = tgt['abs_y']
            else:
                sx = src_cx
                sy = src['abs_y']
                tx = tgt_cx
                ty = tgt['abs_y'] + tgt['h']
        
        sy += offset_val
        ty += offset_val
        sy_plot = -sy
        ty_plot = -ty
        
        ecol = get_edge_color(active_style)
        if src.get('stroke'):
            ecol = src['stroke']

        if abs(dx) > abs(dy):
            mid_x = (sx + tx) / 2 + offset_val
            path_x = [sx, mid_x, mid_x, tx]
            path_y = [sy_plot, sy_plot, ty_plot, ty_plot]
        else:
            mid_y = (sy_plot + ty_plot) / 2 + offset_val
            path_x = [sx, sx, tx, tx]
            path_y = [sy_plot, mid_y, mid_y, ty_plot]
        
        ax.plot(path_x[:-1], path_y[:-1], color=ecol, linewidth=2,
                solid_capstyle='round', solid_joinstyle='round',
                zorder=3)  # Edge lines below nodes
        
        ax.annotate('',
            xy=(path_x[-1], path_y[-1]),
            xytext=(path_x[-2], path_y[-2]),
            arrowprops=dict(
                arrowstyle='-|>',
                color=ecol, lw=2,
                mutation_scale=18
            ),
            zorder=4)  # Arrowheads above edge lines
        
        if edge['label']:
            mid_idx = len(path_x) // 2
            mx = (path_x[mid_idx - 1] + path_x[mid_idx]) / 2
            my = (path_y[mid_idx - 1] + path_y[mid_idx]) / 2
            ax.text(mx, my, edge['label'],
                fontsize=edge_label_fs, color=text_color,
                ha='center', va='center',
                fontweight='bold',
                bbox=dict(boxstyle='round,pad=0.2',
                         facecolor=edge_label_bg,
                         edgecolor=edge_label_border,
                         alpha=0.9),
                zorder=8)  # Edge labels above arrows
    
    # ── Dynamic Legend ────────────────────────────
    category_labels = {
        "compute": "Compute / Services",
        "service": "Microservices",
        "database": "Database / Storage",
        "vector_db": "Vector Database",
        "cache": "In-Memory Cache",
        "queue": "Message Queue / Stream",
        "llm": "LLM / AI Model",
        "gateway": "API Gateway / Proxy",
        "security": "Security / IAM",
        "monitoring": "Monitoring & Metrics",
        "logging": "Logging & Auditing",
        "observability": "Observability",
        "client": "Client Application",
        "external": "External Integration",
        "storage": "Object Storage",
        "analytics": "Data Analytics",
        "framework": "Framework / Engine",
        "embedding": "Embedding Model",
        "retriever": "Information Retriever",
        "agent": "AI Agent / Orchestrator",
        "encoder": "NN Encoder",
        "decoder": "NN Decoder",
        "conv": "Convolutional Layer",
        "pool": "Pooling Layer",
        "dense": "Dense / Fully Connected"
    }

    legend_items = []
    for cat in sorted(list(present_categories)):
        label = category_labels.get(cat, cat.replace('_', ' ').title())
        pal = style_engine_v4.get_palette(cat, active_style)
        legend_items.append(mpatches.Patch(
            facecolor=pal.get('fill', '#dae8fc'),
            edgecolor=pal.get('stroke', '#6c8ebf'),
            label=label
        ))
        
    if legend_items:
        ncol = min(4, len(legend_items))
        ax.legend(
            handles=legend_items,
            loc='lower left',
            fontsize=9,
            ncol=ncol,
            framealpha=0.9,
            facecolor=edge_label_bg,
            edgecolor=edge_label_border,
            labelcolor='white' if is_dark else 'black',
            bbox_to_anchor=(0, -0.05)
        )
    
    # ── Dynamic Title ─────────────────────────────
    title_text = "System Architecture Diagram"
    if slide_title:
        title_text = slide_title
    elif topic:
        title_text = f"{topic} Architecture"
        
    ax.set_title(
        title_text,
        fontsize=18, fontweight='bold',
        color=title_color, pad=15,
        loc='left'
    )
    
    plt.tight_layout(pad=0.5)
    
    out = os.path.join(
        tempfile.gettempdir(),
        f"arch_{uuid.uuid4().hex[:8]}.png"
    )
    plt.savefig(out, dpi=300,
                bbox_inches='tight',
                facecolor=bg_color)
    plt.close()
    
    for p in temp_logo_paths:
        try:
            if os.path.exists(p):
                os.remove(p)
        except OSError:
            pass
            
    if os.path.exists(out):
        sz = os.path.getsize(out)
        print(f"[ARCH] Matplotlib fallback success: {out} ({sz} bytes)")
        return out
    return None


# =====================================================================
# 5. CHARTS  (matplotlib)
# =====================================================================

def render_chart(chart_data: dict) -> Optional[str]:
    """
    Render a chart (bar, pie, or line) using matplotlib.

    chart_data format:
    {
        "type": "bar" | "pie" | "line",
        "labels": ["Q1", "Q2", "Q3", "Q4"],
        "values": [120, 150, 180, 210],
        "title": "Revenue Growth"
    }

    Returns the PNG file path on success, None on failure.
    """
    try:
        import matplotlib
        matplotlib.use("Agg")  # Non-interactive backend
        import matplotlib.pyplot as plt

        # Fix Windows encoding issues with Unicode characters in matplotlib
        matplotlib.rcParams['axes.unicode_minus'] = False
        matplotlib.rcParams['font.family'] = 'DejaVu Sans'

        chart_type = chart_data.get("type", "bar").lower()
        labels = chart_data.get("labels", [])
        values = chart_data.get("values", [])
        title = chart_data.get("title", "")

        if not labels or not values:
            print("[VISUAL] Chart: missing labels or values")
            return None

        print(f"[VISUAL] Rendering {chart_type} chart: '{title}'")

        fig, ax = plt.subplots(figsize=(8, 4.5))

        accent_color = "#534AB7"

        if chart_type == "bar":
            ax.bar(labels, values, color=accent_color, width=0.6, edgecolor="white", linewidth=0.5)
            ax.set_ylabel("")
        elif chart_type == "pie":
            colors = [accent_color, "#7C6DD8", "#A594F9", "#C8BBFD", "#E8E3FF",
                      "#6B5EAA", "#8F84C7", "#B3AAE0"]
            ax.pie(values, labels=labels, colors=colors[:len(values)],
                   autopct="%1.1f%%", startangle=140,
                   textprops={"fontsize": 11, "color": "#333333"})
            ax.set_aspect("equal")
        elif chart_type == "line":
            ax.plot(labels, values, color=accent_color, marker="o", linewidth=2.5,
                    markersize=8, markerfacecolor="white", markeredgecolor=accent_color,
                    markeredgewidth=2)
            ax.fill_between(range(len(labels)), values, alpha=0.08, color=accent_color)
            ax.set_ylabel("")
        else:
            # Default to bar
            ax.bar(labels, values, color=accent_color, width=0.6)

        if title:
            ax.set_title(title, fontsize=14, fontweight="bold", pad=12, color="#333333")

        # Style: clean white, no top/right spines
        if chart_type != "pie":
            ax.spines["top"].set_visible(False)
            ax.spines["right"].set_visible(False)
            ax.spines["left"].set_color("#cccccc")
            ax.spines["bottom"].set_color("#cccccc")
            ax.tick_params(colors="#666666", labelsize=10)
            ax.set_facecolor("white")

        fig.patch.set_facecolor("white")
        fig.tight_layout(pad=1.5)

        path = os.path.join(TMP_DIR, f"chart_{uuid.uuid4().hex[:10]}.png")
        fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
        plt.close(fig)

        print(f"[VISUAL] Chart rendered -> {path}")
        return path

    except Exception as e:
        print(f"[VISUAL] render_chart error: {e}")
        return None


# =====================================================================
# 5. ILLUSTRATIONS  (local SVGs in assets/illustrations/)
# =====================================================================

# Base path for illustration SVGs (relative to project root)
_ILLUSTRATIONS_DIR = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "assets", "illustrations"
)


def get_illustration(keyword: str) -> Optional[str]:
    """
    Find a local SVG illustration matching the keyword from assets/illustrations/.
    Uses fuzzy matching with difflib (cutoff 0.3).
    Returns the full file path on match, None otherwise.
    """
    try:
        print(f"[VISUAL] Searching illustration for: '{keyword}'")

        if not os.path.isdir(_ILLUSTRATIONS_DIR):
            print(f"[VISUAL] Illustrations directory not found: {_ILLUSTRATIONS_DIR}")
            return None

        svg_files = [f for f in os.listdir(_ILLUSTRATIONS_DIR) if f.lower().endswith(".svg")]
        if not svg_files:
            print("[VISUAL] No SVG files found in illustrations directory")
            return None

        # Normalize keyword and filenames for matching
        normalized_kw = keyword.lower().replace(" ", "_").replace("-", "_")
        name_map = {}
        for f in svg_files:
            stem = os.path.splitext(f)[0].lower().replace("-", "_")
            name_map[stem] = f

        # Exact match first
        if normalized_kw in name_map:
            path = os.path.join(_ILLUSTRATIONS_DIR, name_map[normalized_kw])
            print(f"[VISUAL] Illustration exact match -> {path}")
            return path

        # Fuzzy match
        matches = difflib.get_close_matches(normalized_kw, name_map.keys(), n=1, cutoff=0.3)
        if matches:
            path = os.path.join(_ILLUSTRATIONS_DIR, name_map[matches[0]])
            print(f"[VISUAL] Illustration fuzzy match '{matches[0]}' -> {path}")
            return path

        print(f"[VISUAL] No illustration match for '{keyword}'")
        return None

    except Exception as e:
        print(f"[VISUAL] get_illustration error: {e}")
        return None


# =====================================================================
# 6. TABLES  (python-pptx — inline table builder)
# =====================================================================

def build_table(slide, table_data: dict, left, top, width, height):
    """
    Add a styled table to a python-pptx slide.

    table_data format:
    {
        "headers": ["Feature", "Basic", "Pro", "Enterprise"],
        "rows": [
            ["Storage", "5GB", "50GB", "1TB"],
            ["Users", "1", "10", "Unlimited"]
        ]
    }

    Returns the table shape object, or None on failure.
    """
    try:
        from pptx.util import Pt, Inches
        from pptx.dml.color import RGBColor

        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])

        if not headers:
            print("[VISUAL] build_table: no headers provided")
            return None

        num_rows = len(rows) + 1  # +1 for header row
        num_cols = len(headers)

        print(f"[VISUAL] Building table: {num_rows} rows × {num_cols} cols")

        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = table_shape.table

        # ── Header row ──
        header_bg = RGBColor(0x53, 0x4A, 0xB7)  # #534AB7
        header_text_color = RGBColor(0xFF, 0xFF, 0xFF)

        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header_text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_bg

            para = cell.text_frame.paragraphs[0]
            para.word_wrap = True
            if para.runs:
                run = para.runs[0]
            else:
                run = para.add_run()
                run.text = str(header_text)
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = header_text_color
            run.font.name = "Segoe UI"

        # ── Data rows ──
        alt_bg_white = RGBColor(0xFF, 0xFF, 0xFF)
        alt_bg_tinted = RGBColor(0xF5, 0xF5, 0xFF)  # #F5F5FF
        body_text_color = RGBColor(0x33, 0x33, 0x33)

        for row_idx, row_data in enumerate(rows):
            bg_color = alt_bg_white if row_idx % 2 == 0 else alt_bg_tinted

            for col_idx in range(num_cols):
                cell = table.cell(row_idx + 1, col_idx)
                cell_text = str(row_data[col_idx]) if col_idx < len(row_data) else ""
                cell.text = cell_text
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_color

                para = cell.text_frame.paragraphs[0]
                para.word_wrap = True
                if para.runs:
                    run = para.runs[0]
                else:
                    run = para.add_run()
                    run.text = cell_text
                run.font.size = Pt(12)
                run.font.bold = False
                run.font.color.rgb = body_text_color
                run.font.name = "Segoe UI"

        print("[VISUAL] Table built successfully")
        return table_shape

    except Exception as e:
        print(f"[VISUAL] build_table error: {e}")
        return None
