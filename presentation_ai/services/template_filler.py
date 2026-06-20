"""
Template Filler — fills AI-generated content into an existing PPTX template
without modifying any design elements (fonts, colors, positions, sizes).

The template schema acts as the absolute source of truth.
"""

import io
import os
import re
import tempfile
import uuid
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def fill_template(
    template_path: str,
    template_schema: dict,
    slide_contents: list,
    output_dir: str = "outputs",
) -> str:
    """Fill AI-generated content into an original PPTX template file.

    Args:
        template_path:   Path to the original .pptx template file.
        template_schema: Schema dict describing the template structure.
        slide_contents:  List of dicts, one per slide, containing the content.
        output_dir:      Directory where the filled presentation is saved.

    Returns:
        Absolute path to the saved .pptx file.
    """
    prs = Presentation(template_path)
    os.makedirs(output_dir, exist_ok=True)
    output_filename = f"{output_dir}/template_{uuid.uuid4().hex[:10]}.pptx"

    # Get slide dimensions for percentage calculations
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 1. Capture initial shapes data for validation
    initial_shapes_data = {}
    for idx, slide in enumerate(prs.slides):
        shapes_info = []
        for s in slide.shapes:
            font_size = None
            if s.has_text_frame and s.text_frame.paragraphs:
                p = s.text_frame.paragraphs[0]
                for r in p.runs:
                    if r.font.size:
                        font_size = r.font.size.pt
                        break
            shapes_info.append({
                "shape_id": s.shape_id,
                "left": s.left,
                "top": s.top,
                "width": s.width,
                "height": s.height,
                "rotation": getattr(s, "rotation", 0),
                "font_size": font_size
            })
        initial_shapes_data[idx] = shapes_info

    # 2. Fill each slide with matching content
    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx >= len(slide_contents):
            break
        content = slide_contents[slide_idx]
        schema_slide = (
            template_schema["slides"][slide_idx]
            if slide_idx < len(template_schema.get("slides", []))
            else None
        )
        _fill_slide(slide, content, schema_slide, slide_idx, slide_width, slide_height)

    # Remove extra slides (template has more slides than content)
    total_template_slides = len(prs.slides)
    if total_template_slides > len(slide_contents):
        for idx in range(total_template_slides - 1, len(slide_contents) - 1, -1):
            _delete_slide(prs, idx)

    # 3. Validate shape count, coordinates, z-order, font size, and rotation before saving (print warnings, do not raise exceptions)
    for slide_idx in range(len(slide_contents)):
        if slide_idx >= len(prs.slides):
            break
        slide = prs.slides[slide_idx]
        expected_shapes = initial_shapes_data.get(slide_idx, [])
        actual_shapes = list(slide.shapes)

        if len(expected_shapes) != len(actual_shapes):
            print(f"[TEMPLATE WARNING] Shape count changed on slide {slide_idx + 1}. Expected {len(expected_shapes)}, got {len(actual_shapes)}")
            continue

        for idx, exp in enumerate(expected_shapes):
            act = actual_shapes[idx]
            if exp["shape_id"] != act.shape_id:
                print(f"[TEMPLATE WARNING] Z-order or shape identity changed on slide {slide_idx + 1} at position {idx}. Expected shape_id {exp['shape_id']}, got {act.shape_id}")
            if (exp["left"] != act.left or exp["top"] != act.top or 
                exp["width"] != act.width or exp["height"] != act.height):
                print(f"[TEMPLATE WARNING] Shape geometry/coordinates changed on slide {slide_idx + 1} for shape_id {exp['shape_id']}. Expected (L:{exp['left']}, T:{exp['top']}, W:{exp['width']}, H:{exp['height']}), got (L:{act.left}, T:{act.top}, W:{act.width}, H:{act.height})")
            
            exp_rot = exp.get("rotation", 0)
            act_rot = getattr(act, "rotation", 0)
            if exp_rot != act_rot:
                print(f"[TEMPLATE WARNING] Rotation changed on slide {slide_idx + 1} for shape_id {exp['shape_id']}. Expected {exp_rot}, got {act_rot}")
                
            exp_font_size = exp.get("font_size")
            if exp_font_size and act.has_text_frame and act.text_frame.paragraphs:
                p = act.text_frame.paragraphs[0]
                act_font_size = None
                for r in p.runs:
                    if r.font.size:
                        act_font_size = r.font.size.pt
                        break
                if act_font_size and abs(exp_font_size - act_font_size) > 0.1:
                    print(f"[TEMPLATE WARNING] Font size changed on slide {slide_idx + 1} for shape_id {exp['shape_id']}. Expected {exp_font_size}, got {act_font_size}")

    # Apply complex script formatting fix
    try:
        from services.template_filler_v2 import fix_presentation_complex_scripts
        fix_presentation_complex_scripts(prs)
    except Exception as e:
        print(f"[TEMPLATE_FILLER] Warning: failed to apply complex script fix: {e}")

    prs.save(output_filename)
    return os.path.abspath(output_filename)


# ---------------------------------------------------------------------------
# Slide-level filling
# ---------------------------------------------------------------------------

def _fill_slide(slide, content, schema_slide, slide_idx, slide_width, slide_height):
    """Fill a single slide's shapes with content using schema_slide zones."""
    if not schema_slide or "zones" not in schema_slide:
        print(f"[TEMPLATE_FILLER] Warning: No schema zones found for slide {slide_idx + 1}")
        return

    layout = content.get("layout", "content") if content else "content"
    filled_shape_ids = set()

    for zone in schema_slide.get("zones", []):
        role = zone.get("role", "other")
        
        # Skip icon and decorative/slide-wide elements
        role_lower = role.lower()
        is_icon_role = any(k in role_lower for k in ("icon", "logo", "vector", "symbol", "circle", "ellipse"))
        
        dec_keywords = ("footer", "caption", "label", "decorative", "chart", "table", "other", "shape", "arrow", "note", "source", "reference", "metric", "box")
        is_dec_role = any(k in role_lower for k in dec_keywords)
        
        # Skip general 'card' roles but allow cardX_title and cardX_description
        if "card" in role_lower and not ("_title" in role_lower or "_description" in role_lower):
            is_dec_role = True
            
        if is_icon_role or is_dec_role:
            continue

        # Find matching shape using multi-criteria matching
        shape = _find_shape(slide, zone, slide_width, slide_height)
        if shape is None:
            shape = _find_fallback_shape(slide, zone, filled_shape_ids)
            if shape is None:
                print("[TEMPLATE WARNING]")
                print("shape not found")
                print(f"role {role}")
                continue

        filled_shape_ids.add(shape.shape_id)

        val = _get_role_value_with_fallback(role, content, layout)
        if val is None:
            continue
        if isinstance(val, str) and not val.strip():
            continue

        if "image" in role_lower:
            # Replace image bytes via relationship override
            _replace_image_in_shape(slide, shape, val)
            print("[TEMPLATE]")
            print(f"slide {slide_idx + 1}")
            print(f"shape {shape.shape_id}")
            print(f"role {role}")
            print("filled")
        elif shape.has_text_frame:
            # Determine font size for capacity estimation
            font_size = None
            if shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                for r in p.runs:
                    if r.font.size:
                        font_size = r.font.size.pt
                        break
            
            # Estimate capacity
            max_words = _estimate_capacity(shape, font_size)
            
            # Truncate gracefully if text exceeds capacity
            val_str = str(val).strip()
            words_list = val_str.split()
            word_count = len(words_list)
            if word_count > max_words:
                val_str = " ".join(words_list[:max_words]) + "..."
                word_count = len(val_str.split())
                val = val_str
                
            # Fill text frame in-place preserving formatting
            _fill_text_in_shape(shape, val)

            # Print debug log in the exact requested format
            print("[TEMPLATE]")
            print(f"slide {slide_idx + 1}")
            print(f"shape {shape.shape_id}")
            print(f"role {role}")
            print(f"words {word_count}")
            print(f"capacity {max_words}")
            print(f"paragraphs {len(shape.text_frame.paragraphs)}")


# ---------------------------------------------------------------------------
# Shape-level matching
# ---------------------------------------------------------------------------

def _find_shape(slide, zone, slide_width, slide_height):
    """Find the shape in slide that best matches the schema zone definition."""
    shape_id = zone.get("shape_id")
    ph_idx = zone.get("placeholder_idx")
    name = zone.get("name")

    # 1. Match by shape_id (Primary Mechanism)
    if shape_id is not None:
        for s in slide.shapes:
            if s.shape_id == shape_id:
                return s

    # 2. Match by placeholder_idx
    if ph_idx is not None:
        for s in slide.shapes:
            if s.is_placeholder:
                try:
                    if s.placeholder_format.idx == ph_idx:
                        return s
                except Exception:
                    pass

    # 3. Match by shape name
    if name:
        for s in slide.shapes:
            if s.name == name:
                return s

    # 4. Match by coordinate similarity (threshold distance < 5%)
    target_left = zone.get("left_pct", 0)
    target_top = zone.get("top_pct", 0)
    target_width = zone.get("width_pct", 0)
    target_height = zone.get("height_pct", 0)

    best_match = None
    best_dist = float("inf")

    for s in slide.shapes:
        s_left = (s.left / slide_width) * 100 if slide_width else 0
        s_top = (s.top / slide_height) * 100 if slide_height else 0
        s_width = (s.width / slide_width) * 100 if slide_width else 0
        s_height = (s.height / slide_height) * 100 if slide_height else 0

        dist = (abs(s_left - target_left) + 
                abs(s_top - target_top) + 
                abs(s_width - target_width) + 
                abs(s_height - target_height))
        
        if dist < best_dist:
            best_dist = dist
            best_match = s

    if best_dist < 5.0:
        return best_match

    return None


def _find_fallback_shape(slide, zone, filled_shape_ids):
    """Find the first unused shape of a compatible type on the slide as a fallback match."""
    role = zone.get("role", "other")
    is_image = "image" in role.lower()
    
    for s in slide.shapes:
        if s.shape_id in filled_shape_ids:
            continue
            
        if is_image:
            # Check if it has a blip element
            blip_el = None
            if hasattr(s, "element") and s.element is not None:
                blip_el = s.element.find('.//' + qn('a:blip'))
            if blip_el is None and hasattr(s, "_element") and s._element is not None:
                blip_el = s._element.find('.//' + qn('a:blip'))
            if blip_el is not None:
                return s
    return None


# ---------------------------------------------------------------------------
# Capacity estimation helper
# ---------------------------------------------------------------------------

def _estimate_capacity(shape, font_size=None) -> int:
    """Estimate word capacity using template placeholder text count."""
    try:
        if shape.has_text_frame:
            existing_text = shape.text_frame.text or ""
            words = existing_text.split()
            existing_word_count = len(words)
            if not existing_text.strip():
                existing_word_count = 20
            return min(existing_word_count * 2, 40)
    except Exception:
        pass
    return 40


# ---------------------------------------------------------------------------
# Content value lookup and fallbacks
# ---------------------------------------------------------------------------

def _get_role_value_with_fallback(role: str, content: dict, layout: str) -> str:
    """Retrieve content value for a specific zone role from content or zone_content, with no generic fallbacks."""
    if not content:
        return ""

    # 1. Direct role key lookup
    if role in content:
        val = content[role]
        if val:
            if isinstance(val, list):
                return "\n".join(str(v) for v in val).strip()
            return str(val).strip()

    # 2. Check inside zone_content if present
    zone_content = content.get("zone_content", {})
    if isinstance(zone_content, dict) and role in zone_content:
        val = zone_content[role]
        if val:
            if isinstance(val, list):
                return "\n".join(str(v) for v in val).strip()
            return str(val).strip()

    # 3. Image keyword lookup (only for roles containing 'image')
    if "image" in role.lower():
        for key in (f"{role}_keyword", "image_keyword"):
            if key in content and content[key]:
                return str(content[key]).strip()
            if isinstance(zone_content, dict) and key in zone_content and zone_content[key]:
                return str(zone_content[key]).strip()
        # Fallback for image2
        if role == "image2":
            for key in ("image_keyword_2", "image2_keyword"):
                if key in content and content[key]:
                    return str(content[key]).strip()
                if isinstance(zone_content, dict) and key in zone_content and zone_content[key]:
                    return str(zone_content[key]).strip()

    # If no content exists for the role, return empty string instead of fallback chains
    return ""


# ---------------------------------------------------------------------------
# Text content in-place replacement
# ---------------------------------------------------------------------------

def _fill_text_in_shape(shape, text_value):
    """Replace content inside text shape in-place while preserving design, paragraphs, and runs."""
    if not shape.has_text_frame:
        return

    tf = shape.text_frame

    # Split lines into separate elements
    if isinstance(text_value, list):
        lines = [str(line).strip() for line in text_value if str(line).strip()]
    else:
        lines = [line.strip() for line in str(text_value).split("\n") if line.strip()]

    if not lines:
        lines = [""]

    paras = list(tf.paragraphs)
    P = len(paras)

    if P == 0:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "\n".join(lines)
        return

    # Map lines to paragraphs to preserve original paragraph count.
    # If N is original paragraph count, we fill the first N-1 paragraphs with one line each.
    # The remaining lines are joined by '\n' and filled into the N-th paragraph.
    # If there are fewer lines than paragraphs, the extra paragraphs are set to "".
    para_texts = []
    if P <= 1:
        para_texts.append("\n".join(lines))
    else:
        for i in range(P):
            if i < P - 1:
                if i < len(lines):
                    para_texts.append(lines[i])
                else:
                    para_texts.append("")
            else:
                if i < len(lines):
                    para_texts.append("\n".join(lines[i:]))
                else:
                    para_texts.append("")

    # Fill in-place to preserve paragraph/run formatting
    for i, p in enumerate(paras):
        val_text = para_texts[i]
        if p.runs:
            p.runs[0].text = val_text
        else:
            run = p.add_run()
            run.text = val_text


def _extract_paragraph_props(para) -> dict:
    """Extract formatting settings from a paragraph."""
    return {
        "alignment": para.alignment,
        "level": para.level,
        "line_spacing": para.line_spacing,
        "space_before": para.space_before,
        "space_after": para.space_after
    }


def _apply_paragraph_props(para, props: dict):
    """Apply paragraph formatting settings."""
    try:
        para.alignment = props.get("alignment")
    except Exception:
        pass
    try:
        para.level = props.get("level", 0)
    except Exception:
        pass
    try:
        para.line_spacing = props.get("line_spacing")
    except Exception:
        pass
    try:
        para.space_before = props.get("space_before")
    except Exception:
        pass
    try:
        para.space_after = props.get("space_after")
    except Exception:
        pass


def _extract_run_props(run) -> dict:
    """Extract font styling properties from a run."""
    props = {}
    try:
        font = run.font
        if font.name:
            props["name"] = font.name
        if font.size:
            props["size"] = font.size
        props["bold"] = font.bold
        props["italic"] = font.italic
        props["underline"] = font.underline
        
        if font.color:
            props["color_type"] = font.color.type
            if font.color.type == 1:  # RGB
                props["color_rgb"] = font.color.rgb
            elif font.color.type == 2:  # Scheme Color
                props["color_scheme"] = font.color.theme_color
    except Exception:
        pass
    return props


def _apply_run_props(run, props: dict):
    """Apply preserved font styling properties to a run."""
    try:
        font = run.font
        if "name" in props:
            font.name = props["name"]
        if "size" in props:
            font.size = props["size"]
        if "bold" in props:
            font.bold = props["bold"]
        if "italic" in props:
            font.italic = props["italic"]
        if "underline" in props:
            font.underline = props["underline"]
            
        if "color_type" in props:
            if props["color_type"] == 1 and "color_rgb" in props:
                font.color.rgb = props["color_rgb"]
            elif props["color_type"] == 2 and "color_scheme" in props:
                font.color.theme_color = props["color_scheme"]
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Image replacement via XML relationship override
# ---------------------------------------------------------------------------

def _replace_image_in_shape(slide, shape, image_keyword):
    """Overwrite image contents in place by modifying the relationship target blob."""
    if not image_keyword:
        return

    image_path = _download_image_for_template(image_keyword)
    if not image_path or not os.path.exists(image_path):
        return

    try:
        blip_el = None
        # Check standard XML tree for a:blip
        if hasattr(shape, "element") and shape.element is not None:
            blip_el = shape.element.find('.//' + qn('a:blip'))
            
        if blip_el is None and hasattr(shape, "_element") and shape._element is not None:
            blip_el = shape._element.find('.//' + qn('a:blip'))

        if blip_el is not None:
            rId = blip_el.get(qn('r:embed')) or blip_el.rEmbed
            if rId:
                image_part = slide.part.rels[rId].target_part
                with open(image_path, "rb") as img_f:
                    new_bytes = img_f.read()
                image_part.blob = new_bytes
                print(f"[TEMPLATE_FILLER] Overwrote image relation {rId} in shape {shape.shape_id}")
    except Exception as e:
        print(f"[TEMPLATE_FILLER] Warning: Failed to replace image blob directly: {e}")


def _download_image_for_template(keyword):
    """Download an image using the visual engine."""
    try:
        from services.visual_engine import get_hero_image
        return get_hero_image(keyword)
    except ImportError:
        print("[TEMPLATE_FILLER] visual_engine not available, skipping image download")
        return None
    except Exception as e:
        print(f'[TEMPLATE_FILLER] Image download failed for "{keyword}": {e}')
        return None


# ---------------------------------------------------------------------------
# Slide deletion helper
# ---------------------------------------------------------------------------

def _delete_slide(prs, slide_index):
    """Delete a slide by index from the presentation."""
    rId = prs.slides._sldIdLst[slide_index].get(qn("r:id"))
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[slide_index]
